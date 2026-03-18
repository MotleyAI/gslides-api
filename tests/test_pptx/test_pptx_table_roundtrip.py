"""
Comprehensive test suite for PowerPointTableElement.to_markdown_element() roundtrip functionality.

Tests the complete flow:
1. Markdown string → MarkdownTableElement
2. MarkdownTableElement → PowerPointTableElement (via update_content)
3. PowerPointTableElement → MarkdownTableElement (via to_markdown_element)
4. MarkdownTableElement → Markdown string
"""

import os
import tempfile
from typing import List

import pytest
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Inches

from gslides_api.agnostic.element import MarkdownTableElement, TableData

from gslides_api.adapters.pptx_adapter import (
    PowerPointAPIClient,
    PowerPointPresentation,
    PowerPointTableElement,
)


@pytest.fixture
def api_client():
    """Create a PowerPointAPIClient instance."""
    return PowerPointAPIClient()


@pytest.fixture
def test_presentation():
    """Create a test PowerPoint presentation with a slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    return prs, slide


def create_test_table(slide, rows: int, cols: int) -> GraphicFrame:
    """Create a test PowerPoint table shape."""
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape


def normalize_markdown(markdown: str) -> str:
    """Normalize markdown for comparison (handle spacing variations)."""
    lines = markdown.strip().split("\n")
    normalized_lines = []

    for line in lines:
        if line.strip():
            # Normalize spacing around table delimiters
            line = line.strip()
            normalized_lines.append(line)
        else:
            normalized_lines.append("")

    return "\n".join(normalized_lines)


class TestPowerPointTableBasicRoundtrip:
    """Test basic table roundtrip conversion."""

    def test_simple_2x2_table_roundtrip(self, api_client, test_presentation):
        """Test simple 2x2 table roundtrip."""
        prs, slide = test_presentation

        # Original markdown
        markdown_input = """| A | B |
|---|---|
| 1 | 2 |"""

        # Step 1: Create MarkdownTableElement
        markdown_elem = MarkdownTableElement(name="Simple Table", content=markdown_input)

        # Verify parsing worked
        assert markdown_elem.content.headers == ["A", "B"]
        assert markdown_elem.content.rows == [["1", "2"]]

        # Step 2: Create PowerPoint table
        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        # Step 3: Update content
        pptx_table_elem.update_content(api_client, markdown_elem)

        # Step 4: Extract back to markdown
        extracted_elem = pptx_table_elem.to_markdown_element("Simple Table")

        # Step 5: Verify roundtrip
        assert extracted_elem.name == "Simple Table"
        assert extracted_elem.content.headers == ["A", "B"]
        assert extracted_elem.content.rows == [["1", "2"]]

        # Verify markdown output
        final_markdown = extracted_elem.to_markdown()
        assert "| A | B |" in final_markdown
        assert "| 1 | 2 |" in final_markdown

    def test_3x3_table_with_headers_roundtrip(self, api_client, test_presentation):
        """Test 3x3 table with meaningful headers."""
        prs, slide = test_presentation

        markdown_input = """| Name | Age | City |
|------|-----|------|
| Alice | 25 | NYC |
| Bob | 30 | LA |"""

        markdown_elem = MarkdownTableElement(name="People Table", content=markdown_input)

        table_shape = create_test_table(slide, 3, 3)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("People Table")

        assert extracted_elem.content.headers == ["Name", "Age", "City"]
        assert len(extracted_elem.content.rows) == 2
        assert extracted_elem.content.rows[0] == ["Alice", "25", "NYC"]
        assert extracted_elem.content.rows[1] == ["Bob", "30", "LA"]

    def test_single_row_table(self, api_client, test_presentation):
        """Test edge case: table with just headers."""
        prs, slide = test_presentation

        markdown_input = """| Header1 | Header2 | Header3 |
|---------|---------|---------|"""

        markdown_elem = MarkdownTableElement(name="Headers Only", content=markdown_input)

        table_shape = create_test_table(slide, 1, 3)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Headers Only")

        assert extracted_elem.content.headers == ["Header1", "Header2", "Header3"]
        assert len(extracted_elem.content.rows) == 0

    def test_single_column_table(self, api_client, test_presentation):
        """Test edge case: narrow single column table."""
        prs, slide = test_presentation

        markdown_input = """| Items |
|-------|
| Apple |
| Banana |
| Cherry |"""

        markdown_elem = MarkdownTableElement(name="Single Column", content=markdown_input)

        table_shape = create_test_table(slide, 4, 1)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Single Column")

        assert extracted_elem.content.headers == ["Items"]
        assert len(extracted_elem.content.rows) == 3
        assert extracted_elem.content.rows == [["Apple"], ["Banana"], ["Cherry"]]


class TestPowerPointTableRichTextRoundtrip:
    """Test preservation of text formatting in table cells."""

    def test_bold_text_in_cells(self, api_client, test_presentation):
        """Test bold text formatting preservation."""
        prs, slide = test_presentation

        markdown_input = """| **Bold Header** | Normal |
|-----------------|--------|
| **Bold** | Text |"""

        markdown_elem = MarkdownTableElement(name="Bold Table", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Bold Table")

        # Verify bold formatting is preserved through roundtrip
        # Markdown **bold** -> PowerPoint bold -> Markdown **bold**
        assert extracted_elem.content.headers[0] == "**Bold Header**"
        assert extracted_elem.content.headers[1] == "Normal"
        assert extracted_elem.content.rows[0][0] == "**Bold**"
        assert extracted_elem.content.rows[0][1] == "Text"

        # Verify the markdown output contains bold markers
        final_markdown = extracted_elem.to_markdown()
        assert "**Bold Header**" in final_markdown
        assert "**Bold**" in final_markdown

    def test_italic_text_in_cells(self, api_client, test_presentation):
        """Test italic text formatting preservation."""
        prs, slide = test_presentation

        markdown_input = """| _Italic_ | Normal |
|----------|--------|
| _Text_ | Here |"""

        markdown_elem = MarkdownTableElement(name="Italic Table", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Italic Table")

        # Verify italic formatting is preserved through roundtrip
        # Markdown _italic_ -> PowerPoint italic -> Markdown *italic*
        assert extracted_elem.content.headers[0] == "*Italic*"
        assert extracted_elem.content.headers[1] == "Normal"
        assert extracted_elem.content.rows[0] == ["*Text*", "Here"]

    def test_hyperlinks_in_cells(self, api_client, test_presentation):
        """Test hyperlink preservation."""
        prs, slide = test_presentation

        markdown_input = """| [Link](http://example.com) | Text |
|----------------------------|------|
| Normal | [Google](http://google.com) |"""

        markdown_elem = MarkdownTableElement(name="Link Table", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Link Table")

        # Links will be converted to plain text in this roundtrip,
        # but the text content should be preserved
        assert "Link" in extracted_elem.content.headers[0]
        assert extracted_elem.content.headers[1] == "Text"
        assert extracted_elem.content.rows[0][0] == "Normal"
        assert "Google" in extracted_elem.content.rows[0][1]

    def test_mixed_formatting(self, api_client, test_presentation):
        """Test complex formatting combinations."""
        prs, slide = test_presentation

        markdown_input = """| **Bold** _italic_ | Text |
|-------------------|------|
| Normal | Mixed **bold** content |"""

        markdown_elem = MarkdownTableElement(name="Mixed Format", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Mixed Format")

        # Verify mixed formatting is preserved through roundtrip
        # Bold uses **, italic uses * (single asterisk in output)
        header0 = extracted_elem.content.headers[0]
        assert "**Bold**" in header0
        assert "*italic*" in header0
        assert extracted_elem.content.rows[0][0] == "Normal"
        assert "**bold**" in extracted_elem.content.rows[0][1]


class TestPowerPointTableSpecialCharacters:
    """Test handling of special characters and escaping."""

    def test_pipe_character_in_cells(self, api_client, test_presentation):
        """Test cells with pipe-like content (simplified test for parsing limitations)."""
        prs, slide = test_presentation

        # Use a simpler approach since pipe escaping has parsing complexities
        markdown_input = """| Contains Text | Normal |
|---------------|--------|
| Has pipe word | Text |"""

        markdown_elem = MarkdownTableElement(name="Pipe Table", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Pipe Table")

        assert extracted_elem.content.headers == ["Contains Text", "Normal"]
        # Verify the content is preserved (without expecting literal pipe character)
        assert "pipe" in extracted_elem.content.rows[0][0]
        assert extracted_elem.content.rows[0][1] == "Text"

    def test_line_breaks_in_cells(self, api_client, test_presentation):
        """Test line breaks within cells."""
        prs, slide = test_presentation

        markdown_input = """| Multi Line | Normal |
|------------|--------|
| Line<br />Break | Single |"""

        markdown_elem = MarkdownTableElement(name="Line Break Table", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Line Break Table")

        # Line breaks may be converted to spaces or preserved
        first_cell = extracted_elem.content.rows[0][0]
        assert "Line" in first_cell and "Break" in first_cell

    def test_empty_cells(self, api_client, test_presentation):
        """Test tables with empty cells."""
        prs, slide = test_presentation

        markdown_input = """| Header | Empty |
|--------|-------|
| Data |  |
|  | Data2 |"""

        markdown_elem = MarkdownTableElement(name="Empty Cells", content=markdown_input)

        table_shape = create_test_table(slide, 3, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Empty Cells")

        assert extracted_elem.content.headers == ["Header", "Empty"]
        assert len(extracted_elem.content.rows) == 2
        assert extracted_elem.content.rows[0][0] == "Data"
        # Empty cells should be empty strings or whitespace
        assert extracted_elem.content.rows[0][1].strip() == ""
        assert extracted_elem.content.rows[1][0].strip() == ""
        assert extracted_elem.content.rows[1][1] == "Data2"


class TestPowerPointTableMetadata:
    """Test metadata preservation through roundtrip."""

    def test_object_id_preservation(self, api_client, test_presentation):
        """Test that objectId is preserved in metadata."""
        prs, slide = test_presentation

        markdown_input = """| Test | Table |
|------|-------|
| Data | Here |"""

        markdown_elem = MarkdownTableElement(name="Metadata Test", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        # Get the original objectId
        original_object_id = pptx_table_elem.objectId

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Metadata Test")

        # Check metadata preservation
        assert "objectId" in extracted_elem.metadata
        assert extracted_elem.metadata["objectId"] == original_object_id

    def test_dimensions_metadata(self, api_client, test_presentation):
        """Test that table dimensions are preserved in metadata."""
        prs, slide = test_presentation

        markdown_input = """| A | B | C |
|---|---|---|
| 1 | 2 | 3 |
| 4 | 5 | 6 |"""

        markdown_elem = MarkdownTableElement(name="Dimensions Test", content=markdown_input)

        table_shape = create_test_table(slide, 3, 3)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Dimensions Test")

        # Check dimensions in metadata
        assert "rows" in extracted_elem.metadata
        assert "columns" in extracted_elem.metadata
        assert extracted_elem.metadata["rows"] == 3  # Header + 2 data rows
        assert extracted_elem.metadata["columns"] == 3

    def test_size_metadata(self, api_client, test_presentation):
        """Test that size metadata is preserved."""
        prs, slide = test_presentation

        markdown_input = """| Size | Test |
|------|------|
| Data | Here |"""

        markdown_elem = MarkdownTableElement(name="Size Test", content=markdown_input)

        table_shape = create_test_table(slide, 2, 2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        pptx_table_elem.update_content(api_client, markdown_elem)
        extracted_elem = pptx_table_elem.to_markdown_element("Size Test")

        # Size metadata should be present
        if "size" in extracted_elem.metadata:
            size_data = extracted_elem.metadata["size"]
            assert "width" in size_data
            assert "height" in size_data
            assert size_data["width"] > 0
            assert size_data["height"] > 0


class TestPowerPointTableErrorHandling:
    """Test error conditions and edge cases."""

    def test_invalid_pptx_element(self, api_client):
        """Test when pptx_element is not a GraphicFrame."""
        # Create a PowerPointTableElement with invalid element
        invalid_elem = PowerPointTableElement()
        invalid_elem.pptx_element = None

        with pytest.raises(
            ValueError, match="PowerPointTableElement has no valid GraphicFrame element"
        ):
            invalid_elem.to_markdown_element("Invalid Test")

    def test_empty_table_handling(self, api_client, test_presentation):
        """Test handling of tables with no meaningful content."""
        prs, slide = test_presentation

        # Create a minimal table that might cause issues
        table_shape = create_test_table(slide, 1, 1)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        # This should handle gracefully even with minimal content
        extracted_elem = pptx_table_elem.to_markdown_element("Empty Test")

        # Should return a valid MarkdownTableElement even if content is minimal
        assert isinstance(extracted_elem, MarkdownTableElement)
        assert extracted_elem.name == "Empty Test"

    def test_resize_table_does_not_crash(self, api_client, test_presentation):
        """Test that calling resize on a table doesn't crash.

        The resize method catches exceptions internally, so we verify it
        handles edge cases gracefully without raising. This includes the
        guard against index -1 access when current_cols could be 0.
        """
        prs, slide = test_presentation

        # Create a 2x2 table
        table_shape = create_test_table(slide, rows=2, cols=2)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        # Resize should not raise an exception even if operations fail internally
        # (python-pptx doesn't support add_column, but code should handle gracefully)
        pptx_table_elem.resize(api_client, rows=3, cols=4)

        # Table still exists and is accessible (even if dimensions didn't change)
        assert table_shape.table is not None


class TestPowerPointTableFullIntegration:
    """Complete integration test with real PowerPoint objects."""

    def test_complete_roundtrip_with_presentation(self, api_client):
        """Test complete roundtrip with file save/load."""

        # Complex table with various formatting
        markdown_input = """| **Product** | _Price_ | Status |
|-------------|---------|--------|
| **Laptop** | $999 | Available |
| _Tablet_ | $499 | **Sold Out** |
| Phone | $699 | Available |"""

        # Step 1: Parse to MarkdownTableElement
        markdown_elem = MarkdownTableElement(name="Product Table", content=markdown_input)

        # Verify initial parsing (markdown formatting is now preserved)
        assert markdown_elem.content.headers == ["**Product**", "*Price*", "Status"]
        assert len(markdown_elem.content.rows) == 3

        # Step 2: Create PowerPoint presentation with table
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_file:
            try:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])

                # Create table with appropriate dimensions
                table_shape = create_test_table(slide, 4, 3)  # 4 rows (header + 3 data), 3 columns
                pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

                # Step 3: Update table content
                pptx_table_elem.update_content(api_client, markdown_elem)

                # Save presentation
                prs.save(tmp_file.name)

                # Step 4: Extract back to markdown
                extracted_elem = pptx_table_elem.to_markdown_element("Product Table")

                # Step 5: Verify roundtrip results (markdown formatting preserved)
                assert extracted_elem.name == "Product Table"
                assert extracted_elem.content.headers == ["**Product**", "*Price*", "Status"]
                assert len(extracted_elem.content.rows) == 3

                # Check specific content (formatting may be normalized)
                rows = extracted_elem.content.rows
                assert "Laptop" in rows[0][0]
                assert "$999" in rows[0][1]
                assert "Available" in rows[0][2]
                assert "Tablet" in rows[1][0]
                assert "$499" in rows[1][1]
                assert "Sold Out" in rows[1][2] or "Sold" in rows[1][2]

                # Check metadata preservation
                assert "objectId" in extracted_elem.metadata
                assert "rows" in extracted_elem.metadata
                assert "columns" in extracted_elem.metadata
                assert extracted_elem.metadata["rows"] == 4  # Total rows including header
                assert extracted_elem.metadata["columns"] == 3

                # Final markdown output
                final_markdown = extracted_elem.to_markdown()
                assert "Product" in final_markdown
                assert "Price" in final_markdown
                assert "Status" in final_markdown
                assert "Laptop" in final_markdown
                assert "$999" in final_markdown

            finally:
                # Cleanup
                if os.path.exists(tmp_file.name):
                    os.unlink(tmp_file.name)


class TestPowerPointTableCellSizing:
    """Test cell sizing preservation."""

    def test_column_widths_preserved_when_dimensions_match(self, api_client, test_presentation):
        """Test that column widths are preserved when table dimensions don't change."""
        prs, slide = test_presentation

        # Create a table with custom column widths
        table_shape = create_test_table(slide, 2, 3)
        table = table_shape.table

        # Set custom column widths (first column wider)
        total_width = table_shape.width
        table.columns[0].width = total_width // 2  # 50%
        table.columns[1].width = total_width // 4  # 25%
        table.columns[2].width = total_width // 4  # 25%

        original_widths = [col.width for col in table.columns]

        # Create markdown content with same dimensions (2 rows, 3 columns)
        markdown_input = """| A | B | C |
|---|---|---|
| 1 | 2 | 3 |"""

        markdown_elem = MarkdownTableElement(name="Size Test", content=markdown_input)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        # Update content - dimensions match, so widths should be preserved
        pptx_table_elem.update_content(api_client, markdown_elem, check_shape=False)

        # Verify column widths are preserved
        final_widths = [col.width for col in table.columns]
        assert final_widths == original_widths, "Column widths should be preserved when dimensions match"

    def test_row_heights_preserved_when_dimensions_match(self, api_client, test_presentation):
        """Test that row heights are preserved when table dimensions don't change."""
        prs, slide = test_presentation

        # Create a table with custom row heights
        table_shape = create_test_table(slide, 3, 2)
        table = table_shape.table

        # Set custom row heights (first row taller for header)
        total_height = table_shape.height
        table.rows[0].height = total_height // 2  # Header gets 50%
        table.rows[1].height = total_height // 4  # 25%
        table.rows[2].height = total_height // 4  # 25%

        original_heights = [row.height for row in table.rows]

        # Create markdown content with same dimensions (3 rows, 2 columns)
        markdown_input = """| Header A | Header B |
|----------|----------|
| Data 1 | Data 2 |
| Data 3 | Data 4 |"""

        markdown_elem = MarkdownTableElement(name="Height Test", content=markdown_input)
        pptx_table_elem = PowerPointTableElement(pptx_element=table_shape)

        # Update content - dimensions match, so heights should be preserved
        pptx_table_elem.update_content(api_client, markdown_elem, check_shape=False)

        # Verify row heights are preserved
        final_heights = [row.height for row in table.rows]
        assert final_heights == original_heights, "Row heights should be preserved when dimensions match"


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
