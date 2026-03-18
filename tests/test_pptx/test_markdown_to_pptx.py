"""
Test suite for markdown_to_pptx module.

This module tests the conversion of markdown to PowerPoint text frames,
specifically focusing on the _apply_style_to_run function and font size handling.
"""

import pytest
from unittest.mock import Mock, MagicMock

from gslides_api.agnostic.text import (
    FullTextStyle,
    MarkdownRenderableStyle,
    RichStyle,
    AbstractColor,
)

from gslides_api.pptx.markdown_to_pptx import _apply_style_to_run
from gslides_api.adapters.pptx_adapter import _extract_base_style_from_textframe


class TestApplyStyleToRunFontSize:
    """Test font size handling in _apply_style_to_run."""

    def create_mock_run(self):
        """Create a mock run with font properties."""
        mock_run = Mock()
        mock_font = Mock()
        mock_font.bold = None
        mock_font.italic = None
        mock_font.underline = None
        mock_font.strikethrough = None
        mock_font.name = None
        mock_font.size = None
        mock_hyperlink = Mock()
        mock_hyperlink.address = None
        mock_run.font = mock_font
        mock_run.hyperlink = mock_hyperlink
        return mock_run

    def test_font_size_applied_when_set(self):
        """Test that font size is applied when font_size_pt is set."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            rich=RichStyle(font_size_pt=12.0)
        )

        _apply_style_to_run(mock_run, style)

        # Font size should be set
        assert mock_run.font.size is not None
        # The size should be set to Pt(12.0)
        assert mock_run.font.size.pt == 12.0

    def test_font_size_not_applied_when_not_set(self):
        """Test that no font size is applied when font_size_pt is None."""
        mock_run = self.create_mock_run()

        style = FullTextStyle()

        _apply_style_to_run(mock_run, style)

        # Font size should NOT be set
        assert mock_run.font.size is None

    def test_font_size_with_different_pt_values(self):
        """Test various PT font size values are correctly applied."""
        test_sizes = [8.0, 10.0, 12.0, 14.0, 18.0, 24.0, 36.0, 48.0]

        for size in test_sizes:
            mock_run = self.create_mock_run()
            style = FullTextStyle(
                rich=RichStyle(font_size_pt=size)
            )

            _apply_style_to_run(mock_run, style)

            assert mock_run.font.size is not None
            assert mock_run.font.size.pt == size


class TestApplyStyleToRunOtherProperties:
    """Test other style properties in _apply_style_to_run."""

    def create_mock_run(self):
        """Create a mock run with font properties."""
        mock_run = Mock()
        mock_font = Mock()
        mock_font.bold = None
        mock_font.italic = None
        mock_font.underline = None
        mock_font.strikethrough = None
        mock_font.name = None
        mock_font.size = None
        mock_color = Mock()
        mock_color.rgb = None
        mock_font.color = mock_color
        mock_hyperlink = Mock()
        mock_hyperlink.address = None
        mock_run.font = mock_font
        mock_run.hyperlink = mock_hyperlink
        return mock_run

    def test_bold_applied(self):
        """Test that bold is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=True)
        )

        _apply_style_to_run(mock_run, style)

        assert mock_run.font.bold is True

    def test_italic_applied(self):
        """Test that italic is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            markdown=MarkdownRenderableStyle(italic=True)
        )

        _apply_style_to_run(mock_run, style)

        assert mock_run.font.italic is True

    def test_underline_applied(self):
        """Test that underline is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            rich=RichStyle(underline=True)
        )

        _apply_style_to_run(mock_run, style)

        assert mock_run.font.underline is True

    def test_strikethrough_applied(self):
        """Test that strikethrough is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            markdown=MarkdownRenderableStyle(strikethrough=True)
        )

        _apply_style_to_run(mock_run, style)

        assert mock_run.font.strikethrough is True

    def test_font_family_applied(self):
        """Test that font family is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            rich=RichStyle(font_family="Arial")
        )

        _apply_style_to_run(mock_run, style)

        assert mock_run.font.name == "Arial"

    def test_hyperlink_applied(self):
        """Test that hyperlink is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            markdown=MarkdownRenderableStyle(hyperlink="https://example.com")
        )

        _apply_style_to_run(mock_run, style)

        assert mock_run.hyperlink.address == "https://example.com"

    def test_foreground_color_applied(self):
        """Test that foreground color is applied correctly."""
        mock_run = self.create_mock_run()

        style = FullTextStyle(
            rich=RichStyle(
                foreground_color=AbstractColor(red=1.0, green=0.0, blue=0.0)
            )
        )

        _apply_style_to_run(mock_run, style)

        # RGBColor should have been called with the converted tuple
        assert mock_run.font.color.rgb is not None

    def test_none_style_is_handled(self):
        """Test that None style doesn't raise an error."""
        mock_run = self.create_mock_run()

        # This should not raise an error
        _apply_style_to_run(mock_run, None)

        # No properties should be modified
        assert mock_run.font.bold is None
        assert mock_run.font.italic is None

    def test_empty_style_is_handled(self):
        """Test that empty style explicitly sets bold/italic to False.

        Bold and italic are always explicitly set (True or False) to prevent
        inheritance from defRPr (default run properties). If defRPr has bold=True,
        text would appear bold unless we explicitly set bold=False.
        """
        mock_run = self.create_mock_run()

        style = FullTextStyle()

        _apply_style_to_run(mock_run, style)

        # Bold and italic should be explicitly False (not None) to prevent inheritance
        assert mock_run.font.bold is False
        assert mock_run.font.italic is False
        # Other properties should remain unmodified
        assert mock_run.font.size is None


class TestExtractBaseStyleFromTextframe:
    """Test _extract_base_style_from_textframe function."""

    def create_mock_text_frame(
        self,
        text: str = "Hello",
        bold: bool = False,
        italic: bool = False,
        font_name: str = None,
        font_size_pt: float = None,
        color_rgb: tuple = None,
        underline: bool = False,
    ):
        """Create a mock text frame with configurable properties."""
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()

        # Set run text
        mock_run.text = text

        # Set font properties
        mock_font = Mock()
        mock_font.bold = bold
        mock_font.italic = italic
        mock_font.name = font_name
        mock_font.underline = underline

        # Set font size
        if font_size_pt is not None:
            mock_size = Mock()
            mock_size.pt = font_size_pt
            mock_font.size = mock_size
        else:
            mock_font.size = None

        # Set color
        from pptx.enum.dml import MSO_COLOR_TYPE

        mock_color = Mock()
        if color_rgb is not None:
            mock_color.type = MSO_COLOR_TYPE.RGB
            mock_rgb = Mock()
            mock_rgb.__getitem__ = lambda self, idx: color_rgb[idx]
            mock_color.rgb = mock_rgb
        else:
            mock_color.type = None
            mock_color.rgb = None
        mock_font.color = mock_color

        mock_run.font = mock_font
        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]

        return mock_text_frame

    def test_extract_basic_style(self):
        """Test extracting basic style from text frame.

        Note: bold/italic are NOT extracted from base_style. They are always
        False because markdown content should control bold/italic formatting.
        This prevents issues like a bold header making ALL text bold.
        """
        text_frame = self.create_mock_text_frame(
            text="Hello",
            font_size_pt=24.0,
            font_name="Arial",
            bold=True,  # This is intentionally ignored
            italic=False,
        )

        style = _extract_base_style_from_textframe(text_frame)

        assert style is not None
        assert style.rich.font_size_pt == 24.0
        assert style.rich.font_family == "Arial"
        # Bold/italic are always False in base_style - markdown controls these
        assert style.markdown.bold is False
        assert style.markdown.italic is False

    def test_extract_color(self):
        """Test extracting color from text frame."""
        text_frame = self.create_mock_text_frame(
            text="Colored text",
            color_rgb=(128, 0, 255),  # Purple
        )

        style = _extract_base_style_from_textframe(text_frame)

        assert style is not None
        assert style.rich.foreground_color is not None
        # AbstractColor uses 0.0-1.0 scale
        assert abs(style.rich.foreground_color.red - 128 / 255) < 0.01
        assert style.rich.foreground_color.green == 0.0
        assert abs(style.rich.foreground_color.blue - 1.0) < 0.01

    def test_extract_underline(self):
        """Test extracting underline from text frame."""
        text_frame = self.create_mock_text_frame(
            text="Underlined",
            underline=True,
        )

        style = _extract_base_style_from_textframe(text_frame)

        assert style is not None
        assert style.rich.underline is True

    def test_empty_text_frame_returns_none(self):
        """Test that empty text frame returns None."""
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()
        mock_run.text = "   "  # Whitespace only
        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]

        style = _extract_base_style_from_textframe(mock_text_frame)

        assert style is None

    def test_no_runs_returns_none(self):
        """Test that text frame with no runs returns None."""
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_paragraph.runs = []
        mock_text_frame.paragraphs = [mock_paragraph]

        style = _extract_base_style_from_textframe(mock_text_frame)

        assert style is None

    def test_extracts_from_first_non_empty_run(self):
        """Test that style is extracted from first non-whitespace run.

        Note: bold/italic are NOT extracted from base_style - they are always
        False because markdown content should control bold/italic formatting.
        Only RichStyle properties (font, size, color) are extracted.
        """
        mock_text_frame = Mock()
        mock_paragraph = Mock()

        # First run is whitespace
        mock_run1 = Mock()
        mock_run1.text = "   "

        # Second run has content and style
        mock_run2 = Mock()
        mock_run2.text = "Content"
        mock_font2 = Mock()
        mock_font2.bold = True  # Intentionally ignored
        mock_font2.italic = True  # Intentionally ignored
        mock_font2.name = "Times"
        mock_font2.underline = False
        mock_size2 = Mock()
        mock_size2.pt = 18.0
        mock_font2.size = mock_size2
        mock_color2 = Mock()
        mock_color2.type = None
        mock_font2.color = mock_color2
        mock_run2.font = mock_font2

        mock_paragraph.runs = [mock_run1, mock_run2]
        mock_text_frame.paragraphs = [mock_paragraph]

        style = _extract_base_style_from_textframe(mock_text_frame)

        assert style is not None
        # Bold/italic are always False in base_style - markdown controls these
        assert style.markdown.bold is False
        assert style.markdown.italic is False
        # RichStyle properties ARE extracted
        assert style.rich.font_family == "Times"
        assert style.rich.font_size_pt == 18.0


class TestApplyMarkdownToTextframeWithAutoscale:
    """Test apply_markdown_to_textframe with autoscale=True.

    These tests ensure that bold/italic formatting set by markdown parsing
    is preserved after fit_text() is called. The fit_text() API in python-pptx
    sets bold/italic on ALL runs, which would break mixed formatting like
    "**bold** and regular text". We save and restore per-run styles to fix this.
    """

    def test_autoscale_preserves_mixed_bold_italic(self, tmp_path):
        """Test that autoscale preserves mixed bold/italic formatting.

        Markdown: **Gabrielle** is bold, rest is not.
        After autoscale with fit_text(), bold should still be preserved per-run.
        """
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        # Create a test presentation with a text box
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame

        # Set initial text (will be replaced)
        tf.text = "Placeholder"

        # Apply markdown with mixed bold - WITH autoscale
        markdown = "**Gabrielle Aura**\n- Approval rate 28%\n- **17 Booked meetings**\n- Other text"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
            rich=RichStyle(font_size_pt=15.0, font_family="Arial"),
        )
        apply_markdown_to_textframe(markdown, tf, base_style=base_style, autoscale=True)

        # Verify bold is preserved per-run
        # Para 0: "Gabrielle Aura" should be bold
        assert len(tf.paragraphs) >= 4
        para0 = tf.paragraphs[0]
        assert len(para0.runs) >= 1
        assert para0.runs[0].font.bold is True, "Gabrielle Aura should be bold"

        # Para 1: "Approval rate 28%" should NOT be bold
        para1 = tf.paragraphs[1]
        assert len(para1.runs) >= 1
        assert para1.runs[0].font.bold is False, "Approval rate should not be bold"

        # Para 2: "17 Booked meetings" should be bold
        para2 = tf.paragraphs[2]
        assert len(para2.runs) >= 1
        assert para2.runs[0].font.bold is True, "17 Booked meetings should be bold"

        # Para 3: "Other text" should NOT be bold
        para3 = tf.paragraphs[3]
        assert len(para3.runs) >= 1
        assert para3.runs[0].font.bold is False, "Other text should not be bold"

    def test_autoscale_preserves_italic(self, tmp_path):
        """Test that autoscale preserves italic formatting."""
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = "Placeholder"

        # Apply markdown with italic - double newline creates separate paragraphs
        # Note: marko parser creates 3 paragraphs: italic, empty, regular
        markdown = "*italic text*\n\nregular text"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
            rich=RichStyle(font_size_pt=12.0),
        )
        apply_markdown_to_textframe(markdown, tf, base_style=base_style, autoscale=True)

        # Find paragraphs with content
        paras_with_content = [p for p in tf.paragraphs if p.runs]
        assert len(paras_with_content) >= 2

        # First content para: "italic text" should be italic
        assert paras_with_content[0].runs[0].font.italic is True, "italic text should be italic"

        # Second content para: "regular text" should NOT be italic
        assert paras_with_content[1].runs[0].font.italic is False, "regular text should not be italic"

    def test_autoscale_without_autoscale_flag(self, tmp_path):
        """Test that without autoscale=True, bold/italic are also correct.

        This is a baseline test - formatting should work correctly regardless.
        """
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = "Placeholder"

        # Apply markdown WITHOUT autoscale - double newline for separate paragraphs
        # Note: marko parser may create 3 paragraphs: bold, empty, regular
        markdown = "**Bold**\n\nRegular"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
        )
        apply_markdown_to_textframe(markdown, tf, base_style=base_style, autoscale=False)

        # Find paragraphs with content
        paras_with_content = [p for p in tf.paragraphs if p.runs]
        assert len(paras_with_content) >= 2

        # First content para: "Bold" should be bold
        assert paras_with_content[0].runs[0].font.bold is True

        # Second content para: "Regular" should NOT be bold
        assert paras_with_content[1].runs[0].font.bold is False


class TestSoftLineBreakHandling:
    """Test handling of soft line breaks (vertical tab) in text."""

    def test_soft_line_break_creates_proper_xml_break(self, tmp_path):
        """Test that \\x0b (vertical tab) creates proper <a:br> element."""
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = ""

        # Apply text with soft line break (vertical tab)
        text_with_break = "Your results\x0b(Samplead)"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
            rich=RichStyle(font_size_pt=12.0),
        )
        apply_markdown_to_textframe(text_with_break, tf, base_style=base_style)

        # Verify the result - should have two runs with a break between
        para = tf.paragraphs[0]
        assert len(para.runs) == 2, "Should have two runs separated by line break"
        assert para.runs[0].text == "Your results"
        assert para.runs[1].text == "(Samplead)"

        # Verify paragraph text reads back with \x0b (not _x000B_)
        assert para.text == "Your results\x0b(Samplead)", "Paragraph text should contain soft line break"

    def test_escaped_x000b_is_converted_to_break(self, tmp_path):
        """Test that literal '_x000B_' string is converted to proper break."""
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = ""

        # Apply text with escaped soft line break (this is what happens if
        # text was previously corrupted by python-pptx)
        text_with_escaped_break = "Your results_x000B_(Samplead)"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
            rich=RichStyle(font_size_pt=12.0),
        )
        apply_markdown_to_textframe(text_with_escaped_break, tf, base_style=base_style)

        # Verify the result - should NOT contain literal _x000B_
        para = tf.paragraphs[0]
        full_text = "".join(run.text for run in para.runs)
        assert "_x000B_" not in full_text, "Literal _x000B_ should be converted to break"

        # Should have proper line break
        assert para.text == "Your results\x0b(Samplead)", "Should have proper soft line break"

    def test_multiple_soft_line_breaks(self, tmp_path):
        """Test handling of multiple soft line breaks in text."""
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = ""

        # Apply text with multiple soft line breaks
        text_with_breaks = "Line 1\x0bLine 2\x0bLine 3"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
            rich=RichStyle(font_size_pt=12.0),
        )
        apply_markdown_to_textframe(text_with_breaks, tf, base_style=base_style)

        # Verify the result - should have three runs
        para = tf.paragraphs[0]
        assert len(para.runs) == 3, "Should have three runs separated by line breaks"
        assert para.runs[0].text == "Line 1"
        assert para.runs[1].text == "Line 2"
        assert para.runs[2].text == "Line 3"

    def test_soft_line_break_with_bold_formatting(self, tmp_path):
        """Test that formatting is preserved across soft line breaks."""
        from pptx import Presentation
        from pptx.util import Inches
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        tf = textbox.text_frame
        tf.text = ""

        # Apply bold text with soft line break
        text_with_break = "**Line 1\x0bLine 2**"
        base_style = FullTextStyle(
            markdown=MarkdownRenderableStyle(bold=False, italic=False),
            rich=RichStyle(font_size_pt=12.0),
        )
        apply_markdown_to_textframe(text_with_break, tf, base_style=base_style)

        # Verify all runs are bold
        para = tf.paragraphs[0]
        for run in para.runs:
            if run.text.strip():  # Skip empty runs
                assert run.font.bold is True, f"Run '{run.text}' should be bold"
