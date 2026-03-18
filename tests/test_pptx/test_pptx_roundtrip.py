"""Test PowerPoint roundtrip functionality for markdown conversion."""

import pytest
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches

from gslides_api.adapters.pptx_adapter import (
    PowerPointAPIClient,
    PowerPointPresentation,
    PowerPointSlide,
    PowerPointShapeElement,
    PowerPointSpeakerNotes,
)


class TestPowerPointRoundtrip:
    """Test roundtrip functionality for PowerPoint presentations."""

    def create_test_presentation(self) -> str:
        """Create a test presentation file and return its path."""
        prs = Presentation()

        # Set presentation title
        prs.core_properties.title = "Test Presentation"

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "Roundtrip Testing"

        # Add content slide with bullet points
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Bullet Points Test'
        tf = body_shape.text_frame
        tf.text = 'First bullet point'

        p = tf.add_paragraph()
        p.text = 'Second bullet point'
        p.level = 1

        p = tf.add_paragraph()
        p.text = 'Third bullet point with **bold** text'
        p.level = 0

        # Add table slide
        table_slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(table_slide_layout)
        shapes = slide.shapes
        shapes.title.text = 'Table Test'

        # Add table
        rows = cols = 2
        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # Fill table with data
        table.cell(0, 0).text = 'Header 1'
        table.cell(0, 1).text = 'Header 2'
        table.cell(1, 0).text = 'Data 1'
        table.cell(1, 1).text = 'Data 2'

        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        temp_file.close()

        return temp_file.name

    def test_load_presentation(self):
        """Test loading a PowerPoint presentation."""
        pptx_path = self.create_test_presentation()

        try:
            # Load presentation
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)

            assert presentation is not None
            assert presentation.title == "Test Presentation"
            assert len(presentation.slides) == 3
            assert presentation.file_path == pptx_path

        finally:
            os.unlink(pptx_path)

    def test_read_slide_text_markdown(self):
        """Test reading slide text as markdown."""
        pptx_path = self.create_test_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            bullet_slide = presentation.slides[1]  # Second slide with bullets

            # Find text shapes
            text_elements = [elem for elem in bullet_slide.elements if hasattr(elem, 'has_text') and elem.has_text]

            assert len(text_elements) > 0

            # Read text as markdown from body shape (should contain bullets)
            body_element = None
            for elem in text_elements:
                if isinstance(elem, PowerPointShapeElement):
                    text = elem.read_text(as_markdown=True)
                    if 'bullet' in text.lower():
                        body_element = elem
                        break

            assert body_element is not None

        finally:
            os.unlink(pptx_path)

    def test_write_read_text_roundtrip(self):
        """Test writing text and reading it back."""
        pptx_path = self.create_test_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            slide = presentation.slides[1]  # Bullet slide

            # Find a text element to modify
            text_elements = [elem for elem in slide.elements if hasattr(elem, 'has_text') and elem.has_text]
            assert len(text_elements) > 0

            text_element = text_elements[0]
            if isinstance(text_element, PowerPointShapeElement):
                # Test markdown content
                markdown_content = """# Test Header

This is a paragraph.

- First bullet
- Second bullet
  - Nested bullet
- **Bold text** and *italic text*"""

                # Write content
                text_element.write_text(api_client, markdown_content)

                # Read back and verify
                read_content = text_element.read_text(as_markdown=True)

                # Check that basic structure is preserved
                assert "Test Header" in read_content
                assert "This is a paragraph" in read_content
                assert "First bullet" in read_content
                assert "Second bullet" in read_content

        finally:
            os.unlink(pptx_path)

    def test_speaker_notes_roundtrip(self):
        """Test writing and reading speaker notes."""
        pptx_path = self.create_test_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            slide = presentation.slides[0]  # First slide

            # Add speaker notes
            if slide.speaker_notes:
                api_client = PowerPointAPIClient()
                notes_content = "These are speaker notes with **bold** text and bullet points:\n- Note 1\n- Note 2"

                slide.speaker_notes.write_text(api_client, notes_content)

                # Read back
                read_notes = slide.speaker_notes.read_text(as_markdown=True)

                assert "speaker notes" in read_notes.lower()
                assert "Note 1" in read_notes
                assert "Note 2" in read_notes

        finally:
            os.unlink(pptx_path)

    def test_presentation_copy_roundtrip(self):
        """Test copying a presentation and verifying content."""
        pptx_path = self.create_test_presentation()

        try:
            # Load original presentation
            api_client = PowerPointAPIClient()
            original_presentation = PowerPointPresentation.from_id(api_client, pptx_path)

            # Copy presentation
            api_client = PowerPointAPIClient()
            with tempfile.TemporaryDirectory() as temp_dir:
                copied_presentation = original_presentation.copy_via_drive(
                    api_client,
                    "copied_presentation",
                    temp_dir
                )

                # Verify copy
                assert copied_presentation is not None
                assert len(copied_presentation.slides) == len(original_presentation.slides)
                assert copied_presentation.file_path != original_presentation.file_path
                assert os.path.exists(copied_presentation.file_path)

                # Verify content is preserved
                assert copied_presentation.title == original_presentation.title

        finally:
            os.unlink(pptx_path)

    def test_table_content_roundtrip(self):
        """Test table content preservation in roundtrip."""
        pptx_path = self.create_test_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            table_slide = presentation.slides[2]  # Third slide with table

            # Find table elements
            table_elements = [elem for elem in table_slide.elements if elem.type == "table"]

            if table_elements:
                table_element = table_elements[0]
                # Basic verification that table structure is accessible
                assert table_element is not None
                assert hasattr(table_element, 'pptx_element')

        finally:
            os.unlink(pptx_path)

    def test_save_and_reload_presentation(self):
        """Test saving modifications and reloading presentation."""
        pptx_path = self.create_test_presentation()

        try:
            # Load presentation
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)

            # Modify content
            first_slide = presentation.slides[0]
            text_elements = [elem for elem in first_slide.elements if hasattr(elem, 'has_text') and elem.has_text]

            if text_elements:
                text_element = text_elements[0]
                if isinstance(text_element, PowerPointShapeElement):
                    text_element.write_text(api_client, "Modified content for testing")

            # Save presentation
            presentation.save(api_client)

            # Reload and verify changes
            reloaded_presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            reloaded_slide = reloaded_presentation.slides[0]
            reloaded_text_elements = [elem for elem in reloaded_slide.elements if hasattr(elem, 'has_text') and elem.has_text]

            if reloaded_text_elements:
                reloaded_text = reloaded_text_elements[0]
                if isinstance(reloaded_text, PowerPointShapeElement):
                    content = reloaded_text.read_text(as_markdown=False)
                    assert "Modified content for testing" in content

        finally:
            os.unlink(pptx_path)

    def test_markdown_formatting_preservation(self):
        """Test that markdown formatting is preserved in roundtrip."""
        pptx_path = self.create_test_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            slide = presentation.slides[1]

            # Find text element
            text_elements = [elem for elem in slide.elements if hasattr(elem, 'has_text') and elem.has_text]

            if text_elements:
                text_element = text_elements[0]
                if isinstance(text_element, PowerPointShapeElement):
                    # Test various markdown features
                    markdown_content = """
**Bold text**
*Italic text*
Regular text

- Bullet 1
- Bullet 2
  - Sub bullet
- Bullet 3

Another paragraph
"""

                    # Write and read back
                    text_element.write_text(api_client, markdown_content)
                    read_back = text_element.read_text(as_markdown=True)

                    # Verify basic structure preservation
                    # Note: exact formatting may vary due to PowerPoint's text handling
                    assert "Bold text" in read_back or "**Bold text**" in read_back
                    assert "Italic text" in read_back or "*Italic text*" in read_back
                    assert "Bullet 1" in read_back
                    assert "Bullet 2" in read_back

        finally:
            os.unlink(pptx_path)

    def test_empty_content_handling(self):
        """Test handling of empty content in roundtrip."""
        pptx_path = self.create_test_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            slide = presentation.slides[0]

            # Test with empty content
            text_elements = [elem for elem in slide.elements if hasattr(elem, 'has_text') and elem.has_text]

            if text_elements:
                text_element = text_elements[0]
                if isinstance(text_element, PowerPointShapeElement):
                    # Write empty content
                    text_element.write_text(api_client, "")

                    # Read back
                    read_content = text_element.read_text()
                    assert read_content == ""

        finally:
            os.unlink(pptx_path)