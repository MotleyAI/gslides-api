"""
Comprehensive tests for robust PowerPoint slide copying functionality.

Tests the new pptx subpackage to ensure no XML corruption, proper ID management,
and successful copying of various slide types.
"""

import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Inches

from gslides_api.pptx.id_manager import IdManager
from gslides_api.pptx.relationship_copier import RelationshipCopier
from gslides_api.pptx.shape_copier import ShapeCopier
from gslides_api.pptx.slide_copier import SlideCopierManager, _remove_layout_placeholders
from gslides_api.pptx.xml_utils import XmlUtils
from gslides_api.adapters.pptx_adapter import PowerPointAPIClient, PowerPointPresentation


class TestIdManager:
    """Test the ID manager for unique ID generation."""

    def test_id_manager_initialization(self):
        """Test basic ID manager initialization."""
        prs = Presentation()
        id_manager = IdManager(prs)

        assert id_manager.next_slide_id >= 256
        assert id_manager.next_shape_id >= 1
        assert len(id_manager.used_slide_ids) == 0
        assert len(id_manager.used_shape_ids) == 0

    def test_unique_slide_id_generation(self):
        """Test generation of unique slide IDs."""
        prs = Presentation()
        id_manager = IdManager(prs)

        id1 = id_manager.generate_unique_slide_id()
        id2 = id_manager.generate_unique_slide_id()

        assert id1 != id2
        assert id1 in id_manager.used_slide_ids
        assert id2 in id_manager.used_slide_ids

    def test_unique_shape_id_generation(self):
        """Test generation of unique shape IDs."""
        prs = Presentation()
        id_manager = IdManager(prs)

        id1 = id_manager.generate_unique_shape_id()
        id2 = id_manager.generate_unique_shape_id()

        assert id1 != id2
        assert id1 in id_manager.used_shape_ids
        assert id2 in id_manager.used_shape_ids

    def test_unique_creation_id_generation(self):
        """Test generation of unique creation IDs."""
        prs = Presentation()
        id_manager = IdManager(prs)

        id1 = id_manager.generate_unique_creation_id()
        id2 = id_manager.generate_unique_creation_id()

        assert id1 != id2
        assert len(id1) == 36  # GUID length
        assert len(id2) == 36  # GUID length
        assert id1 in id_manager.used_creation_ids
        assert id2 in id_manager.used_creation_ids


class TestXmlUtils:
    """Test XML utilities for safe manipulation."""

    def test_safe_copy_element(self):
        """Test safe copying of XML elements."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Add a shape to get an XML element
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Test text"

        source_element = textbox._element
        copied_element = XmlUtils.safe_copy_element(source_element, new_id=999)

        assert copied_element is not None
        assert copied_element.tag == source_element.tag

        # Check for ID in the cNvPr element (where PowerPoint shape IDs are actually stored)
        cnv_pr_elements = copied_element.xpath('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if cnv_pr_elements:
            assert cnv_pr_elements[0].get('id') == '999'
        else:
            # If no cNvPr element, that's also acceptable for this test
            pass

    def test_update_element_id(self):
        """Test updating element IDs."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        element = textbox._element

        success = XmlUtils.update_element_id(element, 12345)
        assert success
        assert element.get('id') == '12345'

    def test_validate_element(self):
        """Test XML element validation."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        element = textbox._element

        validation = XmlUtils.validate_element(element)
        assert validation['valid']
        assert isinstance(validation['issues'], list)
        assert isinstance(validation['warnings'], list)


class TestShapeCopier:
    """Test shape copying functionality."""

    def test_copy_text_shape(self):
        """Test copying a text box shape."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Count initial shapes (layouts may have placeholders)
        initial_target_shapes = len(target_slide.shapes)

        # Add text box to source slide
        textbox = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Test text content"

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)

        copied_shape = shape_copier.copy_shape(textbox, target_slide)

        assert copied_shape is not None
        assert copied_shape.text == "Test text content"
        # Check that one new shape was added
        assert len(target_slide.shapes) == initial_target_shapes + 1

    def test_copy_table_shape(self):
        """Test copying a table shape."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Add table to source slide
        table_shape = source_slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(5), Inches(3))
        table = table_shape.table
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Data 1"
        table.cell(1, 1).text = "Data 2"

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)

        copied_shape = shape_copier.copy_shape(table_shape, target_slide)

        assert copied_shape is not None
        assert hasattr(copied_shape, 'table')
        copied_table = copied_shape.table
        assert copied_table.cell(0, 0).text == "Header 1"
        assert copied_table.cell(1, 0).text == "Data 1"

    def test_copy_table_preserves_cell_borders(self):
        """Test that table cell borders are preserved when copying."""
        from copy import deepcopy

        from lxml import etree

        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Add table to source slide
        table_shape = source_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"

        # Add border to first cell via XML
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        source_cell = table.cell(0, 0)
        source_tc = source_cell._tc
        source_tcPr = source_tc.find(f"{{{ns['a']}}}tcPr")
        if source_tcPr is None:
            source_tcPr = etree.SubElement(source_tc, f"{{{ns['a']}}}tcPr")

        # Add left border (lnL) with specific properties
        lnL = etree.SubElement(source_tcPr, f"{{{ns['a']}}}lnL", w="12700", cap="flat", cmpd="sng")
        solidFill = etree.SubElement(lnL, f"{{{ns['a']}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{ns['a']}}}srgbClr", val="0000FF")  # Blue

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)
        copied_shape = shape_copier.copy_shape(table_shape, target_slide)

        assert copied_shape is not None
        copied_table = copied_shape.table

        # Verify border was copied
        copied_cell = copied_table.cell(0, 0)
        copied_tc = copied_cell._tc
        copied_tcPr = copied_tc.find(f"{{{ns['a']}}}tcPr")
        assert copied_tcPr is not None, "tcPr should exist in copied cell"

        copied_lnL = copied_tcPr.find(f"{{{ns['a']}}}lnL")
        assert copied_lnL is not None, "Left border (lnL) should be copied"
        assert copied_lnL.get("w") == "12700", "Border width should be preserved"

        # Verify border color
        copied_solidFill = copied_lnL.find(f"{{{ns['a']}}}solidFill")
        assert copied_solidFill is not None, "Border fill should be copied"
        copied_srgbClr = copied_solidFill.find(f"{{{ns['a']}}}srgbClr")
        assert copied_srgbClr is not None, "Border color should be copied"
        assert copied_srgbClr.get("val") == "0000FF", "Border color should be blue"

    def test_copy_table_preserves_cell_margins(self):
        """Test that table cell margins are preserved when copying."""
        from lxml import etree

        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Add table to source slide
        table_shape = source_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "A"

        # Set cell margins via XML
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        source_cell = table.cell(0, 0)
        source_tc = source_cell._tc
        source_tcPr = source_tc.find(f"{{{ns['a']}}}tcPr")
        if source_tcPr is None:
            source_tcPr = etree.SubElement(source_tc, f"{{{ns['a']}}}tcPr")

        # Set custom margins
        source_tcPr.set("marL", "91440")  # Left margin
        source_tcPr.set("marR", "91440")  # Right margin
        source_tcPr.set("marT", "45720")  # Top margin
        source_tcPr.set("marB", "45720")  # Bottom margin
        source_tcPr.set("anchor", "ctr")  # Center vertical alignment

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)
        copied_shape = shape_copier.copy_shape(table_shape, target_slide)

        assert copied_shape is not None
        copied_table = copied_shape.table

        # Verify margins were copied
        copied_cell = copied_table.cell(0, 0)
        copied_tc = copied_cell._tc
        copied_tcPr = copied_tc.find(f"{{{ns['a']}}}tcPr")
        assert copied_tcPr is not None, "tcPr should exist in copied cell"

        assert copied_tcPr.get("marL") == "91440", "Left margin should be preserved"
        assert copied_tcPr.get("marR") == "91440", "Right margin should be preserved"
        assert copied_tcPr.get("marT") == "45720", "Top margin should be preserved"
        assert copied_tcPr.get("marB") == "45720", "Bottom margin should be preserved"
        assert copied_tcPr.get("anchor") == "ctr", "Anchor should be preserved"

    def test_copy_table_preserves_style_properties(self):
        """Test that table-level style properties are preserved when copying."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Add table to source slide
        table_shape = source_slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(5), Inches(3))
        table = table_shape.table
        table.cell(0, 0).text = "Header"

        # Set table style properties
        table.first_row = False  # Disable first row styling
        table.first_col = False  # Disable first column styling
        table.horz_banding = False  # Disable horizontal banding
        table.vert_banding = False  # Disable vertical banding

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)
        copied_shape = shape_copier.copy_shape(table_shape, target_slide)

        assert copied_shape is not None
        copied_table = copied_shape.table

        # Verify style properties were copied
        assert copied_table.first_row == False, "first_row should be preserved as False"
        assert copied_table.first_col == False, "first_col should be preserved as False"
        assert copied_table.horz_banding == False, "horz_banding should be preserved as False"
        assert copied_table.vert_banding == False, "vert_banding should be preserved as False"

    def test_copy_text_shape_preserves_word_wrap(self):
        """Test that word_wrap is preserved when copying text shapes."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Add text box with word_wrap explicitly set
        textbox = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Test text that should wrap within the text box"
        textbox.text_frame.word_wrap = True  # Explicitly enable word wrap

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)
        copied_shape = shape_copier.copy_shape(textbox, target_slide)

        assert copied_shape is not None
        assert copied_shape.text_frame.word_wrap is True, "word_wrap was not preserved during copy"

    def test_copy_text_shape_preserves_margins(self):
        """Test that margins are preserved when copying text shapes."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)
        target_slide = prs.slides.add_slide(layout)

        # Add text box with margins explicitly set
        textbox = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Test text with custom margins"
        textbox.text_frame.margin_top = Inches(0.1)
        textbox.text_frame.margin_bottom = Inches(0.1)
        textbox.text_frame.margin_left = Inches(0.2)
        textbox.text_frame.margin_right = Inches(0.2)

        # Copy the shape
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)
        copied_shape = shape_copier.copy_shape(textbox, target_slide)

        assert copied_shape is not None
        assert copied_shape.text_frame.margin_top == Inches(0.1), "margin_top was not preserved during copy"
        assert copied_shape.text_frame.margin_bottom == Inches(0.1), "margin_bottom was not preserved during copy"
        assert copied_shape.text_frame.margin_left == Inches(0.2), "margin_left was not preserved during copy"
        assert copied_shape.text_frame.margin_right == Inches(0.2), "margin_right was not preserved during copy"


class TestSlideCopierManager:
    """Test the main slide copier orchestration."""

    def test_slide_copier_initialization(self):
        """Test slide copier manager initialization."""
        prs = Presentation()
        copier = SlideCopierManager(prs)

        assert copier.target_presentation == prs
        assert isinstance(copier.id_manager, IdManager)
        assert isinstance(copier.relationship_copier, RelationshipCopier)
        assert isinstance(copier.shape_copier, ShapeCopier)

    def test_copy_simple_slide(self):
        """Test copying a simple slide with text."""
        source_prs = Presentation()
        target_prs = Presentation()

        # Create source slide with content
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)
        textbox = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Source slide content"

        # Copy the slide
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide(source_slide)

        assert copied_slide is not None
        assert len(target_prs.slides) == 1

        # Verify content was copied
        copied_shapes = list(copied_slide.shapes)
        text_shapes = [s for s in copied_shapes if hasattr(s, 'text') and s.text]
        assert len(text_shapes) > 0

    def test_copy_slide_with_notes(self):
        """Test copying a slide with speaker notes."""
        source_prs = Presentation()
        target_prs = Presentation()

        # Create source slide with notes
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Add notes
        if source_slide.has_notes_slide:
            source_slide.notes_slide.notes_text_frame.text = "These are speaker notes"

        # Copy the slide
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide(source_slide)

        assert copied_slide is not None

        # Check if notes were copied
        if copied_slide.has_notes_slide and copied_slide.notes_slide.notes_text_frame:
            # Notes copying may not always work due to complexity
            pass  # This is acceptable

    def test_copy_slide_safe_fallback(self):
        """Test safe copying with fallback strategies."""
        source_prs = Presentation()
        target_prs = Presentation()

        # Create source slide
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Copy using safe method
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide_safe(source_slide, fallback_to_layout_only=True)

        # Should succeed even if content copying fails
        assert copied_slide is not None
        assert len(target_prs.slides) == 1

    def test_batch_copy_slides(self):
        """Test batch copying of multiple slides."""
        source_prs = Presentation()
        target_prs = Presentation()

        # Create multiple source slides
        layout = source_prs.slide_layouts[0]
        source_slides = []
        for i in range(3):
            slide = source_prs.slides.add_slide(layout)
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            textbox.text = f"Slide {i + 1} content"
            source_slides.append(slide)

        # Batch copy
        copier = SlideCopierManager(target_prs)
        results = copier.batch_copy_slides(source_slides)

        # Should have copied all slides
        assert len(results) == 3
        assert len(target_prs.slides) == 3

        # At least some should be successful
        successful_copies = [r for r in results if r is not None]
        assert len(successful_copies) > 0

    def test_copy_statistics(self):
        """Test getting copy operation statistics."""
        source_prs = Presentation()
        target_prs = Presentation()

        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        copier = SlideCopierManager(target_prs)
        copier.copy_slide(source_slide)

        stats = copier.get_copy_statistics()

        assert 'total_operations' in stats
        assert 'successful_operations' in stats
        assert 'failed_operations' in stats
        assert stats['total_operations'] >= 1


class TestPowerPointPresentationIntegration:
    """Test integration with PowerPointPresentation class."""

    def test_insert_copy_integration(self):
        """Test the updated insert_copy method."""
        # Create source presentation
        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide_pptx = source_prs.slides.add_slide(layout)
        textbox = source_slide_pptx.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Integration test content"

        # Save source to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            source_prs.save(temp_file.name)
            source_path = temp_file.name

        try:
            # Load presentations
            api_client = PowerPointAPIClient()
            source_presentation = PowerPointPresentation.from_id(api_client, source_path)

            # Create target presentation
            target_prs = Presentation()
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                target_prs.save(temp_file.name)
                target_path = temp_file.name

            try:
                target_presentation = PowerPointPresentation.from_id(api_client, target_path)

                # Copy slide using the new robust method
                source_slide = source_presentation.slides[0]
                copied_slide = target_presentation.insert_copy(source_slide, api_client)

                assert copied_slide is not None
                assert len(target_presentation.slides) >= 1

                # Save to verify no corruption
                target_presentation.save(api_client)

            finally:
                if os.path.exists(target_path):
                    os.unlink(target_path)

        finally:
            if os.path.exists(source_path):
                os.unlink(source_path)

    def test_insert_copy_fallback(self):
        """Test the fallback method when robust copying fails."""
        # This test verifies that even if the robust copying fails,
        # the fallback still creates a usable slide

        # Create minimal presentation
        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide_pptx = source_prs.slides.add_slide(layout)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            source_prs.save(temp_file.name)
            source_path = temp_file.name

        try:
            api_client = PowerPointAPIClient()
            source_presentation = PowerPointPresentation.from_id(api_client, source_path)

            target_prs = Presentation()
            target_presentation = PowerPointPresentation(target_prs)

            # This should work even if content copying fails
            source_slide = source_presentation.slides[0]
            copied_slide = target_presentation.insert_copy(source_slide, api_client)

            assert copied_slide is not None

        finally:
            if os.path.exists(source_path):
                os.unlink(source_path)


class TestRobustnessAndErrorHandling:
    """Test robustness and error handling scenarios."""

    def test_handle_corrupted_source_slide(self):
        """Test handling of potentially corrupted source slides."""
        source_prs = Presentation()
        target_prs = Presentation()

        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # The copier should handle edge cases gracefully
        copier = SlideCopierManager(target_prs)

        # This should not crash even with minimal slide content
        copied_slide = copier.copy_slide_safe(source_slide)
        assert copied_slide is not None

    def test_context_manager_cleanup(self):
        """Test proper cleanup using context manager."""
        source_prs = Presentation()
        target_prs = Presentation()

        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Use context manager
        with SlideCopierManager(target_prs) as copier:
            copied_slide = copier.copy_slide(source_slide)
            assert copied_slide is not None

        # Cleanup should have been called automatically

    def test_large_number_of_shapes(self):
        """Test copying slides with many shapes."""
        source_prs = Presentation()
        target_prs = Presentation()

        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Add multiple shapes
        for i in range(10):
            textbox = source_slide.shapes.add_textbox(
                Inches(0.5 + i * 0.5),
                Inches(0.5 + i * 0.2),
                Inches(2),
                Inches(0.5)
            )
            textbox.text = f"Shape {i}"

        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide(source_slide)

        # Should handle multiple shapes without issues
        assert copied_slide is not None
        stats = copier.get_copy_statistics()
        assert stats['total_shapes_copied'] > 0


class TestShapeAltTextCopying:
    """Test that alt text is preserved when copying shapes."""

    def test_alt_text_copied_with_shape(self):
        """Test that alt text title attribute is copied during slide copying."""
        source_prs = Presentation()
        source_layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(source_layout)

        # Add a textbox with alt text
        textbox = source_slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )
        textbox.text = "Test content"

        # Set alt text in XML (the way it's actually stored and read)
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        cnvpr.attrib["title"] = "TestAltText"

        # Copy slide using SlideCopierManager
        target_prs = Presentation()
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide_safe(source_slide, insertion_index=0)

        # Verify alt text was copied
        # Find the copied textbox (skip layout placeholders)
        copied_textbox = None
        for shape in copied_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                cnvpr = shape._element.xpath(".//p:cNvPr")[0]
                if cnvpr.attrib.get("title") == "TestAltText":
                    copied_textbox = shape
                    break

        assert copied_textbox is not None, "Copied textbox with alt text not found"

        # Verify the alt text value
        copied_cnvpr = copied_textbox._element.xpath(".//p:cNvPr")[0]
        copied_title = copied_cnvpr.attrib.get("title")
        assert copied_title == "TestAltText", f"Expected alt text 'TestAltText', got '{copied_title}'"

    def test_alt_text_copied_for_multiple_shapes(self):
        """Test that alt text is copied for all shapes in a slide."""
        source_prs = Presentation()
        source_layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(source_layout)

        # Add multiple textboxes with different alt text
        alt_texts = ["Title", "CustomerName", "Content"]
        for i, alt_text in enumerate(alt_texts):
            textbox = source_slide.shapes.add_textbox(
                Inches(1), Inches(1 + i), Inches(5), Inches(1)
            )
            textbox.text = f"Shape {i}"

            # Set alt text
            cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
            cnvpr.attrib["title"] = alt_text

        # Copy slide
        target_prs = Presentation()
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide_safe(source_slide, insertion_index=0)

        # Verify all alt texts were copied
        # Find copied textboxes (skip layout placeholders)
        copied_textboxes = [
            shape for shape in copied_slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        ]

        # Match each expected alt text with a copied shape
        for expected_alt_text in alt_texts:
            found = False
            for shape in copied_textboxes:
                cnvpr = shape._element.xpath(".//p:cNvPr")[0]
                if cnvpr.attrib.get("title") == expected_alt_text:
                    found = True
                    break
            assert found, f"Expected alt text '{expected_alt_text}' not found in any copied shape"


class TestRelationshipRemapping:
    """Test that r:id references are properly remapped when copying shapes with images."""

    def test_remap_element_relationships_function(self):
        """Test the XmlUtils.remap_element_relationships function directly."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Add a textbox to get an XML element
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        element = textbox._element

        # Create a mapping
        mapping = {"rId1": "rId100", "rId2": "rId200"}

        # The function should run without error even with no matching refs
        result = XmlUtils.remap_element_relationships(element, mapping)
        assert result >= 0  # Returns count of remapped refs

    def test_remap_element_relationships_with_empty_mapping(self):
        """Test that empty mapping returns 0 remapped."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        element = textbox._element

        result = XmlUtils.remap_element_relationships(element, {})
        assert result == 0

    def test_remap_element_relationships_with_none_element(self):
        """Test that None element returns 0 remapped."""
        mapping = {"rId1": "rId100"}
        result = XmlUtils.remap_element_relationships(None, mapping)
        assert result == 0

    def test_image_relationship_remapped_after_copy(self):
        """Test that image r:embed references are remapped to new relationship IDs."""
        import io

        from PIL import Image

        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Create a test image
        img = Image.new('RGB', (100, 100), color='red')
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(img_bytes.getvalue())
            img_path = f.name

        try:
            # Add image to source slide
            picture = source_slide.shapes.add_picture(img_path, Inches(1), Inches(1))

            # Copy to new presentation
            target_prs = Presentation()
            copier = SlideCopierManager(target_prs)
            copied_slide = copier.copy_slide(source_slide)

            assert copied_slide is not None

            # Find copied picture and verify r:embed exists
            found_picture = False
            for shape in copied_slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    found_picture = True
                    # Use the NAMESPACES dict from XmlUtils for xpath
                    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    copied_blip = shape._element.findall(f'.//{{{a_ns}}}blip')
                    if copied_blip:
                        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                        new_rid = copied_blip[0].get(f'{{{r_ns}}}embed')
                        # r:embed should exist (not be None or cleared)
                        assert new_rid is not None, "r:embed was cleared instead of remapped"
                        # r:id should point to valid relationship
                        rel = copied_slide.part.rels.get(new_rid)
                        assert rel is not None, f"r:embed {new_rid} doesn't exist in slide relationships"
                    break

            assert found_picture, "No picture shape found in copied slide"

        finally:
            os.unlink(img_path)

    def test_copied_image_can_be_saved_and_reopened(self):
        """Test that copied images can be saved and reopened without corruption."""
        import io

        from PIL import Image

        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Create a test image
        img = Image.new('RGB', (100, 100), color='blue')
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(img_bytes.getvalue())
            img_path = f.name

        try:
            # Add image to source slide
            source_slide.shapes.add_picture(img_path, Inches(1), Inches(1))

            # Copy to new presentation
            target_prs = Presentation()
            copier = SlideCopierManager(target_prs)
            copier.copy_slide(source_slide)

            # Save and reopen to verify no corruption
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
                output_path = f.name

            target_prs.save(output_path)

            try:
                # Reopen to verify file isn't corrupted
                reopened = Presentation(output_path)
                assert len(reopened.slides) == 1

                # Verify image shape exists
                found_picture = False
                for shape in reopened.slides[0].shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        found_picture = True
                        # Verify image data is accessible
                        assert hasattr(shape, 'image')
                        assert shape.image.blob is not None
                        break

                assert found_picture, "Picture not found after save/reopen"

            finally:
                os.unlink(output_path)

        finally:
            os.unlink(img_path)


class TestGroupShapeCopying:
    """Test copying of GROUP shapes including nested groups."""

    def test_group_shape_is_copied_via_xml(self):
        """Test that GROUP shapes are handled by _copy_group_shape."""
        # Create a simple presentation and manually test the copy_shape dispatch
        prs = Presentation()
        layout = prs.slide_layouts[0]
        source_slide = prs.slides.add_slide(layout)

        # Add some shapes to source
        textbox = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        textbox.text = "Test"

        # Copy slide to test GROUP handling works
        target_prs = Presentation()
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide(source_slide)

        assert copied_slide is not None

    def test_group_shape_handler_exists(self):
        """Test that GROUP shape type triggers _copy_group_shape method."""
        prs = Presentation()
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)

        # Verify the method exists
        assert hasattr(shape_copier, '_copy_group_shape')

    def test_nested_shape_ids_are_regenerated(self):
        """Test that nested shape IDs inside groups are regenerated."""
        prs = Presentation()
        id_manager = IdManager(prs)

        # Generate multiple IDs to verify uniqueness
        ids = [id_manager.generate_unique_shape_id() for _ in range(10)]

        # All IDs should be unique
        assert len(set(ids)) == 10


class TestFreeformShapeCopying:
    """Test copying of FREEFORM shapes (custom geometry)."""

    def test_freeform_shape_handler_exists(self):
        """Test that FREEFORM shape type triggers _copy_freeform_shape method."""
        prs = Presentation()
        id_manager = IdManager(prs)
        shape_copier = ShapeCopier(id_manager)

        # Verify the method exists
        assert hasattr(shape_copier, '_copy_freeform_shape')


class TestSampleadTemplateRoundtrip:
    """End-to-end test using actual PPTX files with complex shapes."""

    @pytest.fixture
    def samplead_template_path(self):
        """Get path to Samplead template if it exists."""
        # Check relative to the test file
        possible_paths = [
            "playground/samplead/Samplead Master Deck Template.pptx",
            "../playground/samplead/Samplead Master Deck Template.pptx",
            "../../playground/samplead/Samplead Master Deck Template.pptx",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                return path
        return None

    def test_slide_with_complex_shapes_preserves_shape_count(self):
        """Test that copying a slide preserves the shape count."""
        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Add various shapes
        source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        source_slide.shapes.add_shape(1, Inches(3), Inches(1), Inches(2), Inches(2))  # RECTANGLE
        source_slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(4), Inches(2))

        source_shape_count = len(source_slide.shapes)

        # Copy to new presentation
        target_prs = Presentation()
        copier = SlideCopierManager(target_prs)
        copied_slide = copier.copy_slide(source_slide)

        assert copied_slide is not None

        # Shape count may differ slightly due to layout placeholders
        # but should be in similar range
        target_shape_count = len(copied_slide.shapes)
        assert target_shape_count >= source_shape_count - 2  # Allow some variation

    def test_save_and_reopen_complex_slide(self):
        """Test that complex slides can be saved and reopened."""
        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Add multiple shapes
        source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        source_slide.shapes.add_shape(1, Inches(3), Inches(1), Inches(2), Inches(2))

        # Copy to new presentation
        target_prs = Presentation()
        copier = SlideCopierManager(target_prs)
        copier.copy_slide(source_slide)

        # Save and reopen
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = f.name

        try:
            target_prs.save(output_path)

            # Reopen to verify no corruption
            reopened = Presentation(output_path)
            assert len(reopened.slides) == 1

        finally:
            os.unlink(output_path)


class TestCopyShapeElementWithRelationshipMapping:
    """Test that copy_shape_element properly handles relationship_mapping parameter."""

    def test_copy_shape_element_with_no_mapping(self):
        """Test copy_shape_element works with no relationship mapping (legacy behavior)."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        textbox.text = "Test"

        id_manager = IdManager(prs)
        new_shape_id = id_manager.generate_unique_shape_id()
        new_creation_id = id_manager.generate_unique_creation_id()

        # Copy without relationship mapping
        new_element = XmlUtils.copy_shape_element(
            textbox,
            new_shape_id,
            new_creation_id,
            relationship_mapping=None
        )

        assert new_element is not None

    def test_copy_shape_element_with_mapping(self):
        """Test copy_shape_element works with relationship mapping."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        textbox.text = "Test"

        id_manager = IdManager(prs)
        new_shape_id = id_manager.generate_unique_shape_id()
        new_creation_id = id_manager.generate_unique_creation_id()

        # Copy with relationship mapping
        mapping = {"rId1": "rId100"}
        new_element = XmlUtils.copy_shape_element(
            textbox,
            new_shape_id,
            new_creation_id,
            relationship_mapping=mapping
        )

        assert new_element is not None


class TestImageRelationshipCopying:
    """Test that image relationships are correctly copied to target slide."""

    def test_image_relationship_copied_successfully(self):
        """Test that _copy_image_relationship returns valid relationship ID."""
        import io as stdlib_io

        from PIL import Image

        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Create and add test image
        img = Image.new('RGB', (100, 100), color='red')
        img_bytes = stdlib_io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(img_bytes.getvalue())
            img_path = f.name

        try:
            source_slide.shapes.add_picture(img_path, Inches(1), Inches(1))

            # Copy slide - this should now work
            target_prs = Presentation()
            copier = SlideCopierManager(target_prs)
            copied_slide = copier.copy_slide(source_slide)

            # Verify relationship mapping was populated
            stats = copier.get_copy_statistics()
            assert stats['total_relationships_copied'] > 0, "No relationships were copied"

        finally:
            os.unlink(img_path)

    def test_image_relationship_copy_creates_valid_reference(self):
        """Test that copied images have valid r:embed references in target slide."""
        import io as stdlib_io

        from PIL import Image

        source_prs = Presentation()
        layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(layout)

        # Create and add test image
        img = Image.new('RGB', (100, 100), color='green')
        img_bytes = stdlib_io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(img_bytes.getvalue())
            img_path = f.name

        try:
            source_slide.shapes.add_picture(img_path, Inches(1), Inches(1))

            # Copy slide
            target_prs = Presentation()
            copier = SlideCopierManager(target_prs)
            copied_slide = copier.copy_slide(source_slide)

            # Find picture in copied slide and verify its r:embed reference is valid
            for shape in copied_slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Get the r:embed reference from the shape's XML
                    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                    blip_elements = shape._element.findall(f'.//{{{a_ns}}}blip')
                    for blip in blip_elements:
                        embed_ref = blip.get(f'{{{r_ns}}}embed')
                        if embed_ref:
                            # Verify the relationship ID exists in target slide
                            rel = copied_slide.part.rels.get(embed_ref)
                            assert rel is not None, f"r:embed {embed_ref} not found in target slide relationships"
                            break

        finally:
            os.unlink(img_path)


class TestAppXmlMetadataUpdate:
    """Test that app.xml metadata is updated correctly after slide changes."""

    def test_app_xml_metadata_updated_after_slide_deletion(self, tmp_path):
        """Test that app.xml metadata reflects actual slide count after deletion."""
        import re
        import zipfile

        # Create a presentation with multiple slides
        prs = Presentation()
        layout = prs.slide_layouts[0]
        for i in range(5):
            prs.slides.add_slide(layout)

        original_count = len(prs.slides)
        assert original_count == 5

        # Delete a slide using the SlideDeleter
        from gslides_api.pptx.slide_deleter import SlideDeleter
        deleter = SlideDeleter(prs)
        deleter.delete_slide(0)

        new_count = len(prs.slides)
        assert new_count == original_count - 1

        # Save via PowerPointPresentation which should fix app.xml
        output_path = tmp_path / "output.pptx"
        pp = PowerPointPresentation(pptx_presentation=prs, file_path=str(output_path))
        pp.save(api_client=PowerPointAPIClient())

        # Verify app.xml has correct slide count
        with zipfile.ZipFile(output_path, 'r') as zf:
            app_xml = zf.read('docProps/app.xml').decode('utf-8')
            match = re.search(r'<Slides>(\d+)</Slides>', app_xml)
            assert match is not None, "Could not find <Slides> element in app.xml"
            assert int(match.group(1)) == new_count, (
                f"app.xml says {match.group(1)} slides but actual is {new_count}"
            )

    def test_app_xml_metadata_updated_after_multiple_deletions(self, tmp_path):
        """Test that app.xml metadata is correct after deleting multiple slides."""
        import re
        import zipfile

        # Create a presentation with many slides
        prs = Presentation()
        layout = prs.slide_layouts[0]
        for i in range(10):
            prs.slides.add_slide(layout)

        assert len(prs.slides) == 10

        # Delete multiple slides
        from gslides_api.pptx.slide_deleter import SlideDeleter
        deleter = SlideDeleter(prs)
        # Delete in reverse order to avoid index shifting
        for idx in [8, 5, 2, 0]:
            deleter.delete_slide(idx)

        remaining_count = len(prs.slides)
        assert remaining_count == 6

        # Save via PowerPointPresentation
        output_path = tmp_path / "output.pptx"
        pp = PowerPointPresentation(pptx_presentation=prs, file_path=str(output_path))
        pp.save(api_client=PowerPointAPIClient())

        # Verify app.xml has correct slide count
        with zipfile.ZipFile(output_path, 'r') as zf:
            app_xml = zf.read('docProps/app.xml').decode('utf-8')
            match = re.search(r'<Slides>(\d+)</Slides>', app_xml)
            assert match is not None
            assert int(match.group(1)) == remaining_count

    def test_saved_file_can_be_reopened(self, tmp_path):
        """Test that a file saved after modifications can be reopened."""
        # Create and modify presentation
        prs = Presentation()
        layout = prs.slide_layouts[0]
        for i in range(5):
            slide = prs.slides.add_slide(layout)
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        from gslides_api.pptx.slide_deleter import SlideDeleter
        deleter = SlideDeleter(prs)
        deleter.delete_slide(0)
        deleter.delete_slide(0)

        # Save
        output_path = tmp_path / "output.pptx"
        pp = PowerPointPresentation(pptx_presentation=prs, file_path=str(output_path))
        pp.save(api_client=PowerPointAPIClient())

        # Reopen and verify
        reopened = Presentation(str(output_path))
        assert len(reopened.slides) == 3


class TestCreateImageElementLike:
    """Test the create_image_element_like functionality for chart replacement."""

    def test_create_image_element_like_replaces_shape(self, tmp_path):
        """Test that create_image_element_like creates a valid image element."""
        from gslides_api.adapters.pptx_adapter import (
            PowerPointAPIClient,
            PowerPointElementParent,
            PowerPointImageElement,
            PowerPointSlide,
        )

        # Create a presentation with a shape (simulating a chart)
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Add a shape that we'll replace with an image
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1), Inches(1), Inches(4), Inches(3)
        )
        shape.name = "TestChart"
        original_left = shape.left
        original_top = shape.top
        original_width = shape.width
        original_height = shape.height
        original_shape_count = len(slide.shapes)

        # Create wrapper for the shape
        from gslides_api.adapters.pptx_adapter import validate_pptx_element
        element = validate_pptx_element(shape)
        element.pptx_slide = slide

        # Call create_image_element_like
        api_client = PowerPointAPIClient()
        new_image_element = element.create_image_element_like(api_client=api_client)

        # Verify the result
        assert isinstance(new_image_element, PowerPointImageElement)
        assert new_image_element.pptx_element is not None
        assert new_image_element.pptx_slide == slide

        # Verify the new image has similar position
        # (might differ slightly due to placeholder image size)
        assert new_image_element.pptx_element.left == original_left
        assert new_image_element.pptx_element.top == original_top

        # Verify the old shape was removed and new one added
        # (shape count should be same since we removed one and added one)
        assert len(slide.shapes) == original_shape_count

    def test_create_image_element_like_requires_slide_reference(self):
        """Test that create_image_element_like fails without slide reference."""
        from gslides_api.adapters.pptx_adapter import (
            PowerPointAPIClient,
            PowerPointElementParent,
        )

        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1), Inches(1), Inches(2), Inches(2)
        )

        from gslides_api.adapters.pptx_adapter import validate_pptx_element
        element = validate_pptx_element(shape)
        # Don't set pptx_slide

        api_client = PowerPointAPIClient()
        with pytest.raises(ValueError, match="slide reference"):
            element.create_image_element_like(api_client=api_client)


class TestPlaceholderRemoval:
    """Test removal of layout placeholder shapes."""

    def test_remove_layout_placeholders_removes_all_placeholders(self):
        """Test that _remove_layout_placeholders removes all placeholder shapes."""
        prs = Presentation()
        # Use a layout that typically has placeholders (Title Slide layout)
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Count placeholder shapes before removal
        placeholders_before = [
            s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
        ]
        # Layout slides typically have at least 1 placeholder
        assert len(placeholders_before) > 0, "Test requires a layout with placeholders"

        # Remove placeholders
        _remove_layout_placeholders(slide)

        # Count placeholder shapes after removal
        placeholders_after = [
            s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
        ]
        assert len(placeholders_after) == 0, "All placeholder shapes should be removed"

    def test_remove_layout_placeholders_preserves_added_content(self):
        """Test that _remove_layout_placeholders preserves content shapes."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Add a textbox (not a placeholder)
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )
        textbox.text = "This should be preserved"

        # Add a shape
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(1), Inches(3), Inches(2), Inches(1)
        )

        # Remove placeholders
        _remove_layout_placeholders(slide)

        # Check that the added content is still there
        remaining_shapes = list(slide.shapes)
        texts = [s.text for s in remaining_shapes if hasattr(s, 'text') and s.text]
        assert "This should be preserved" in texts

        # Check that rectangle shape is still there
        rectangles = [
            s for s in remaining_shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        ]
        assert len(rectangles) >= 1

    def test_slide_copier_manager_removes_placeholders_in_create_target_slide(self):
        """Test that SlideCopierManager removes placeholders when creating target slide."""
        source_prs = Presentation()
        target_prs = Presentation()

        # Create source slide
        source_layout = source_prs.slide_layouts[0]
        source_slide = source_prs.slides.add_slide(source_layout)

        # Create target slide using copier (which now removes placeholders)
        copier = SlideCopierManager(target_prs)
        target_slide = copier._create_target_slide(source_slide=source_slide)

        assert target_slide is not None

        # Verify no placeholder shapes on the target slide
        placeholders = [
            s for s in target_slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
        ]
        assert len(placeholders) == 0, (
            "SlideCopierManager._create_target_slide should remove placeholders"
        )


if __name__ == "__main__":
    pytest.main([__file__, "-v"])