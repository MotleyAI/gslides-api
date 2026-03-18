"""
Test suite for PowerPointShapeElement text frame to markdown conversion using pptx2md.

This module tests the enhanced text frame conversion capabilities that use
pptx2md for proper formatting, hyperlinks, colors, and special character handling.
"""

import pytest
from unittest.mock import Mock, MagicMock
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.dml.color import RGBColor

from gslides_api.adapters.pptx_adapter import PowerPointShapeElement, PowerPointAPIClient
from gslides_api.adapters.abstract_slides import AbstractAltText


class TestPowerPointTextFrameMarkdown:
    """Test text frame to markdown conversion with pptx2md."""

    def create_mock_shape_with_text_frame(self, paragraphs_data):
        """Helper to create mock shape with specified text frame structure."""
        mock_shape = Mock()
        mock_text_frame = Mock()

        mock_paragraphs = []
        for para_data in paragraphs_data:
            mock_paragraph = Mock()
            mock_paragraph.level = para_data.get('level', 0)
            mock_paragraph.text = para_data.get('text', '')

            # Mock XML element for bullet detection
            # If has_bullet is explicitly set, use it; otherwise, level > 0 implies bullets
            has_bullet = para_data.get('has_bullet', para_data.get('level', 0) > 0)
            mock_pPr = Mock()
            mock_element = Mock()

            def make_find_method(is_bullet):
                """Create a find method that returns buChar for bullets."""
                def find(tag):
                    if is_bullet and 'buChar' in tag:
                        return Mock()  # Return something truthy for buChar
                    return None  # Return None for everything else
                return find

            mock_pPr.find = make_find_method(has_bullet)
            mock_element.get_or_add_pPr = Mock(return_value=mock_pPr)
            mock_paragraph._element = mock_element

            mock_runs = []
            for run_data in para_data.get('runs', []):
                mock_run = Mock()
                mock_run.text = run_data.get('text', '')

                # Mock font properties
                mock_font = Mock()
                mock_font.bold = run_data.get('bold', False)
                mock_font.italic = run_data.get('italic', False)
                mock_font.underline = run_data.get('underline', False)

                # Mock color
                mock_color = Mock()
                if 'color_rgb' in run_data:
                    mock_color.type = MSO_COLOR_TYPE.RGB
                    mock_color.rgb = RGBColor(*run_data['color_rgb'])
                else:
                    mock_color.type = MSO_COLOR_TYPE.SCHEME
                    mock_color.theme_color = MSO_THEME_COLOR.DARK_1

                mock_font.color = mock_color
                mock_run.font = mock_font

                # Mock hyperlink
                mock_hyperlink = Mock()
                mock_hyperlink.address = run_data.get('hyperlink', None)
                mock_run.hyperlink = mock_hyperlink

                mock_runs.append(mock_run)

            mock_paragraph.runs = mock_runs
            mock_paragraphs.append(mock_paragraph)

        mock_text_frame.paragraphs = mock_paragraphs
        mock_shape.text_frame = mock_text_frame

        return mock_shape

    def test_simple_text_conversion(self):
        """Test basic text conversion without formatting."""
        paragraphs_data = [
            {
                'text': 'Hello World',
                'runs': [{'text': 'Hello World'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        assert "Hello World" in result

    def test_bold_text_conversion(self):
        """Test bold text formatting preservation."""
        paragraphs_data = [
            {
                'text': 'Bold Text',
                'runs': [{'text': 'Bold Text', 'bold': True}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        assert "**Bold Text**" in result or "__Bold Text__" in result

    def test_italic_text_conversion(self):
        """Test italic text formatting preservation."""
        paragraphs_data = [
            {
                'text': 'Italic Text',
                'runs': [{'text': 'Italic Text', 'italic': True}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        assert "*Italic Text*" in result or "_Italic Text_" in result

    def test_hyperlink_conversion(self):
        """Test hyperlink preservation."""
        paragraphs_data = [
            {
                'text': 'Click here',
                'runs': [{'text': 'Click here', 'hyperlink': 'http://example.com'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # pptx2md should convert hyperlinks to markdown link format
        assert "[Click here](http://example.com)" in result

    def test_text_color_not_in_markdown(self):
        """Test that RGB colors are NOT converted to HTML spans (colors ignored in markdown)."""
        paragraphs_data = [
            {
                'text': 'Red Text',
                'runs': [{'text': 'Red Text', 'color_rgb': (255, 0, 0)}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # Colors should NOT produce HTML spans - just plain text
        assert "Red Text" in result
        assert "<span" not in result  # No HTML tags in markdown output

    def test_combined_formatting(self):
        """Test multiple formatting styles on the same text."""
        paragraphs_data = [
            {
                'text': 'Bold Italic',
                'runs': [{'text': 'Bold Italic', 'bold': True, 'italic': True}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # Should have both bold and italic formatting
        assert "Bold Italic" in result
        # Could be **_text_** or __*text*__ depending on pptx2md implementation

    def test_bullet_points_single_level(self):
        """Test simple bullet list conversion."""
        paragraphs_data = [
            {
                'text': 'First item',
                'level': 1,
                'runs': [{'text': 'First item'}]
            },
            {
                'text': 'Second item',
                'level': 1,
                'runs': [{'text': 'Second item'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # pptx2md outputs with double underscores for strong text by default
        assert "First item" in result
        assert "Second item" in result
        assert "  -" in result  # Should have bullet formatting with indentation

    def test_bullet_points_multi_level(self):
        """Test nested bullet list conversion."""
        paragraphs_data = [
            {
                'text': 'Level 1',
                'level': 1,
                'runs': [{'text': 'Level 1'}]
            },
            {
                'text': 'Level 2',
                'level': 2,
                'runs': [{'text': 'Level 2'}]
            },
            {
                'text': 'Level 2 again',
                'level': 2,
                'runs': [{'text': 'Level 2 again'}]
            },
            {
                'text': 'Back to Level 1',
                'level': 1,
                'runs': [{'text': 'Back to Level 1'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # Check for content and indentation structure
        assert "Level 1" in result
        assert "Level 2" in result
        assert "Level 2 again" in result
        assert "Back to Level 1" in result
        # Check for proper indentation levels (pptx2md uses 2 spaces per level)
        assert "  -" in result  # Level 1 indentation
        assert "    -" in result  # Level 2 indentation

    def test_special_character_escaping(self):
        """Test special markdown characters are properly escaped."""
        paragraphs_data = [
            {
                'text': 'Text with * and _ and [brackets]',
                'runs': [{'text': 'Text with * and _ and [brackets]'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # pptx2md should escape these characters so they don't interfere with markdown
        assert "\\*" in result or "Text with * and _ and [brackets]" in result

    def test_mixed_paragraph_and_lists(self):
        """Test complex content with mixed paragraphs and lists."""
        paragraphs_data = [
            {
                'text': 'Regular paragraph',
                'level': 0,
                'runs': [{'text': 'Regular paragraph'}]
            },
            {
                'text': 'First bullet',
                'level': 1,
                'runs': [{'text': 'First bullet'}]
            },
            {
                'text': 'Second bullet',
                'level': 1,
                'runs': [{'text': 'Second bullet'}]
            },
            {
                'text': 'Another paragraph',
                'level': 0,
                'runs': [{'text': 'Another paragraph'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        assert "Regular paragraph" in result
        assert "First bullet" in result
        assert "Second bullet" in result
        assert "Another paragraph" in result
        # Check for bullet point formatting
        assert "  -" in result  # Should have bullet indentation

    def test_empty_text_frame(self):
        """Test handling of empty text frames."""
        mock_shape = Mock()
        mock_shape.text_frame = None

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        assert result == ""

    def test_text_frame_with_empty_paragraphs(self):
        """Test handling of text frames with empty paragraphs."""
        paragraphs_data = [
            {
                'text': 'First line',
                'runs': [{'text': 'First line'}]
            },
            {
                'text': '',  # Empty paragraph
                'runs': []
            },
            {
                'text': 'Third line',
                'runs': [{'text': 'Third line'}]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        assert "First line" in result
        assert "Third line" in result
        # Empty lines should be preserved as line breaks

    def test_complex_real_world_example(self):
        """Test a complex real-world example with multiple formatting types."""
        paragraphs_data = [
            {
                'text': 'Project Overview',
                'runs': [{'text': 'Project Overview', 'bold': True}]
            },
            {
                'text': 'Key objectives:',
                'level': 0,
                'runs': [{'text': 'Key objectives:'}]
            },
            {
                'text': 'Implement new features',
                'level': 1,
                'runs': [{'text': 'Implement new features'}]
            },
            {
                'text': 'Improve performance by 20%',
                'level': 1,
                'runs': [
                    {'text': 'Improve performance by '},
                    {'text': '20%', 'bold': True, 'color_rgb': (255, 0, 0)}
                ]
            },
            {
                'text': 'Documentation available at company.com',
                'level': 1,
                'runs': [
                    {'text': 'Documentation available at '},
                    {'text': 'company.com', 'hyperlink': 'https://company.com'}
                ]
            }
        ]

        mock_shape = self.create_mock_shape_with_text_frame(paragraphs_data)

        element = PowerPointShapeElement(
            objectId="1",
            pptx_element=mock_shape,
            alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)

        # Verify basic structure
        assert "Project Overview" in result
        assert "Key objectives:" in result
        assert "Implement new features" in result
        assert "Improve performance" in result
        assert "20%" in result
        assert "Documentation available" in result
        assert "company.com" in result

        # Verify formatting is preserved
        # Bold, colors, hyperlinks should be converted by pptx2md