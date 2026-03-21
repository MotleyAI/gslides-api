"""
Tests for PowerPoint adapter alt text functionality.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from gslides_api.adapters.pptx_adapter import PowerPointAPIClient, validate_pptx_element


class TestPowerPointAltText:
    """Test alt text functionality in PowerPoint adapter."""

    def test_set_alt_text_sets_xml_attribute(self):
        """Test that set_alt_text sets the XML p:cNvPr/@title attribute."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Set alt text
        element.set_alt_text(api_client, title="TestTitle")

        # Verify XML attribute is set
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        assert cnvpr.attrib.get("title") == "TestTitle"

    def test_alt_text_reading_matches_setting(self):
        """Test that reading alt text returns what was set."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()

        # Set via set_alt_text
        element = validate_pptx_element(textbox)
        element.set_alt_text(api_client, title="TestAltText")

        # Re-read as new element (simulates reading from saved file)
        element2 = validate_pptx_element(textbox)
        assert element2.alt_text.title == "TestAltText"

    def test_set_alt_text_with_special_characters(self):
        """Test that alt text with special characters is handled correctly."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Test with special characters
        special_title = "Test & Title <with> \"special\" 'chars'"
        element.set_alt_text(api_client, title=special_title)

        # Verify it's preserved
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        assert cnvpr.attrib.get("title") == special_title

    def test_set_alt_text_multiple_times(self):
        """Test that alt text can be updated multiple times."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Set first time
        element.set_alt_text(api_client, title="FirstTitle")
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        assert cnvpr.attrib.get("title") == "FirstTitle"

        # Update
        element.set_alt_text(api_client, title="SecondTitle")
        assert cnvpr.attrib.get("title") == "SecondTitle"

    def test_set_alt_text_empty_string(self):
        """Test that empty string alt text is handled."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Set to empty string (should still set the attribute)
        element.set_alt_text(api_client, title="")

        # Verify empty string is set (not None)
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        title = cnvpr.attrib.get("title")
        assert title == "", f"Expected empty string, got {repr(title)}"

    def test_set_alt_text_description_sets_xml_attribute(self):
        """Test that set_alt_text sets the XML p:cNvPr/@descr attribute for description."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Set alt text description
        element.set_alt_text(api_client, description="This is the chart description")

        # Verify XML descr attribute is set
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        assert cnvpr.attrib.get("descr") == "This is the chart description"

    def test_alt_text_description_reading_matches_setting(self):
        """Test that reading alt text description returns what was set."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()

        # Set via set_alt_text
        element = validate_pptx_element(textbox)
        element.set_alt_text(api_client, description="Test Description")

        # Re-read as new element (simulates reading from saved file)
        element2 = validate_pptx_element(textbox)
        assert element2.alt_text.description == "Test Description"

    def test_set_alt_text_both_title_and_description(self):
        """Test that both title and description can be set together."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Set both title and description
        element.set_alt_text(
            api_client,
            title="Chart Title",
            description="Chart showing weekly sales by region",
        )

        # Verify both XML attributes are set
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        assert cnvpr.attrib.get("title") == "Chart Title"
        assert cnvpr.attrib.get("descr") == "Chart showing weekly sales by region"

        # Verify reading them back works
        element2 = validate_pptx_element(textbox)
        assert element2.alt_text.title == "Chart Title"
        assert element2.alt_text.description == "Chart showing weekly sales by region"

    def test_alt_text_description_only_without_title(self):
        """Test that description can be set without title."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        api_client = PowerPointAPIClient()
        element = validate_pptx_element(textbox)

        # Initially no alt text
        assert element.alt_text.title is None
        assert element.alt_text.description is None

        # Set only description
        element.set_alt_text(api_client, description="Only description set")

        # Verify only descr is set, title remains unset
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        assert cnvpr.attrib.get("title") is None
        assert cnvpr.attrib.get("descr") == "Only description set"

        # Verify alt_text model is updated
        assert element.alt_text.description == "Only description set"

    def test_read_pre_existing_description_from_xml(self):
        """Test reading description that was already set in XML (simulating loading from file)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )

        # Manually set the descr attribute in XML (as if loaded from file)
        cnvpr = textbox._element.xpath(".//p:cNvPr")[0]
        cnvpr.attrib["title"] = "Pre-existing Title"
        cnvpr.attrib["descr"] = "Pre-existing Description from PowerPoint"

        # Now read via validate_pptx_element
        element = validate_pptx_element(textbox)

        # Both should be read correctly
        assert element.alt_text.title == "Pre-existing Title"
        assert element.alt_text.description == "Pre-existing Description from PowerPoint"


class TestPresentationIdPropagation:
    """Test that presentation_id is properly propagated to slides and elements."""

    def test_presentation_id_propagates_to_slides_on_init(self):
        """Test that creating PowerPointPresentation propagates presentation_id to slides."""
        from gslides_api.adapters.pptx_adapter import PowerPointPresentation, PowerPointSlide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Add speaker notes (required by PowerPointSlide)
        slide.notes_slide.notes_text_frame.text = "Test Slide"

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text_frame.text = "Test Text"

        # Create PowerPointPresentation with a file path
        presentation = PowerPointPresentation(prs, file_path="/test/path/presentation.pptx")

        # Verify presentation_id is set on all slides
        for ppt_slide in presentation.slides:
            assert ppt_slide.presentation_id == "/test/path/presentation.pptx"
            # Verify presentation_id is propagated to all elements
            for element in ppt_slide.elements:
                assert element.presentation_id == "/test/path/presentation.pptx"

    def test_setting_slide_presentation_id_propagates_to_elements(self):
        """Test that setting presentation_id on a slide propagates to its elements."""
        from gslides_api.adapters.pptx_adapter import PowerPointSlide

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Test Slide"

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text_frame.text = "Test Text"

        ppt_slide = PowerPointSlide(slide)

        # Elements should start with empty presentation_id
        for element in ppt_slide.elements:
            assert element.presentation_id == ""

        # Set presentation_id on slide
        ppt_slide.presentation_id = "/new/path/test.pptx"

        # Verify all elements now have the presentation_id
        for element in ppt_slide.elements:
            assert element.presentation_id == "/new/path/test.pptx"

    def test_insert_copy_sets_presentation_id(self, tmp_path):
        """Test that insert_copy sets presentation_id on the copied slide."""
        from gslides_api.adapters.pptx_adapter import (
            PowerPointAPIClient,
            PowerPointPresentation,
            PowerPointSlide,
        )

        # Create source presentation
        src_prs = Presentation()
        src_slide = src_prs.slides.add_slide(src_prs.slide_layouts[0])
        src_slide.notes_slide.notes_text_frame.text = "Source Slide"
        src_textbox = src_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        src_textbox.text_frame.text = "Source Text"

        # Save source presentation
        src_path = str(tmp_path / "source.pptx")
        src_prs.save(src_path)

        # Create target presentation
        tgt_prs = Presentation()
        tgt_slide = tgt_prs.slides.add_slide(tgt_prs.slide_layouts[0])
        tgt_slide.notes_slide.notes_text_frame.text = "Target Slide"

        tgt_path = str(tmp_path / "target.pptx")
        tgt_prs.save(tgt_path)

        # Load presentations
        api_client = PowerPointAPIClient()
        source = PowerPointPresentation.from_id(api_client, src_path)
        target = PowerPointPresentation.from_id(api_client, tgt_path)

        # Insert a copy
        new_slide = target.insert_copy(
            source_slide=source.slides[0],
            api_client=api_client,
        )

        # Verify the new slide has the target presentation_id
        assert new_slide.presentation_id == tgt_path
        # Verify elements also have the correct presentation_id
        for element in new_slide.elements:
            assert element.presentation_id == tgt_path


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
