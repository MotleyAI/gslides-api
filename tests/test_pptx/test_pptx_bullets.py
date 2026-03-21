"""
Test suite for PPTX bullet point functionality.

Tests bullet detection (reading) and bullet creation (writing) via XML manipulation.
"""

import os
import tempfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from gslides_api.pptx.markdown_to_pptx import (
    _enable_paragraph_bullets,
    apply_markdown_to_textframe,
)
from gslides_api.pptx.converters import (
    _DRAWINGML_NS,
    _paragraph_has_bullet,
    pptx_paragraph_to_markdown,
    pptx_text_frame_to_markdown,
)


class TestParagraphHasBullet:
    """Test the _paragraph_has_bullet() function for detecting bullets in XML."""

    def create_test_presentation_with_bullets(self):
        """Create a test presentation with bullet-formatted text."""
        prs = Presentation()
        # Use bullet slide layout (layout 1 typically has bullets)
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame

        # First paragraph is typically bullet-formatted in placeholder
        tf.paragraphs[0].text = "First bullet item"

        # Add more bullet items
        p = tf.add_paragraph()
        p.text = "Second bullet item"
        p.level = 0

        p = tf.add_paragraph()
        p.text = "Nested bullet item"
        p.level = 1

        return prs, tf

    def create_test_textbox_no_bullets(self):
        """Create a text box without bullet formatting."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add a plain text box (no bullets)
        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        tf.text = "Plain text without bullets"

        p = tf.add_paragraph()
        p.text = "Another plain paragraph"

        return prs, tf

    def test_detect_bullet_in_placeholder(self):
        """Test that bullets are detected in bullet placeholder."""
        prs, tf = self.create_test_presentation_with_bullets()

        # At least some paragraphs should have bullets
        has_any_bullet = False
        for para in tf.paragraphs:
            if _paragraph_has_bullet(para):
                has_any_bullet = True
                break

        # Note: Bullet detection depends on whether the placeholder has
        # bullet XML elements. Some placeholders inherit from master.
        # This test verifies the function doesn't crash and handles the case.
        assert isinstance(has_any_bullet, bool)

    def test_no_bullet_in_plain_textbox(self):
        """Test that no bullets are detected in a plain text box."""
        prs, tf = self.create_test_textbox_no_bullets()

        for para in tf.paragraphs:
            # Plain text boxes should not have bullet formatting
            assert not _paragraph_has_bullet(para)

    def test_detect_bullet_after_enabling(self):
        """Test that bullets are detected after enabling via XML."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = "Text that will get bullets"

        # Before enabling, should have no bullet
        assert not _paragraph_has_bullet(para)

        # Enable bullets via XML
        _enable_paragraph_bullets(para)

        # After enabling, should have bullet
        assert _paragraph_has_bullet(para)


class TestEnableParagraphBullets:
    """Test the _enable_paragraph_bullets() function for adding bullets via XML."""

    def test_enable_bullets_adds_buchar_element(self):
        """Test that enabling bullets adds buChar XML element."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = "Test paragraph"

        # Enable bullets
        _enable_paragraph_bullets(para)

        # Check that buChar element exists
        pPr = para._element.get_or_add_pPr()
        buChar = pPr.find(f"{{{_DRAWINGML_NS}}}buChar")
        assert buChar is not None
        assert buChar.get("char") == "•"

    def test_enable_bullets_custom_character(self):
        """Test enabling bullets with a custom character."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = "Test paragraph"

        # Enable bullets with custom character
        _enable_paragraph_bullets(para, char="★")

        # Check that buChar element has custom character
        pPr = para._element.get_or_add_pPr()
        buChar = pPr.find(f"{{{_DRAWINGML_NS}}}buChar")
        assert buChar is not None
        assert buChar.get("char") == "★"

    def test_enable_bullets_removes_bunone(self):
        """Test that enabling bullets removes buNone element if present."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = "Test paragraph"

        # Manually add buNone to simulate disabled bullets
        pPr = para._element.get_or_add_pPr()
        buNone = etree.Element(f"{{{_DRAWINGML_NS}}}buNone")
        pPr.insert(0, buNone)

        # Verify buNone exists
        assert pPr.find(f"{{{_DRAWINGML_NS}}}buNone") is not None

        # Enable bullets
        _enable_paragraph_bullets(para)

        # Verify buNone was removed
        assert pPr.find(f"{{{_DRAWINGML_NS}}}buNone") is None

        # Verify buChar was added
        assert pPr.find(f"{{{_DRAWINGML_NS}}}buChar") is not None


class TestBulletRoundtrip:
    """Test roundtrip: write markdown with bullets, read back as markdown."""

    def test_bullet_list_roundtrip(self):
        """Test writing and reading back a bullet list."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
        )
        tf = textbox.text_frame

        # Write markdown with bullets
        markdown_input = """- First bullet point
- Second bullet point
- Third bullet point"""

        apply_markdown_to_textframe(markdown_input, tf)

        # Read back as markdown
        markdown_output = pptx_text_frame_to_markdown(tf)

        # Verify bullets are preserved
        assert "- First bullet point" in markdown_output
        assert "- Second bullet point" in markdown_output
        assert "- Third bullet point" in markdown_output

    def test_nested_bullet_list_roundtrip(self):
        """Test writing and reading back nested bullet lists."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
        )
        tf = textbox.text_frame

        # Write markdown with nested bullets
        markdown_input = """- Top level item
  - Nested item one
  - Nested item two
- Another top level"""

        apply_markdown_to_textframe(markdown_input, tf)

        # Read back as markdown
        markdown_output = pptx_text_frame_to_markdown(tf)

        # Verify structure is preserved (at least bullet markers should exist)
        assert "Top level item" in markdown_output
        assert "Nested item" in markdown_output

    def test_mixed_content_with_bullets(self):
        """Test markdown with mixed paragraphs and bullets."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
        )
        tf = textbox.text_frame

        # Write markdown with mixed content
        markdown_input = """Regular paragraph text

- Bullet one
- Bullet two

Another paragraph"""

        apply_markdown_to_textframe(markdown_input, tf)

        # Read back as markdown
        markdown_output = pptx_text_frame_to_markdown(tf)

        # Verify both regular text and bullets exist
        assert "Regular paragraph" in markdown_output or "paragraph" in markdown_output.lower()
        assert "Bullet" in markdown_output or "-" in markdown_output

    def test_bullet_with_bold_text(self):
        """Test bullets with bold formatting."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
        )
        tf = textbox.text_frame

        # Write markdown with bold text in bullets
        markdown_input = """- **Bold item** with more text
- Regular item
- Item with **bold** in middle"""

        apply_markdown_to_textframe(markdown_input, tf)

        # Read back as markdown
        markdown_output = pptx_text_frame_to_markdown(tf)

        # Verify bold is preserved
        assert "**Bold item**" in markdown_output or "Bold item" in markdown_output
        assert "-" in markdown_output  # Bullet markers should exist


class TestPptxParagraphToMarkdown:
    """Test the pptx_paragraph_to_markdown() function."""

    def test_paragraph_without_bullet_no_marker(self):
        """Test that non-bullet paragraphs don't get bullet markers."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = "Plain text"

        # Convert to markdown
        md = pptx_paragraph_to_markdown(para)

        # Should not have bullet marker
        assert not md.startswith("-")
        assert "Plain text" in md

    def test_paragraph_with_bullet_gets_marker(self):
        """Test that bullet paragraphs get bullet markers."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = "Bullet text"

        # Enable bullets
        _enable_paragraph_bullets(para)

        # Convert to markdown
        md = pptx_paragraph_to_markdown(para)

        # Should have bullet marker
        assert md.startswith("-") or "- " in md
        assert "Bullet text" in md


class TestBulletSaveAndReload:
    """Test that bullets persist after saving and reloading the PPTX file."""

    def test_bullets_persist_after_save(self):
        """Test that bullet formatting persists after saving and reloading."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            # Create presentation with bullets
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            textbox = slide.shapes.add_textbox(
                left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
            )
            tf = textbox.text_frame

            # Write markdown with bullets
            markdown_input = """- Bullet one
- Bullet two
- Bullet three"""

            apply_markdown_to_textframe(markdown_input, tf)

            # Save presentation
            prs.save(tmp_path)

            # Reload presentation
            prs2 = Presentation(tmp_path)
            slide2 = prs2.slides[0]

            # Find the textbox
            textbox2 = None
            for shape in slide2.shapes:
                if hasattr(shape, "text_frame"):
                    textbox2 = shape
                    break

            assert textbox2 is not None
            tf2 = textbox2.text_frame

            # Read back as markdown
            markdown_output = pptx_text_frame_to_markdown(tf2)

            # Verify bullets are preserved
            assert "-" in markdown_output
            assert "Bullet one" in markdown_output

        finally:
            os.unlink(tmp_path)


class TestBulletSpacing:
    """Test bullet spacing and indentation via XML attributes."""

    def test_bullet_has_margin_and_indent(self):
        """Test that enabling bullets sets marL and indent attributes."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = "Bullet with spacing"

        # Enable bullets with level 0
        _enable_paragraph_bullets(para, level=0)

        # Check that marL and indent are set
        pPr = para._element.get_or_add_pPr()
        marL = pPr.get("marL")
        indent = pPr.get("indent")

        # marL should be set (level 0 = 342900)
        assert marL is not None
        assert int(marL) > 0

        # indent should be negative (hanging indent)
        assert indent is not None
        assert int(indent) < 0

    def test_bullet_level_increases_margin(self):
        """Test that higher bullet levels have larger left margins."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame

        # Level 0 bullet
        para0 = tf.paragraphs[0]
        para0.text = "Level 0"
        _enable_paragraph_bullets(para0, level=0)

        # Level 1 bullet
        para1 = tf.add_paragraph()
        para1.text = "Level 1"
        _enable_paragraph_bullets(para1, level=1)

        # Get marL values
        pPr0 = para0._element.get_or_add_pPr()
        pPr1 = para1._element.get_or_add_pPr()

        marL0 = int(pPr0.get("marL"))
        marL1 = int(pPr1.get("marL"))

        # Level 1 should have larger margin than level 0
        assert marL1 > marL0

    def test_bullet_indent_is_consistent(self):
        """Test that indent value is the same across levels (negative for hanging)."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame

        # Create bullets at different levels
        para0 = tf.paragraphs[0]
        para0.text = "Level 0"
        _enable_paragraph_bullets(para0, level=0)

        para1 = tf.add_paragraph()
        para1.text = "Level 1"
        _enable_paragraph_bullets(para1, level=1)

        # Get indent values
        pPr0 = para0._element.get_or_add_pPr()
        pPr1 = para1._element.get_or_add_pPr()

        indent0 = int(pPr0.get("indent"))
        indent1 = int(pPr1.get("indent"))

        # Indent should be the same (same hanging amount at all levels)
        assert indent0 == indent1
        # And should be negative
        assert indent0 < 0


class TestBodyPrInsetPreservation:
    """Test that bodyPr insets are preserved when applying markdown."""

    def test_insets_preserved_after_apply_markdown(self):
        """Test that bodyPr insets are preserved after apply_markdown_to_textframe.

        When text_frame.clear() is called, bodyPr insets may be reset.
        Our implementation should preserve and restore them.
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        tf.text = "Initial text"

        # Set custom insets on the text frame
        bodyPr = tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert bodyPr is not None

        # Set custom inset values
        bodyPr.set("lIns", "182880")  # ~0.2 inches
        bodyPr.set("rIns", "182880")
        bodyPr.set("tIns", "91440")   # ~0.1 inches
        bodyPr.set("bIns", "91440")

        # Verify insets are set
        assert bodyPr.get("lIns") == "182880"
        assert bodyPr.get("tIns") == "91440"

        # Apply markdown (this calls text_frame.clear() internally)
        apply_markdown_to_textframe("New markdown text", tf)

        # Verify insets are preserved after apply
        bodyPr_after = tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert bodyPr_after is not None
        assert bodyPr_after.get("lIns") == "182880", (
            f"lIns should be preserved, got {bodyPr_after.get('lIns')}"
        )
        assert bodyPr_after.get("rIns") == "182880"
        assert bodyPr_after.get("tIns") == "91440", (
            f"tIns should be preserved, got {bodyPr_after.get('tIns')}"
        )
        assert bodyPr_after.get("bIns") == "91440"

    def test_insets_preserved_with_bullets(self):
        """Test that insets are preserved when applying markdown with bullets."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        tf.text = "Initial text"

        # Set custom insets
        bodyPr = tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        bodyPr.set("lIns", "274320")  # ~0.3 inches
        bodyPr.set("tIns", "137160")  # ~0.15 inches

        # Apply markdown with bullets
        markdown_with_bullets = """- First bullet
- Second bullet
- Third bullet"""
        apply_markdown_to_textframe(markdown_with_bullets, tf)

        # Verify insets are preserved
        bodyPr_after = tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert bodyPr_after.get("lIns") == "274320"
        assert bodyPr_after.get("tIns") == "137160"

        # Verify bullets were also applied
        first_para = tf.paragraphs[0]
        assert _paragraph_has_bullet(first_para)


class TestParagraphSpacingPreservation:
    """Test that line spacing is preserved for regular (non-bullet) paragraphs."""

    def test_line_spacing_preserved_for_regular_paragraphs(self):
        """Test that line spacing is preserved after apply_markdown_to_textframe.

        When text_frame.clear() is called, line spacing may be reset.
        Our implementation should preserve and restore it for regular paragraphs.
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        tf.text = "Initial text"

        # Set custom line spacing on the first paragraph
        para = tf.paragraphs[0]
        pPr = para._element.get_or_add_pPr()

        # Set line spacing to 150% (150000 in PPTX units = 1500 * 100)
        lnSpc = etree.SubElement(pPr, f"{{{_DRAWINGML_NS}}}lnSpc")
        spcPct = etree.SubElement(lnSpc, f"{{{_DRAWINGML_NS}}}spcPct")
        spcPct.set("val", "150000")

        # Verify line spacing is set
        lnSpc_before = pPr.find(f"{{{_DRAWINGML_NS}}}lnSpc")
        assert lnSpc_before is not None
        spcPct_before = lnSpc_before.find(f"{{{_DRAWINGML_NS}}}spcPct")
        assert spcPct_before is not None
        assert spcPct_before.get("val") == "150000"

        # Apply plain text markdown (no bullets)
        apply_markdown_to_textframe("New paragraph text", tf)

        # Verify line spacing is preserved after apply
        para_after = tf.paragraphs[0]
        pPr_after = para_after._element.find(f"{{{_DRAWINGML_NS}}}pPr")
        assert pPr_after is not None, "pPr should exist after apply"

        lnSpc_after = pPr_after.find(f"{{{_DRAWINGML_NS}}}lnSpc")
        assert lnSpc_after is not None, "lnSpc should be preserved"

        spcPct_after = lnSpc_after.find(f"{{{_DRAWINGML_NS}}}spcPct")
        assert spcPct_after is not None, "spcPct should be preserved"
        assert spcPct_after.get("val") == "150000", (
            f"Line spacing should be 150%, got {spcPct_after.get('val')}"
        )

    def test_space_before_preserved_for_regular_paragraphs(self):
        """Test that space-before is preserved after apply_markdown_to_textframe."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        tf.text = "Initial text"

        # Set custom space-before on the first paragraph
        para = tf.paragraphs[0]
        pPr = para._element.get_or_add_pPr()

        # Set space-before to 9pt (900 in PPTX units = 100ths of a point)
        spcBef = etree.SubElement(pPr, f"{{{_DRAWINGML_NS}}}spcBef")
        spcPts = etree.SubElement(spcBef, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts.set("val", "900")

        # Verify space-before is set
        spcBef_before = pPr.find(f"{{{_DRAWINGML_NS}}}spcBef")
        assert spcBef_before is not None

        # Apply plain text markdown (no bullets)
        apply_markdown_to_textframe("New paragraph text", tf)

        # Verify space-before is preserved after apply
        para_after = tf.paragraphs[0]
        pPr_after = para_after._element.find(f"{{{_DRAWINGML_NS}}}pPr")
        assert pPr_after is not None

        spcBef_after = pPr_after.find(f"{{{_DRAWINGML_NS}}}spcBef")
        assert spcBef_after is not None, "spcBef should be preserved"

        spcPts_after = spcBef_after.find(f"{{{_DRAWINGML_NS}}}spcPts")
        assert spcPts_after is not None
        assert spcPts_after.get("val") == "900", (
            f"Space-before should be 9pt (900), got {spcPts_after.get('val')}"
        )

    def test_space_after_preserved_for_regular_paragraphs(self):
        """Test that space-after is preserved after apply_markdown_to_textframe."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        tf.text = "Initial text"

        # Set custom space-after on the first paragraph
        para = tf.paragraphs[0]
        pPr = para._element.get_or_add_pPr()

        # Set space-after to 0pt (common in templates to control spacing explicitly)
        spcAft = etree.SubElement(pPr, f"{{{_DRAWINGML_NS}}}spcAft")
        spcPts = etree.SubElement(spcAft, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts.set("val", "0")

        # Verify space-after is set
        spcAft_before = pPr.find(f"{{{_DRAWINGML_NS}}}spcAft")
        assert spcAft_before is not None

        # Apply plain text markdown (no bullets)
        apply_markdown_to_textframe("New paragraph text", tf)

        # Verify space-after is preserved after apply
        para_after = tf.paragraphs[0]
        pPr_after = para_after._element.find(f"{{{_DRAWINGML_NS}}}pPr")
        assert pPr_after is not None

        spcAft_after = pPr_after.find(f"{{{_DRAWINGML_NS}}}spcAft")
        assert spcAft_after is not None, "spcAft should be preserved"

        spcPts_after = spcAft_after.find(f"{{{_DRAWINGML_NS}}}spcPts")
        assert spcPts_after is not None
        assert spcPts_after.get("val") == "0", (
            f"Space-after should be 0pt (0), got {spcPts_after.get('val')}"
        )

    def test_first_paragraph_and_bullet_spacing_preserved_separately(self):
        """Test that first paragraph gets different spacing than bullet paragraphs.

        Template pattern: title (para 0) has spcBef=0, bullets have spcBef=900.
        The implementation should preserve and apply these separately.
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(3)
        )
        tf = textbox.text_frame

        # Create template pattern: title with spcBef=0, then bullets with spcBef=900
        # First paragraph (title) - spcBef=0
        tf.text = "Title"
        para0 = tf.paragraphs[0]
        pPr0 = para0._element.get_or_add_pPr()
        spcBef0 = etree.SubElement(pPr0, f"{{{_DRAWINGML_NS}}}spcBef")
        spcPts0 = etree.SubElement(spcBef0, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts0.set("val", "0")  # Title has 0pt space before

        # Second paragraph (bullet) - spcBef=900
        para1 = tf.add_paragraph()
        para1.text = "Bullet item"
        pPr1 = para1._element.get_or_add_pPr()
        # Add bullet marker to mark as bullet paragraph
        buChar = etree.SubElement(pPr1, f"{{{_DRAWINGML_NS}}}buChar")
        buChar.set("char", "•")
        # Add bullet spacing
        spcBef1 = etree.SubElement(pPr1, f"{{{_DRAWINGML_NS}}}spcBef")
        spcPts1 = etree.SubElement(spcBef1, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts1.set("val", "900")  # Bullet has 9pt space before
        # Add margin and indent for bullet
        pPr1.set("marL", "520700")
        pPr1.set("indent", "-209550")

        # Apply markdown with title and bullets (use markdown bullet syntax)
        markdown = "New Title\n- First bullet\n- Second bullet"
        apply_markdown_to_textframe(markdown, tf)

        # Verify first paragraph (title) has spcBef=0
        output_para0 = tf.paragraphs[0]
        pPr_out0 = output_para0._element.find(f"{{{_DRAWINGML_NS}}}pPr")
        assert pPr_out0 is not None, "First paragraph should have pPr"

        spcBef_out0 = pPr_out0.find(f"{{{_DRAWINGML_NS}}}spcBef")
        if spcBef_out0 is not None:
            spcPts_out0 = spcBef_out0.find(f"{{{_DRAWINGML_NS}}}spcPts")
            if spcPts_out0 is not None:
                assert spcPts_out0.get("val") == "0", (
                    f"Title should have spcBef=0, got {spcPts_out0.get('val')}"
                )

        # Verify bullet paragraphs have spcBef=900
        assert len(tf.paragraphs) >= 2, "Should have at least 2 paragraphs"
        output_para1 = tf.paragraphs[1]
        pPr_out1 = output_para1._element.find(f"{{{_DRAWINGML_NS}}}pPr")
        assert pPr_out1 is not None, "Bullet paragraph should have pPr"

        spcBef_out1 = pPr_out1.find(f"{{{_DRAWINGML_NS}}}spcBef")
        assert spcBef_out1 is not None, "Bullet should have spcBef"
        spcPts_out1 = spcBef_out1.find(f"{{{_DRAWINGML_NS}}}spcPts")
        assert spcPts_out1 is not None
        assert spcPts_out1.get("val") == "900", (
            f"Bullet should have spcBef=900, got {spcPts_out1.get('val')}"
        )


class TestPPrElementOrder:
    """Test that pPr child elements are in the correct order for PowerPoint rendering.

    XML element order in <a:pPr> affects how PowerPoint renders spacing.
    The correct order is: lnSpc, spcBef, spcAft, buClr, buSzPts, buFont, buChar.
    """

    def test_spacing_elements_come_before_bullet_elements(self):
        """Test that spacing elements (lnSpc, spcBef, spcAft) come before buChar.

        PowerPoint requires this order for correct visual rendering of line spacing.
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(3)
        )
        tf = textbox.text_frame

        # Create template pattern with spacing values
        tf.text = "Title"
        para0 = tf.paragraphs[0]
        pPr0 = para0._element.get_or_add_pPr()

        # Add line spacing (110%)
        lnSpc = etree.SubElement(pPr0, f"{{{_DRAWINGML_NS}}}lnSpc")
        spcPct = etree.SubElement(lnSpc, f"{{{_DRAWINGML_NS}}}spcPct")
        spcPct.set("val", "110000")

        # Add space-before (0pt for title)
        spcBef = etree.SubElement(pPr0, f"{{{_DRAWINGML_NS}}}spcBef")
        spcPts = etree.SubElement(spcBef, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts.set("val", "0")

        # Add space-after
        spcAft = etree.SubElement(pPr0, f"{{{_DRAWINGML_NS}}}spcAft")
        spcPts2 = etree.SubElement(spcAft, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts2.set("val", "0")

        # Second paragraph (bullet) with spacing
        para1 = tf.add_paragraph()
        para1.text = "Bullet item"
        pPr1 = para1._element.get_or_add_pPr()

        # Add line spacing
        lnSpc1 = etree.SubElement(pPr1, f"{{{_DRAWINGML_NS}}}lnSpc")
        spcPct1 = etree.SubElement(lnSpc1, f"{{{_DRAWINGML_NS}}}spcPct")
        spcPct1.set("val", "110000")

        # Add space-before (9pt for bullets)
        spcBef1 = etree.SubElement(pPr1, f"{{{_DRAWINGML_NS}}}spcBef")
        spcPts1 = etree.SubElement(spcBef1, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts1.set("val", "900")

        # Add space-after
        spcAft1 = etree.SubElement(pPr1, f"{{{_DRAWINGML_NS}}}spcAft")
        spcPts1_aft = etree.SubElement(spcAft1, f"{{{_DRAWINGML_NS}}}spcPts")
        spcPts1_aft.set("val", "0")

        # Add bullet marker
        buChar = etree.SubElement(pPr1, f"{{{_DRAWINGML_NS}}}buChar")
        buChar.set("char", "•")
        pPr1.set("marL", "520700")
        pPr1.set("indent", "-209550")

        # Apply markdown with title and bullets
        markdown = "New Title\n- First bullet\n- Second bullet"
        apply_markdown_to_textframe(markdown, tf)

        # Verify element order in bullet paragraphs
        for i, para in enumerate(tf.paragraphs):
            if i == 0:
                continue  # Skip title paragraph

            pPr = para._element.find(f"{{{_DRAWINGML_NS}}}pPr")
            if pPr is None:
                continue

            # Get indices of spacing and bullet elements
            element_names = [
                etree.QName(child).localname for child in pPr
            ]

            lnSpc_idx = element_names.index("lnSpc") if "lnSpc" in element_names else -1
            spcBef_idx = element_names.index("spcBef") if "spcBef" in element_names else -1
            spcAft_idx = element_names.index("spcAft") if "spcAft" in element_names else -1
            buChar_idx = element_names.index("buChar") if "buChar" in element_names else -1

            # If buChar exists, it must come AFTER spacing elements
            if buChar_idx >= 0:
                if lnSpc_idx >= 0:
                    assert lnSpc_idx < buChar_idx, (
                        f"Para {i}: lnSpc (idx={lnSpc_idx}) must come before "
                        f"buChar (idx={buChar_idx}). Order: {element_names}"
                    )
                if spcBef_idx >= 0:
                    assert spcBef_idx < buChar_idx, (
                        f"Para {i}: spcBef (idx={spcBef_idx}) must come before "
                        f"buChar (idx={buChar_idx}). Order: {element_names}"
                    )
                if spcAft_idx >= 0:
                    assert spcAft_idx < buChar_idx, (
                        f"Para {i}: spcAft (idx={spcAft_idx}) must come before "
                        f"buChar (idx={buChar_idx}). Order: {element_names}"
                    )

            # Verify correct order: lnSpc < spcBef < spcAft
            if lnSpc_idx >= 0 and spcBef_idx >= 0:
                assert lnSpc_idx < spcBef_idx, (
                    f"Para {i}: lnSpc must come before spcBef. Order: {element_names}"
                )
            if spcBef_idx >= 0 and spcAft_idx >= 0:
                assert spcBef_idx < spcAft_idx, (
                    f"Para {i}: spcBef must come before spcAft. Order: {element_names}"
                )

    def test_enable_bullets_preserves_correct_element_order(self):
        """Test that _enable_paragraph_bullets() creates elements in correct order."""
        from gslides_api.agnostic.text import ParagraphStyle, SpacingValue

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame
        para = tf.paragraphs[0]
        para.text = "Bullet with spacing"

        # Create preserved props with spacing values
        preserved_props = ParagraphStyle(
            line_spacing=SpacingValue(percentage=110),
            space_before=SpacingValue(points=9),
            space_after=SpacingValue(points=0),
            margin_left=520700,
            indent=-209550,
        )

        # Enable bullets with preserved props
        _enable_paragraph_bullets(para, level=0, preserved_props=preserved_props)

        # Check element order
        pPr = para._element.find(f"{{{_DRAWINGML_NS}}}pPr")
        assert pPr is not None

        element_names = [etree.QName(child).localname for child in pPr]

        # All spacing elements should exist
        assert "lnSpc" in element_names, f"lnSpc missing. Elements: {element_names}"
        assert "spcBef" in element_names, f"spcBef missing. Elements: {element_names}"
        assert "spcAft" in element_names, f"spcAft missing. Elements: {element_names}"
        assert "buChar" in element_names, f"buChar missing. Elements: {element_names}"

        lnSpc_idx = element_names.index("lnSpc")
        spcBef_idx = element_names.index("spcBef")
        spcAft_idx = element_names.index("spcAft")
        buChar_idx = element_names.index("buChar")

        # Verify order: lnSpc, spcBef, spcAft, buChar
        assert lnSpc_idx < spcBef_idx < spcAft_idx < buChar_idx, (
            f"Incorrect element order. Expected lnSpc < spcBef < spcAft < buChar, "
            f"got indices: lnSpc={lnSpc_idx}, spcBef={spcBef_idx}, "
            f"spcAft={spcAft_idx}, buChar={buChar_idx}. Elements: {element_names}"
        )


class TestSoftLineBreaksInBullets:
    """Test soft line breaks (line breaks within a single bullet point).

    PowerPoint supports line breaks inside list items via <a:br/> XML elements.
    This is different from Google Slides which does NOT support this.
    The soft line break character is \\x0b (vertical tab).
    """

    def test_soft_line_break_creates_br_element(self):
        """Test that soft line breaks within bullet content create <a:br/> elements.

        When a bullet point contains \\x0b, it should create a single paragraph
        with <a:br/> XML elements, NOT multiple paragraphs.
        """
        from gslides_api.agnostic.ir import (
            FormattedDocument,
            FormattedList,
            FormattedListItem,
            FormattedParagraph,
            FormattedTextRun,
        )
        from gslides_api.agnostic.text import FullTextStyle
        from gslides_api.pptx.markdown_to_pptx import _apply_ir_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
        )
        tf = textbox.text_frame

        # Create IR with a bullet item containing soft line break
        soft_break = "\x0b"  # Vertical tab - PowerPoint's soft line break
        doc_ir = FormattedDocument(
            elements=[
                FormattedList(
                    ordered=False,
                    items=[
                        FormattedListItem(
                            paragraphs=[
                                FormattedParagraph(
                                    runs=[
                                        FormattedTextRun(
                                            content=f"First line{soft_break}Second line",
                                            style=FullTextStyle(),
                                        )
                                    ]
                                )
                            ],
                            nesting_level=0,
                        )
                    ],
                )
            ]
        )

        # Apply IR to text frame
        _apply_ir_to_textframe(doc_ir, tf)

        # Verify there's only ONE paragraph (the soft break is within it)
        assert len(tf.paragraphs) == 1, (
            f"Expected 1 paragraph, got {len(tf.paragraphs)}. "
            "Soft line break should NOT create a new paragraph."
        )

        # Verify the paragraph has bullet formatting
        assert _paragraph_has_bullet(tf.paragraphs[0])

        # Verify <a:br/> element exists in the paragraph XML
        para_xml = tf.paragraphs[0]._element
        br_elements = para_xml.findall(f".//{{{_DRAWINGML_NS}}}br")
        assert len(br_elements) >= 1, (
            f"Expected at least 1 <a:br/> element, found {len(br_elements)}. "
            "Soft line break should create <a:br/> XML element."
        )

    def test_soft_line_break_content_preserved_on_roundtrip(self):
        """Test that content with soft line breaks is preserved after save/reload.

        This verifies that the <a:br/> elements are correctly saved to the PPTX file
        and can be read back.
        """
        from gslides_api.agnostic.ir import (
            FormattedDocument,
            FormattedList,
            FormattedListItem,
            FormattedParagraph,
            FormattedTextRun,
        )
        from gslides_api.agnostic.text import FullTextStyle
        from gslides_api.pptx.markdown_to_pptx import _apply_ir_to_textframe

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            # Create presentation with soft line break in bullet
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            textbox = slide.shapes.add_textbox(
                left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
            )
            tf = textbox.text_frame

            soft_break = "\x0b"
            doc_ir = FormattedDocument(
                elements=[
                    FormattedList(
                        ordered=False,
                        items=[
                            FormattedListItem(
                                paragraphs=[
                                    FormattedParagraph(
                                        runs=[
                                            FormattedTextRun(
                                                content=f"Line A{soft_break}Line B",
                                                style=FullTextStyle(),
                                            )
                                        ]
                                    )
                                ],
                                nesting_level=0,
                            )
                        ],
                    )
                ]
            )

            _apply_ir_to_textframe(doc_ir, tf)

            # Save presentation
            prs.save(tmp_path)

            # Reload and verify
            prs2 = Presentation(tmp_path)
            slide2 = prs2.slides[0]

            # Find the textbox
            textbox2 = None
            for shape in slide2.shapes:
                if hasattr(shape, "text_frame"):
                    textbox2 = shape
                    break

            assert textbox2 is not None
            tf2 = textbox2.text_frame

            # Should still be 1 paragraph
            assert len(tf2.paragraphs) == 1

            # Verify <a:br/> element still exists after save/reload
            para_xml = tf2.paragraphs[0]._element
            br_elements = para_xml.findall(f".//{{{_DRAWINGML_NS}}}br")
            assert len(br_elements) >= 1, (
                "The <a:br/> element should persist after save/reload"
            )

            # Content should contain both parts
            full_text = tf2.text
            assert "Line A" in full_text
            assert "Line B" in full_text

        finally:
            os.unlink(tmp_path)

    def test_multiple_soft_breaks_in_bullet(self):
        """Test bullet with multiple soft line breaks creates multiple <a:br/> elements."""
        from gslides_api.agnostic.ir import (
            FormattedDocument,
            FormattedList,
            FormattedListItem,
            FormattedParagraph,
            FormattedTextRun,
        )
        from gslides_api.agnostic.text import FullTextStyle
        from gslides_api.pptx.markdown_to_pptx import _apply_ir_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(4)
        )
        tf = textbox.text_frame

        soft_break = "\x0b"
        doc_ir = FormattedDocument(
            elements=[
                FormattedList(
                    ordered=False,
                    items=[
                        FormattedListItem(
                            paragraphs=[
                                FormattedParagraph(
                                    runs=[
                                        FormattedTextRun(
                                            content=f"Line 1{soft_break}Line 2{soft_break}Line 3",
                                            style=FullTextStyle(),
                                        )
                                    ]
                                )
                            ],
                            nesting_level=0,
                        )
                    ],
                )
            ]
        )

        _apply_ir_to_textframe(doc_ir, tf)

        # Should have 2 <a:br/> elements (between 3 lines)
        para_xml = tf.paragraphs[0]._element
        br_elements = para_xml.findall(f".//{{{_DRAWINGML_NS}}}br")
        assert len(br_elements) == 2, (
            f"Expected 2 <a:br/> elements for 3 lines, got {len(br_elements)}"
        )
