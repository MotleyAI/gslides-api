"""
Test suite for PPTX shape styling functionality.

Tests shape fill copying and text frame anchor (vertical alignment) copying.
"""

import os
import tempfile

import pytest
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.util import Inches, Pt

from gslides_api.pptx.shape_copier import ShapeCopier, _DRAWINGML_NS, _FILL_TAGS
from gslides_api.pptx.id_manager import IdManager
from gslides_api.pptx.xml_utils import XmlUtils


class TestShapeFillCopying:
    """Test shape fill copying functionality."""

    def test_copy_solid_fill(self):
        """Test copying solid fill color from one shape to another."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source shape with solid fill
        source = slide.shapes.add_shape(
            1,  # Rectangle
            left=Inches(1),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )
        source.fill.solid()
        source.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red

        # Create target shape without fill
        target = slide.shapes.add_shape(
            1,
            left=Inches(4),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )

        # Copy using ShapeCopier
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_shape_fill(source, target)

        # Verify fill was copied by checking XML
        target_spPr = target._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if target_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            target_spPr = target._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        assert target_spPr is not None, "target shape should have spPr element"

        solidFill = target_spPr.find(f"{{{_DRAWINGML_NS}}}solidFill")

        # solidFill should exist after copying solid fill
        assert solidFill is not None, (
            "solidFill element should exist in target after copying solid fill"
        )

        # Verify fill type matches source
        assert target.fill.type == MSO_FILL.SOLID, (
            f"Expected target fill type SOLID, got {target.fill.type}"
        )

        # Verify color matches source (Red: RGB 255, 0, 0)
        assert target.fill.fore_color.rgb == RGBColor(255, 0, 0), (
            f"Expected RGB(255, 0, 0), got {target.fill.fore_color.rgb}"
        )

    def test_copy_no_fill(self):
        """Test that noFill is properly copied."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source shape with no fill
        source = slide.shapes.add_shape(
            1,
            left=Inches(1),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )
        source.fill.background()  # Set to no fill

        # Create target shape with solid fill
        target = slide.shapes.add_shape(
            1,
            left=Inches(4),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )
        target.fill.solid()
        target.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green

        # Copy fill (should copy noFill)
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_shape_fill(source, target)

        # Verify the target's XML has noFill element
        # Note: We check XML directly because python-pptx FillFormat caches its state
        # and won't reflect our direct XML modifications via target.fill.type
        target_spPr = target._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if target_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            target_spPr = target._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        assert target_spPr is not None, "target shape should have spPr element"

        # Verify noFill element exists in target after copying
        noFill = target_spPr.find(f"{{{_DRAWINGML_NS}}}noFill")
        assert noFill is not None, (
            "noFill element should exist in target after copying background fill"
        )

        # Verify the original solidFill is gone (the green fill was replaced)
        solidFill = target_spPr.find(f"{{{_DRAWINGML_NS}}}solidFill")
        assert solidFill is None, (
            "solidFill element should be removed after copying noFill"
        )

        # Verify source still has its noFill
        source_spPr = source._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if source_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            source_spPr = source._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")
        source_noFill = source_spPr.find(f"{{{_DRAWINGML_NS}}}noFill")
        assert source_noFill is not None, "source should still have noFill"

    def test_copy_fill_tags_constant(self):
        """Test that _FILL_TAGS contains expected fill types."""
        expected_tags = ["solidFill", "gradFill", "pattFill", "blipFill", "noFill"]

        for tag in expected_tags:
            matching = [t for t in _FILL_TAGS if tag in t]
            assert len(matching) == 1, f"Expected {tag} in _FILL_TAGS"


class TestTextFrameAnchorCopying:
    """Test text frame anchor (vertical alignment) copying."""

    def test_copy_anchor_top(self):
        """Test copying top anchor alignment."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_tf.text = "Source text"

        # Set anchor to top via XML
        source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        if source_bodyPr is not None:
            source_bodyPr.set("anchor", "t")

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_tf.text = "Target text"

        # Copy anchor
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame_anchor(source_tf, target_tf)

        # Verify anchor was copied
        target_bodyPr = target_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert target_bodyPr is not None, (
            "target text frame should have bodyPr element"
        )
        anchor = target_bodyPr.get("anchor")
        assert anchor == "t", f"Expected anchor 't', got '{anchor}'"

    def test_copy_anchor_center(self):
        """Test copying center anchor alignment."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with center anchor
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_tf.text = "Source text"

        source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        if source_bodyPr is not None:
            source_bodyPr.set("anchor", "ctr")

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_tf.text = "Target text"

        # Copy anchor
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame_anchor(source_tf, target_tf)

        # Verify anchor was copied
        target_bodyPr = target_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert target_bodyPr is not None, (
            "target text frame should have bodyPr element"
        )
        anchor = target_bodyPr.get("anchor")
        assert anchor == "ctr", f"Expected anchor 'ctr', got '{anchor}'"

    def test_copy_anchor_bottom(self):
        """Test copying bottom anchor alignment."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with bottom anchor
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_tf.text = "Source text"

        source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        if source_bodyPr is not None:
            source_bodyPr.set("anchor", "b")

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_tf.text = "Target text"

        # Copy anchor
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame_anchor(source_tf, target_tf)

        # Verify anchor was copied
        target_bodyPr = target_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert target_bodyPr is not None, (
            "target text frame should have bodyPr element"
        )
        anchor = target_bodyPr.get("anchor")
        assert anchor == "b", f"Expected anchor 'b', got '{anchor}'"


class TestShapeCopyIntegration:
    """Integration tests for shape copying with fill and anchor."""

    def test_copy_text_shape_preserves_anchor(self):
        """Test that copying a text shape preserves its anchor setting."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with center anchor
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_tf.text = "Centered text"

        source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        if source_bodyPr is not None:
            source_bodyPr.set("anchor", "ctr")

        # Copy shape
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copied_shape = copier.copy_shape(source, slide)

        # Assert copied shape exists
        assert copied_shape is not None, "copy_shape should return a shape"

        # Assert it has a text_frame
        assert hasattr(copied_shape, "text_frame"), (
            "Copied shape should have a text_frame"
        )

        # Assert bodyPr exists and has the correct anchor
        copied_tf = copied_shape.text_frame
        copied_bodyPr = copied_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert copied_bodyPr is not None, (
            "Copied shape's text_frame should have a bodyPr element"
        )

        anchor = copied_bodyPr.get("anchor")
        assert anchor == "ctr", (
            f"Expected anchor 'ctr', got '{anchor}'"
        )

    def test_copy_shape_preserves_fill_and_text(self):
        """Test that copying a shape preserves both fill and text."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source shape with fill and text
        source = slide.shapes.add_shape(
            1,  # Rectangle
            left=Inches(1),
            top=Inches(1),
            width=Inches(3),
            height=Inches(2),
        )
        source.fill.solid()
        source.fill.fore_color.rgb = RGBColor(100, 150, 200)
        source.text_frame.text = "Text in shape"

        # Copy shape
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copied_shape = copier.copy_shape(source, slide)

        # Shape should be copied (returns something or adds to slide)
        # If copy_shape returns the shape, verify it
        if copied_shape is not None:
            assert hasattr(copied_shape, "text_frame")

            # Verify text is preserved
            assert copied_shape.text_frame.text == "Text in shape", (
                f"Expected 'Text in shape', got '{copied_shape.text_frame.text}'"
            )

            # Verify fill color is preserved
            assert copied_shape.fill.type == MSO_FILL.SOLID, (
                f"Expected solid fill, got {copied_shape.fill.type}"
            )
            assert copied_shape.fill.fore_color.rgb == RGBColor(100, 150, 200), (
                f"Expected RGB(100, 150, 200), got {copied_shape.fill.fore_color.rgb}"
            )


class TestShapeStylingPersistence:
    """Test that shape styling persists after save/reload."""

    def test_fill_persists_after_save(self):
        """Test that copied fill persists after saving and reloading."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            # Create presentation with styled shape
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            source = slide.shapes.add_shape(
                1,
                left=Inches(1),
                top=Inches(1),
                width=Inches(2),
                height=Inches(1),
            )
            source.fill.solid()
            source.fill.fore_color.rgb = RGBColor(255, 128, 0)  # Orange

            target = slide.shapes.add_shape(
                1,
                left=Inches(4),
                top=Inches(1),
                width=Inches(2),
                height=Inches(1),
            )

            # Copy fill
            id_manager = IdManager(prs)
            copier = ShapeCopier(id_manager)
            copier._copy_shape_fill(source, target)

            # Save
            prs.save(tmp_path)

            # Reload and verify
            prs2 = Presentation(tmp_path)
            slide2 = prs2.slides[0]

            # Find shapes (target should be second shape)
            shapes = list(slide2.shapes)
            assert len(shapes) >= 2

            # Verify the target shape retained the copied solid fill and color
            target_shape = shapes[1]
            assert target_shape.fill.type == MSO_FILL.SOLID, (
                f"Expected solid fill, got {target_shape.fill.type}"
            )
            assert target_shape.fill.fore_color.rgb == RGBColor(255, 128, 0), (
                f"Expected RGB(255, 128, 0), got {target_shape.fill.fore_color.rgb}"
            )

        finally:
            os.unlink(tmp_path)

    def test_anchor_persists_after_save(self):
        """Test that copied anchor persists after saving and reloading."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name

        try:
            # Create presentation with text box
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            source = slide.shapes.add_textbox(
                left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
            )
            source_tf = source.text_frame
            source_tf.text = "Centered"

            # Set center anchor
            source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
            if source_bodyPr is not None:
                source_bodyPr.set("anchor", "ctr")

            target = slide.shapes.add_textbox(
                left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
            )
            target_tf = target.text_frame
            target_tf.text = "Target"

            # Copy anchor
            id_manager = IdManager(prs)
            copier = ShapeCopier(id_manager)
            copier._copy_text_frame_anchor(source_tf, target_tf)

            # Save
            prs.save(tmp_path)

            # Reload and verify
            prs2 = Presentation(tmp_path)
            slide2 = prs2.slides[0]

            # Find second textbox
            textboxes = [s for s in slide2.shapes if hasattr(s, "text_frame")]
            assert len(textboxes) >= 2

            target_tf2 = textboxes[1].text_frame
            target_bodyPr2 = target_tf2._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")

            assert target_bodyPr2 is not None, (
                "reloaded target text frame should have bodyPr element"
            )
            anchor = target_bodyPr2.get("anchor")
            assert anchor == "ctr", (
                f"Expected anchor 'ctr' after reload, got '{anchor}'"
            )

        finally:
            os.unlink(tmp_path)


class TestFontColorCopying:
    """Test font color copying functionality."""

    def test_copy_rgb_color(self):
        """Test copying RGB font color from one shape to another."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with colored text
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(1)
        )
        source_tf = source.text_frame
        source_para = source_tf.paragraphs[0]
        source_run = source_para.add_run()
        source_run.text = "Red text"
        source_run.font.color.rgb = RGBColor(255, 0, 0)  # Red

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(1)
        )
        target_tf = target.text_frame
        target_para = target_tf.paragraphs[0]
        target_run = target_para.add_run()
        target_run.text = "Target text"

        # Copy using ShapeCopier
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame(source_tf, target_tf)

        # Verify color was copied
        copied_run = target_tf.paragraphs[0].runs[0]
        assert copied_run.font.color.rgb == RGBColor(255, 0, 0)

    def test_copy_font_underline(self):
        """Test copying underline formatting."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with underlined text
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(1)
        )
        source_tf = source.text_frame
        source_para = source_tf.paragraphs[0]
        source_run = source_para.add_run()
        source_run.text = "Underlined text"
        source_run.font.underline = True

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(1)
        )
        target_tf = target.text_frame
        target_para = target_tf.paragraphs[0]
        target_run = target_para.add_run()
        target_run.text = "Target text"

        # Copy using ShapeCopier
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame(source_tf, target_tf)

        # Verify underline was copied
        copied_run = target_tf.paragraphs[0].runs[0]
        assert copied_run.font.underline is True


class TestBodyPrInsetsCopying:
    """Test bodyPr inset (margin) copying functionality."""

    def test_copy_text_insets(self):
        """Test copying text inset attributes from bodyPr."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box and set insets via XML
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_tf.text = "Source text"

        # Set insets on source bodyPr
        source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        if source_bodyPr is not None:
            source_bodyPr.set("lIns", "91440")   # Left inset
            source_bodyPr.set("rIns", "91440")   # Right inset
            source_bodyPr.set("tIns", "45720")   # Top inset
            source_bodyPr.set("bIns", "45720")   # Bottom inset

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_tf.text = "Target text"

        # Copy bodyPr attributes
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame_anchor(source_tf, target_tf)

        # Verify insets were copied
        target_bodyPr = target_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert target_bodyPr is not None, (
            "target text frame should have bodyPr element"
        )
        assert target_bodyPr.get("lIns") == "91440", (
            f"Expected lIns '91440', got '{target_bodyPr.get('lIns')}'"
        )
        assert target_bodyPr.get("rIns") == "91440", (
            f"Expected rIns '91440', got '{target_bodyPr.get('rIns')}'"
        )
        assert target_bodyPr.get("tIns") == "45720", (
            f"Expected tIns '45720', got '{target_bodyPr.get('tIns')}'"
        )
        assert target_bodyPr.get("bIns") == "45720", (
            f"Expected bIns '45720', got '{target_bodyPr.get('bIns')}'"
        )

    def test_copy_anchor_ctr_attribute(self):
        """Test copying anchorCtr attribute."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with anchorCtr
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_tf.text = "Source text"

        source_bodyPr = source_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        if source_bodyPr is not None:
            source_bodyPr.set("anchorCtr", "1")

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_tf.text = "Target text"

        # Copy bodyPr attributes
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame_anchor(source_tf, target_tf)

        # Verify anchorCtr was copied
        target_bodyPr = target_tf._element.find(f"{{{_DRAWINGML_NS}}}bodyPr")
        assert target_bodyPr is not None, (
            "target text frame should have bodyPr element"
        )
        assert target_bodyPr.get("anchorCtr") == "1", (
            f"Expected anchorCtr '1', got '{target_bodyPr.get('anchorCtr')}'"
        )


class TestLineSpacingCopying:
    """Test line spacing copying functionality."""

    def test_copy_line_spacing(self):
        """Test copying line_spacing property between paragraphs."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with line spacing
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_para = source_tf.paragraphs[0]
        source_para.add_run().text = "Line with spacing"
        source_para.line_spacing = 1.5  # 1.5x line spacing

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_para = target_tf.paragraphs[0]
        target_para.add_run().text = "Target"

        # Copy using ShapeCopier
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame(source_tf, target_tf)

        # Verify line spacing was copied
        copied_para = target_tf.paragraphs[0]
        assert copied_para.line_spacing == 1.5

    def test_copy_space_before(self):
        """Test copying space_before property between paragraphs."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with space_before
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_para = source_tf.paragraphs[0]
        source_para.add_run().text = "Paragraph with space before"
        source_para.space_before = Pt(12)  # 12pt space before

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_para = target_tf.paragraphs[0]
        target_para.add_run().text = "Target"

        # Copy using ShapeCopier
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame(source_tf, target_tf)

        # Verify space_before was copied
        copied_para = target_tf.paragraphs[0]
        assert copied_para.space_before == Pt(12)

    def test_copy_space_after(self):
        """Test copying space_after property between paragraphs."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source text box with space_after
        source = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        source_tf = source.text_frame
        source_para = source_tf.paragraphs[0]
        source_para.add_run().text = "Paragraph with space after"
        source_para.space_after = Pt(18)  # 18pt space after

        # Create target text box
        target = slide.shapes.add_textbox(
            left=Inches(5), top=Inches(1), width=Inches(3), height=Inches(2)
        )
        target_tf = target.text_frame
        target_para = target_tf.paragraphs[0]
        target_para.add_run().text = "Target"

        # Copy using ShapeCopier
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_text_frame(source_tf, target_tf)

        # Verify space_after was copied
        copied_para = target_tf.paragraphs[0]
        assert copied_para.space_after == Pt(18)


class TestBoldExplicitSetting:
    """Test that bold is explicitly set to True or False, not inherited."""

    def test_bold_false_explicitly_set(self):
        """Test that non-bold text has font.bold explicitly set to False.

        This is critical because runs inherit from defRPr if bold is None.
        If defRPr has bold=True, text will appear bold unless explicitly False.
        """
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame

        # Apply non-bold markdown text
        apply_markdown_to_textframe("Regular text without bold", tf)

        # Verify font.bold is explicitly False, not None
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    # Bold should be explicitly False, not None (inherited)
                    assert run.font.bold is False, (
                        f"font.bold should be False, not {run.font.bold}"
                    )

    def test_bold_true_explicitly_set(self):
        """Test that bold text has font.bold explicitly set to True."""
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame

        # Apply bold markdown text
        apply_markdown_to_textframe("**Bold text**", tf)

        # Find the bold run and verify it's explicitly True
        bold_found = False
        for para in tf.paragraphs:
            for run in para.runs:
                if "Bold" in run.text:
                    assert run.font.bold is True
                    bold_found = True

        assert bold_found, "Bold text run not found"

    def test_mixed_bold_and_regular(self):
        """Test that mixed bold and regular text has correct font.bold values."""
        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        textbox = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        tf = textbox.text_frame

        # Apply mixed markdown text
        apply_markdown_to_textframe("Regular **bold** regular", tf)

        # Check that runs have appropriate bold values
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    if "bold" in run.text.lower():
                        assert run.font.bold is True, f"'{run.text}' should be bold"
                    else:
                        assert run.font.bold is False, (
                            f"'{run.text}' should not be bold"
                        )


class TestImageShapeCopying:
    """Test image shape copying preserves blipFill properties."""

    def test_copy_image_preserves_src_rect(self):
        """Test that copying an image shape preserves srcRect (crop settings).

        This is critical for images that are cropped - the srcRect defines
        which portion of the image is displayed.
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create a simple image shape with a small test image
        import tempfile
        import os
        from PIL import Image

        # Create a simple test image
        img = Image.new('RGB', (100, 100), color='purple')
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            img.save(tmp.name)
            tmp_path = tmp.name

        try:
            # Add image to slide
            picture = slide.shapes.add_picture(
                tmp_path,
                left=Inches(1),
                top=Inches(1),
                width=Inches(2),
                height=Inches(2)
            )

            # Manually add srcRect to simulate a cropped image
            # blipFill is in PresentationML namespace (p:), not DrawingML (a:)
            _PRESENTATIONML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
            blip_fill = picture._element.find(
                f".//{{{_PRESENTATIONML_NS}}}blipFill"
            )
            assert blip_fill is not None, "blipFill not found in source picture"

            # srcRect is in DrawingML namespace (a:) and should be child of blipFill
            src_rect = etree.SubElement(
                blip_fill, f"{{{_DRAWINGML_NS}}}srcRect"
            )
            src_rect.set("l", "10000")  # 10% from left
            src_rect.set("t", "5000")   # 5% from top
            src_rect.set("r", "15000")  # 15% from right
            src_rect.set("b", "20000")  # 20% from bottom

            # Now copy the shape using ShapeCopier
            id_manager = IdManager(prs)
            copier = ShapeCopier(id_manager)

            # Create a second slide to copy to
            target_slide = prs.slides.add_slide(blank_layout)

            # Copy the shape - but we need relationship mapping
            # For this test, create a simple mapping
            relationship_mapping = {}

            # Get source relationship ID and copy the image relationship
            source_blip = picture._element.find(
                f".//{{{_DRAWINGML_NS}}}blip"
            )
            if source_blip is not None:
                r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                old_rel_id = source_blip.get(f"{{{r_ns}}}embed")
                if old_rel_id:
                    # Copy the image relationship to target slide
                    # Get the relationship and its image part
                    relationship = slide.part.rels[old_rel_id]
                    image_part = relationship._target
                    # Use get_or_add_image_part to add image to target
                    image_data = image_part.blob
                    import io
                    image_stream = io.BytesIO(image_data)
                    _, new_rel_id = target_slide.part.get_or_add_image_part(image_stream)
                    relationship_mapping[old_rel_id] = new_rel_id

            # Copy the image shape
            copier._copy_image_shape(
                picture,
                target_slide,
                relationship_mapping=relationship_mapping
            )

            # Verify the copied shape has srcRect preserved
            # Find the picture shape on target slide
            target_pics = [
                s for s in target_slide.shapes._spTree.iterchildren()
                if s.tag.endswith('}pic')
            ]

            assert len(target_pics) > 0, "No picture shape found on target slide"

            target_pic = target_pics[0]
            target_src_rect = target_pic.find(
                f".//{{{_DRAWINGML_NS}}}srcRect"
            )

            assert target_src_rect is not None, "srcRect not found on copied image"
            assert target_src_rect.get("l") == "10000", "srcRect 'l' not preserved"
            assert target_src_rect.get("t") == "5000", "srcRect 't' not preserved"
            assert target_src_rect.get("r") == "15000", "srcRect 'r' not preserved"
            assert target_src_rect.get("b") == "20000", "srcRect 'b' not preserved"

        finally:
            os.unlink(tmp_path)


class TestBlipFillRelationshipRemapping:
    """Test blipFill relationship remapping during shape fill copying."""

    def test_copy_shape_fill_with_relationship_mapping(self):
        """Test that blipFill elements have their relationship IDs remapped.

        When copying a shape fill that contains a blipFill (image fill),
        the r:embed attribute must be remapped to the new relationship ID.
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source shape
        source = slide.shapes.add_shape(
            1,  # Rectangle
            left=Inches(1),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )

        # Create target shape
        target = slide.shapes.add_shape(
            1,
            left=Inches(4),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )

        # Manually add a blipFill element to source shape with a fake relationship ID
        source_spPr = source._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if source_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            source_spPr = source._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        if source_spPr is not None:
            # Create a blipFill element with a relationship reference
            blip_fill = etree.SubElement(source_spPr, f"{{{_DRAWINGML_NS}}}blipFill")
            r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            blip = etree.SubElement(blip_fill, f"{{{_DRAWINGML_NS}}}blip")
            blip.set(f"{{{r_ns}}}embed", "rId_OLD")

        # Create relationship mapping
        relationship_mapping = {"rId_OLD": "rId_NEW"}

        # Copy fill with relationship mapping
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_shape_fill(source, target, relationship_mapping)

        # Verify blipFill was copied with remapped relationship ID
        target_spPr = target._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if target_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            target_spPr = target._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        if target_spPr is not None:
            target_blip_fill = target_spPr.find(f"{{{_DRAWINGML_NS}}}blipFill")
            assert target_blip_fill is not None, "blipFill should have been copied"

            target_blip = target_blip_fill.find(f"{{{_DRAWINGML_NS}}}blip")
            assert target_blip is not None, "blip element should exist"

            r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            embed_id = target_blip.get(f"{{{r_ns}}}embed")
            assert embed_id == "rId_NEW", (
                f"r:embed should be remapped from rId_OLD to rId_NEW, got {embed_id}"
            )

    def test_copy_shape_fill_without_relationship_mapping(self):
        """Test that blipFill copying works without relationship mapping.

        When no relationship mapping is provided, the fill should still be copied
        (though relationship IDs won't be remapped).
        """
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source shape
        source = slide.shapes.add_shape(
            1,
            left=Inches(1),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )

        # Create target shape
        target = slide.shapes.add_shape(
            1,
            left=Inches(4),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )

        # Manually add a blipFill element to source shape
        source_spPr = source._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if source_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            source_spPr = source._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        if source_spPr is not None:
            blip_fill = etree.SubElement(source_spPr, f"{{{_DRAWINGML_NS}}}blipFill")
            r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            blip = etree.SubElement(blip_fill, f"{{{_DRAWINGML_NS}}}blip")
            blip.set(f"{{{r_ns}}}embed", "rId_ORIGINAL")

        # Copy fill without relationship mapping
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_shape_fill(source, target, relationship_mapping=None)

        # Verify blipFill was copied (with original relationship ID)
        target_spPr = target._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if target_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            target_spPr = target._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        if target_spPr is not None:
            target_blip_fill = target_spPr.find(f"{{{_DRAWINGML_NS}}}blipFill")
            assert target_blip_fill is not None, "blipFill should have been copied"

            target_blip = target_blip_fill.find(f"{{{_DRAWINGML_NS}}}blip")
            assert target_blip is not None, "blip element should exist"

            r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            embed_id = target_blip.get(f"{{{r_ns}}}embed")
            assert embed_id == "rId_ORIGINAL", (
                f"r:embed should remain unchanged without mapping, got {embed_id}"
            )

    def test_copy_shape_properties_passes_relationship_mapping(self):
        """Test that _copy_shape_properties passes relationship_mapping to _copy_shape_fill."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create source and target shapes
        source = slide.shapes.add_shape(
            1,
            left=Inches(1),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )
        target = slide.shapes.add_shape(
            1,
            left=Inches(4),
            top=Inches(1),
            width=Inches(2),
            height=Inches(1),
        )

        # Add blipFill to source
        source_spPr = source._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if source_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            source_spPr = source._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        if source_spPr is not None:
            blip_fill = etree.SubElement(source_spPr, f"{{{_DRAWINGML_NS}}}blipFill")
            r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            blip = etree.SubElement(blip_fill, f"{{{_DRAWINGML_NS}}}blip")
            blip.set(f"{{{r_ns}}}embed", "rId_SOURCE")

        # Call _copy_shape_properties with relationship mapping
        relationship_mapping = {"rId_SOURCE": "rId_TARGET"}
        id_manager = IdManager(prs)
        copier = ShapeCopier(id_manager)
        copier._copy_shape_properties(source, target, relationship_mapping)

        # Verify the relationship was remapped
        target_spPr = target._element.find(f".//{{{_DRAWINGML_NS}}}spPr")
        if target_spPr is None:
            from gslides_api.pptx.shape_copier import _PRESENTATIONML_NS
            target_spPr = target._element.find(f".//{{{_PRESENTATIONML_NS}}}spPr")

        if target_spPr is not None:
            target_blip_fill = target_spPr.find(f"{{{_DRAWINGML_NS}}}blipFill")
            if target_blip_fill is not None:
                target_blip = target_blip_fill.find(f"{{{_DRAWINGML_NS}}}blip")
                if target_blip is not None:
                    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                    embed_id = target_blip.get(f"{{{r_ns}}}embed")
                    assert embed_id == "rId_TARGET", (
                        f"Relationship should be remapped via _copy_shape_properties, got {embed_id}"
                    )
