"""Test PowerPoint adapter implementation."""

import os
import tempfile
from unittest.mock import MagicMock, Mock, patch

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.slide import Slide

from gslides_api.agnostic.units import OutputUnit

from gslides_api.adapters.abstract_slides import AbstractAltText, AbstractThumbnail
from gslides_api.adapters.pptx_adapter import (
    PowerPointAPIClient,
    PowerPointElementParent,
    PowerPointImageElement,
    PowerPointPresentation,
    PowerPointShapeElement,
    PowerPointSlide,
    PowerPointSpeakerNotes,
    PowerPointTableElement,
    pptx_element_discriminator,
    validate_pptx_element,
)


class TestPowerPointAPIClient:
    """Test PowerPoint API client implementation."""

    def test_init(self):
        """Test API client initialization."""
        client = PowerPointAPIClient()
        assert client.auto_flush is True

    def test_auto_flush_property(self):
        """Test auto_flush property getter and setter."""
        client = PowerPointAPIClient()
        assert client.auto_flush is True

        client.auto_flush = False
        assert client.auto_flush is False

    def test_flush_batch_update(self):
        """Test flush_batch_update does nothing (no-op for filesystem)."""
        client = PowerPointAPIClient()
        # Should not raise any exception
        client.flush_batch_update()

    def test_copy_presentation(self):
        """Test presentation copying."""
        client = PowerPointAPIClient()

        with tempfile.TemporaryDirectory() as temp_dir:
            # Create a dummy presentation file
            source_path = os.path.join(temp_dir, "source.pptx")
            with open(source_path, "w") as f:
                f.write("dummy content")

            # Test copying to same folder
            result = client.copy_presentation(source_path, "copy_title")

            assert result["name"] == "copy_title"
            assert result["id"].endswith("copy_title.pptx")
            assert os.path.exists(result["id"])

    def test_copy_presentation_with_folder(self):
        """Test presentation copying to specific folder."""
        client = PowerPointAPIClient()

        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = os.path.join(temp_dir, "source.pptx")
            dest_folder = os.path.join(temp_dir, "dest")
            os.makedirs(dest_folder)

            with open(source_path, "w") as f:
                f.write("dummy content")

            result = client.copy_presentation(source_path, "copy_title", dest_folder)

            assert result["name"] == "copy_title"
            assert result["parents"] == [dest_folder]
            assert os.path.exists(result["id"])

    def test_copy_presentation_file_not_found(self):
        """Test copying non-existent presentation raises error."""
        client = PowerPointAPIClient()

        with pytest.raises(FileNotFoundError):
            client.copy_presentation("non_existent.pptx", "copy_title")

    def test_create_folder(self):
        """Test folder creation."""
        client = PowerPointAPIClient()

        with tempfile.TemporaryDirectory() as temp_dir:
            folder_name = "test_folder"
            result = client.create_folder(folder_name, parent_folder_id=temp_dir)

            expected_path = os.path.join(temp_dir, folder_name)
            assert result["id"] == expected_path
            assert result["name"] == folder_name
            assert result["parents"] == [temp_dir]
            assert os.path.exists(expected_path)

    def test_create_folder_ignore_existing(self):
        """Test folder creation with ignore_existing=True."""
        client = PowerPointAPIClient()

        with tempfile.TemporaryDirectory() as temp_dir:
            folder_name = "test_folder"
            folder_path = os.path.join(temp_dir, folder_name)
            os.makedirs(folder_path)  # Create folder first

            # Should not raise error
            result = client.create_folder(folder_name, parent_folder_id=temp_dir)
            assert result["id"] == folder_path

    def test_delete_file(self):
        """Test file deletion."""
        client = PowerPointAPIClient()

        with tempfile.TemporaryDirectory() as temp_dir:
            test_file = os.path.join(temp_dir, "test.pptx")
            with open(test_file, "w") as f:
                f.write("test content")

            assert os.path.exists(test_file)
            client.delete_file(test_file)
            assert not os.path.exists(test_file)

    def test_delete_nonexistent_file(self):
        """Test deleting non-existent file doesn't raise error."""
        client = PowerPointAPIClient()
        # Should not raise any exception
        client.delete_file("non_existent_file.pptx")

    def test_credentials_methods(self):
        """Test credential methods (no-op for filesystem)."""
        client = PowerPointAPIClient()

        # Should not raise errors
        client.set_credentials(None)
        assert client.get_credentials() is None

    @patch("gslides_api.adapters.pptx_adapter.Presentation")
    def test_replace_text(self, mock_presentation_class):
        """Test text replacement in slides."""
        client = PowerPointAPIClient()

        # Mock presentation and slides
        mock_prs = Mock()
        mock_slide = Mock()
        mock_shape = Mock()
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()

        mock_run.text = "Hello World"
        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_shape.text_frame = mock_text_frame
        mock_slide.shapes = [mock_shape]
        mock_prs.slides = [mock_slide]

        mock_presentation_class.return_value = mock_prs

        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            client.replace_text(["0"], "Hello", "Hi", temp_file.name)

            # Verify text was replaced
            assert mock_run.text == "Hi World"
            mock_prs.save.assert_called_once_with(temp_file.name)

    def test_get_default_api_client(self):
        """Test getting default API client."""
        client = PowerPointAPIClient.get_default_api_client()
        assert isinstance(client, PowerPointAPIClient)


class TestPowerPointSpeakerNotes:
    """Test PowerPoint speaker notes implementation."""

    def test_init(self):
        """Test speaker notes initialization."""
        mock_notes_slide = Mock()
        notes = PowerPointSpeakerNotes(mock_notes_slide)
        assert notes.notes_slide == mock_notes_slide

    def test_read_text_empty(self):
        """Test reading text from empty notes."""
        mock_notes_slide = Mock()
        mock_notes_slide.notes_text_frame = None

        notes = PowerPointSpeakerNotes(mock_notes_slide)
        assert notes.read_text() == ""

    def test_read_text_with_content(self):
        """Test reading text from notes with content."""
        mock_notes_slide = Mock()
        mock_text_frame = Mock()
        mock_text_frame.text = "Speaker notes content"
        mock_notes_slide.notes_text_frame = mock_text_frame

        notes = PowerPointSpeakerNotes(mock_notes_slide)
        assert notes.read_text(as_markdown=False) == "Speaker notes content"

    def test_write_text(self):
        """Test writing text to speaker notes."""
        mock_notes_slide = Mock()
        mock_text_frame = Mock()
        mock_notes_slide.notes_text_frame = mock_text_frame

        notes = PowerPointSpeakerNotes(mock_notes_slide)
        api_client = PowerPointAPIClient()

        notes.write_text(api_client, "New content")

        mock_text_frame.clear.assert_called_once()
        assert mock_text_frame.text == "New content"


class TestElementDiscriminator:
    """Test element discriminator functionality."""

    def test_discriminator_shape_element(self):
        """Test discriminator with shape element."""
        mock_element = Mock()
        mock_element.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE

        result = pptx_element_discriminator(mock_element)
        assert result == "shape"

    def test_discriminator_image_element(self):
        """Test discriminator with image element."""
        mock_element = Mock()
        mock_element.shape_type = MSO_SHAPE_TYPE.PICTURE
        # Picture elements don't have text_frame attribute
        del mock_element.text_frame

        result = pptx_element_discriminator(mock_element)
        assert result == "image"

    def test_discriminator_table_element(self):
        """Test discriminator with table element."""
        mock_element = Mock()
        mock_element.shape_type = MSO_SHAPE_TYPE.TABLE
        # Table elements don't have text_frame attribute
        del mock_element.text_frame

        result = pptx_element_discriminator(mock_element)
        assert result == "table"

    def test_discriminator_placeholder_with_text(self):
        """Test discriminator with placeholder element that has text."""
        mock_element = Mock()
        mock_element.shape_type = MSO_SHAPE_TYPE.PLACEHOLDER
        mock_element.text_frame = Mock()

        result = pptx_element_discriminator(mock_element)
        assert result == "shape"

    def test_discriminator_generic_element(self):
        """Test discriminator with generic element."""
        mock_element = Mock()
        mock_element.shape_type = MSO_SHAPE_TYPE.LINE
        # Generic LINE elements don't have text_frame attribute
        del mock_element.text_frame

        result = pptx_element_discriminator(mock_element)
        assert result == "generic"


class TestPowerPointElementParent:
    """Test PowerPoint element parent class."""

    def test_convert_from_pptx_element(self):
        """Test converting from pptx element."""
        # Test with a dictionary input (already converted case)
        test_data = {
            "objectId": "123",
            "alt_text": {"title": "Test Shape"},
            "type": "generic",
            "pptx_element": "mock_element",
        }

        result = PowerPointElementParent.convert_from_pptx_element(test_data)

        assert result["objectId"] == "123"
        assert result["alt_text"]["title"] == "Test Shape"
        assert result["type"] == "generic"
        assert result["pptx_element"] == "mock_element"

    def test_absolute_size(self):
        """Test getting absolute size."""
        mock_shape = Mock()
        mock_shape.width = 914400  # 1 inch in EMU
        mock_shape.height = 914400

        element = PowerPointElementParent(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        width, height = element.absolute_size(units=OutputUnit.IN)
        assert abs(width - 1.0) < 0.01  # Close to 1 inch
        assert abs(height - 1.0) < 0.01

    def test_absolute_position(self):
        """Test getting absolute position."""
        mock_shape = Mock()
        mock_shape.left = 914400  # 1 inch in EMU
        mock_shape.top = 914400

        element = PowerPointElementParent(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        left, top = element.absolute_position(units=OutputUnit.IN)
        assert abs(left - 1.0) < 0.01  # Close to 1 inch
        assert abs(top - 1.0) < 0.01

    def test_set_alt_text(self):
        """Test setting alt text."""
        mock_shape = Mock()
        mock_shape.name = "Original Name"

        element = PowerPointElementParent(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        api_client = PowerPointAPIClient()
        element.set_alt_text(api_client, title="New Title")

        assert mock_shape.name == "New Title"
        assert element.alt_text.title == "New Title"


class TestPowerPointShapeElement:
    """Test PowerPoint shape element implementation."""

    def test_has_text_with_text_frame(self):
        """Test has_text property with text frame."""
        mock_shape = Mock()
        mock_shape.text_frame = Mock()
        mock_shape.text_frame.text = "Some text content"

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        assert element.has_text

    def test_has_text_without_text_frame(self):
        """Test has_text property without text frame."""
        mock_shape = Mock()
        mock_shape.text_frame = None

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        assert element.has_text is False

    def test_read_text_markdown(self):
        """Test reading text as markdown with bold formatting."""
        mock_shape = Mock()
        mock_text_frame = Mock()
        mock_paragraph = Mock()
        mock_run = Mock()

        # Configure the mock run
        mock_run.text = "Hello World"
        mock_run.font.bold = True
        mock_run.font.italic = False
        mock_run.hyperlink.address = None

        # Configure paragraph with iterable runs
        mock_paragraph.runs = [mock_run]
        mock_paragraph.text.strip.return_value = "Hello World"
        mock_paragraph.level = 0

        mock_text_frame.paragraphs = [mock_paragraph]
        mock_shape.text_frame = mock_text_frame

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        result = element.read_text(as_markdown=True)
        # Bold text should be formatted with markdown
        assert "**Hello World**" in result

    def test_has_text_frame_with_empty_text(self):
        """Test has_text_frame returns True for empty text boxes."""
        mock_shape = Mock()
        mock_shape.text_frame = Mock()
        mock_shape.text_frame.text = ""  # Empty text

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        assert element.has_text_frame is True  # Can contain text
        assert element.has_text is False  # But has no content

    def test_has_text_frame_without_text_frame(self):
        """Test has_text_frame returns False when no text_frame exists."""
        mock_shape = Mock()
        mock_shape.text_frame = None

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        assert element.has_text_frame is False

    @patch("gslides_api.pptx.markdown_to_pptx.apply_markdown_to_textframe")
    def test_write_text_to_empty_text_box(self, mock_apply_md):
        """Test write_text works on empty text boxes (the bug fix)."""
        mock_shape = Mock()
        mock_text_frame = Mock()
        mock_text_frame.text = ""  # Empty text box
        # Mock paragraphs for _extract_base_style_from_textframe (empty paragraphs = no style)
        mock_text_frame.paragraphs = []
        mock_shape.text_frame = mock_text_frame

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        api_client = PowerPointAPIClient()
        element.write_text(api_client, content="New content", autoscale=False)

        # Should have called apply_markdown_to_textframe, not skipped
        # base_style will be None for empty text box (no existing runs to extract from)
        mock_apply_md.assert_called_once_with(
            markdown_text="New content",
            text_frame=mock_text_frame,
            base_style=None,
            autoscale=False,
        )

    @patch("gslides_api.pptx.markdown_to_pptx.apply_markdown_to_textframe")
    def test_write_text_skipped_without_text_frame(self, mock_apply_md):
        """Test write_text is skipped when shape has no text_frame."""
        mock_shape = Mock()
        mock_shape.text_frame = None

        element = PowerPointShapeElement(
            objectId="1", pptx_element=mock_shape, alt_text=AbstractAltText()
        )

        api_client = PowerPointAPIClient()
        element.write_text(api_client, content="New content", autoscale=False)

        # Should NOT have called apply_markdown_to_textframe
        mock_apply_md.assert_not_called()


class TestValidatePptxElement:
    """Test element validation functionality."""

    def test_validate_shape_element(self):
        """Test validating shape element."""
        mock_shape = Mock()
        mock_shape.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        mock_shape.shape_id = 123
        mock_shape.name = "Test Shape"
        mock_shape.element.attrib = {"id": "123"}
        mock_shape._element.getparent.return_value = None

        result = validate_pptx_element(mock_shape)
        assert isinstance(result, PowerPointShapeElement)

    def test_validate_image_element(self):
        """Test validating image element."""
        mock_shape = Mock()
        mock_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        mock_shape.shape_id = 123
        mock_shape.name = "Test Image"
        mock_shape.element.attrib = {"id": "123"}
        mock_shape._element.getparent.return_value = None
        # Picture elements don't have text_frame attribute
        del mock_shape.text_frame

        result = validate_pptx_element(mock_shape)
        assert isinstance(result, PowerPointImageElement)

    def test_validate_table_element(self):
        """Test validating table element."""
        mock_shape = Mock()
        mock_shape.shape_type = MSO_SHAPE_TYPE.TABLE
        mock_shape.shape_id = 123
        mock_shape.name = "Test Table"
        mock_shape.element.attrib = {"id": "123"}
        mock_shape._element.getparent.return_value = None
        # Table elements don't have text_frame attribute
        del mock_shape.text_frame

        result = validate_pptx_element(mock_shape)
        assert isinstance(result, PowerPointTableElement)


class TestAbstractThumbnail:
    """Test AbstractThumbnail model."""

    def test_thumbnail_without_content(self):
        """Test creating thumbnail without content."""
        thumbnail = AbstractThumbnail(
            contentUrl="https://example.com/thumbnail.png",
            width=800,
            height=600,
            mime_type="image/png",
        )
        assert thumbnail.contentUrl == "https://example.com/thumbnail.png"
        assert thumbnail.width == 800
        assert thumbnail.height == 600
        assert thumbnail.mime_type == "image/png"
        assert thumbnail.content is None

    def test_thumbnail_with_content(self):
        """Test creating thumbnail with content."""
        png_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100
        thumbnail = AbstractThumbnail(
            contentUrl="https://example.com/thumbnail.png",
            width=800,
            height=600,
            mime_type="image/png",
            content=png_bytes,
        )
        assert thumbnail.content == png_bytes
        assert isinstance(thumbnail.content, bytes)


class TestPowerPointSlideThumbnail:
    """Test PowerPoint slide thumbnail functionality."""

    def _create_mock_pptx_slide(self):
        """Create a properly mocked python-pptx slide."""
        mock_slide = Mock()
        mock_slide.slide_id = 256
        mock_slide.slide_layout = Mock()
        mock_slide.shapes = []
        mock_slide.has_notes_slide = True
        mock_notes = Mock()
        mock_notes.notes_text_frame = Mock()
        mock_notes.notes_text_frame.text = "Test Slide"
        mock_slide.notes_slide = mock_notes
        return mock_slide

    @patch("gslides_api.adapters.pptx_adapter.render_slide_to_image")
    def test_thumbnail_without_include_data(self, mock_render):
        """Test thumbnail with include_data=False returns content=None."""
        # Create a minimal valid PNG bytes
        png_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100

        mock_render.return_value = png_bytes

        # Create mock slide and presentation
        mock_slide = self._create_mock_pptx_slide()

        mock_prs = Mock()
        mock_prs.slides = [mock_slide]

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            temp_path = temp_file.name

        try:
            # Create PowerPointSlide - the __init__ doesn't pass kwargs to parent
            # So we need to set presentation_id and pptx_presentation manually
            element = PowerPointSlide(pptx_slide=mock_slide)
            element.presentation_id = temp_path
            element.pptx_presentation = mock_prs

            api_client = PowerPointAPIClient()
            thumbnail = element.thumbnail(
                api_client=api_client, size="MEDIUM", include_data=False
            )

            assert thumbnail.content is None
            assert thumbnail.mime_type == "image/png"
            assert thumbnail.width > 0
            assert thumbnail.height > 0
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

    @patch("gslides_api.adapters.pptx_adapter.render_slide_to_image")
    def test_thumbnail_with_include_data(self, mock_render):
        """Test thumbnail with include_data=True returns content with bytes."""
        # Create a minimal valid PNG bytes
        png_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100

        mock_render.return_value = png_bytes

        # Create mock slide and presentation
        mock_slide = self._create_mock_pptx_slide()

        # Create a list-like mock for slides that properly iterates
        mock_prs = Mock()
        mock_prs.slides = [mock_slide]  # Same slide object so slide_id matches

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            temp_path = temp_file.name

        try:
            # Create PowerPointSlide - the __init__ doesn't pass kwargs to parent
            # So we need to set presentation_id and pptx_presentation manually
            element = PowerPointSlide(pptx_slide=mock_slide)
            element.presentation_id = temp_path
            element.pptx_presentation = mock_prs

            # Verify the slide index lookup works
            assert element._get_slide_index() == 0

            api_client = PowerPointAPIClient()
            thumbnail = element.thumbnail(
                api_client=api_client, size="MEDIUM", include_data=True
            )

            assert thumbnail.content is not None
            assert isinstance(thumbnail.content, bytes)
            assert thumbnail.content == png_bytes
            assert thumbnail.mime_type == "image/png"
        finally:
            if os.path.exists(temp_path):
                os.unlink(temp_path)

    def test_thumbnail_placeholder_on_missing_file(self):
        """Test thumbnail returns placeholder when presentation file doesn't exist."""
        mock_slide = self._create_mock_pptx_slide()

        mock_prs = Mock()
        mock_prs.slides = [mock_slide]

        # Create PowerPointSlide with non-existent file path
        element = PowerPointSlide(
            pptx_slide=mock_slide,
            presentation_id="/nonexistent/path.pptx",
            pptx_presentation=mock_prs,
        )

        api_client = PowerPointAPIClient()
        thumbnail = element.thumbnail(
            api_client=api_client, size="MEDIUM", include_data=True
        )

        # Should return placeholder without content
        assert thumbnail.contentUrl == "placeholder_thumbnail.png"
        assert thumbnail.width == 320
        assert thumbnail.height == 240
        assert thumbnail.content is None
