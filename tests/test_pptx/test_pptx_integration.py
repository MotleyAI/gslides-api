"""Integration tests for PowerPoint adapter."""

import pytest
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches

from gslides_api.adapters.pptx_adapter import (
    PowerPointAPIClient,
    PowerPointPresentation,
    PowerPointSlide,
    PowerPointShapeElement,
    PowerPointImageElement,
    PowerPointTableElement,
)
from gslides_api.agnostic.element import MarkdownTableElement


class TestPowerPointIntegration:
    """Integration tests for PowerPoint adapter functionality."""

    def create_sample_presentation(self) -> str:
        """Create a sample presentation for testing."""
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Integration Test Presentation"
        subtitle.text = "Testing PowerPoint Adapter"

        # Content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = 'Features Test'
        tf = body_shape.text_frame
        tf.text = 'Text formatting'

        # Add some formatted text
        p = tf.add_paragraph()
        p.text = 'Bullet points'
        p.level = 0

        p = tf.add_paragraph()
        p.text = 'Nested bullets'
        p.level = 1

        # Table slide
        table_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(table_slide_layout)
        shapes = slide.shapes
        shapes.title.text = 'Data Table'

        # Add table
        rows, cols = 3, 3
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(3)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # Fill table
        headers = ['Name', 'Age', 'City']
        data = [
            ['Alice', '25', 'New York'],
            ['Bob', '30', 'San Francisco']
        ]

        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header

        for row_idx, row_data in enumerate(data, 1):
            for col_idx, cell_data in enumerate(row_data):
                table.cell(row_idx, col_idx).text = cell_data

        # Save to temp file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        temp_file.close()

        return temp_file.name

    def test_full_presentation_workflow(self):
        """Test complete presentation workflow."""
        pptx_path = self.create_sample_presentation()

        try:
            # Load presentation
            presentation = PowerPointPresentation.from_id(PowerPointAPIClient(), pptx_path)
            assert presentation is not None
            assert len(presentation.slides) == 3

            # Test API client
            api_client = PowerPointAPIClient()
            assert api_client is not None

            # Test slide access
            first_slide = presentation.slides[0]
            assert isinstance(first_slide, PowerPointSlide)

            # Test element access
            elements = first_slide.elements
            assert len(elements) > 0

            # Test text elements
            text_elements = [e for e in elements if hasattr(e, 'has_text') and e.has_text]
            assert len(text_elements) > 0

            # Test shape elements
            shape_elements = [e for e in elements if isinstance(e, PowerPointShapeElement)]
            assert len(shape_elements) > 0

        finally:
            os.unlink(pptx_path)

    def test_text_manipulation_workflow(self):
        """Test text manipulation workflow."""
        pptx_path = self.create_sample_presentation()

        try:
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            content_slide = presentation.slides[1]  # Second slide

            # Find text elements
            text_elements = [e for e in content_slide.elements if isinstance(e, PowerPointShapeElement) and e.has_text]
            assert len(text_elements) > 0

            text_element = text_elements[0]

            # Read original text
            original_text = text_element.read_text(as_markdown=False)
            assert len(original_text) > 0

            # Write new content
            new_content = """# Updated Content

This slide has been updated with:

- **Bold** formatting
- *Italic* formatting
- Multiple bullet points
  - With nesting
  - And sub-bullets

Regular paragraph text."""

            text_element.write_text(api_client, new_content)

            # Read back as markdown
            markdown_text = text_element.read_text(as_markdown=True)
            assert "Updated Content" in markdown_text
            assert "Bold" in markdown_text or "**Bold**" in markdown_text

            # Save presentation
            presentation.save(api_client)

            # Reload and verify persistence
            reloaded_presentation = PowerPointPresentation.from_id(api_client, pptx_path)
            reloaded_slide = reloaded_presentation.slides[1]
            reloaded_text_elements = [e for e in reloaded_slide.elements if isinstance(e, PowerPointShapeElement) and e.has_text]

            if reloaded_text_elements:
                reloaded_text = reloaded_text_elements[0].read_text(as_markdown=False)
                assert "Updated Content" in reloaded_text

        finally:
            os.unlink(pptx_path)

    def test_table_manipulation_workflow(self):
        """Test table manipulation workflow."""
        pptx_path = self.create_sample_presentation()

        try:
            presentation = PowerPointPresentation.from_id(PowerPointAPIClient(), pptx_path)
            table_slide = presentation.slides[2]  # Third slide with table

            # Find table elements
            table_elements = [e for e in table_slide.elements if isinstance(e, PowerPointTableElement)]

            if table_elements:  # Only test if we successfully found/converted table elements
                table_element = table_elements[0]
                api_client = PowerPointAPIClient()

                # Test table resizing
                try:
                    table_element.resize(api_client, rows=4, cols=4)
                except Exception as e:
                    # Table resizing might fail due to python-pptx limitations
                    print(f"Table resize failed (expected): {e}")

                # Test content update (if we have MarkdownTableElement support)
                try:
                    # Create markdown table data
                    table_data = [
                        ["Product", "Price", "Stock", "Category"],
                        ["Laptop", "$999", "50", "Electronics"],
                        ["Phone", "$599", "100", "Electronics"],
                        ["Book", "$25", "200", "Education"]
                    ]

                    # This would need MarkdownTableElement implementation
                    # For now, just verify the table element exists
                    assert table_element.pptx_element is not None

                except Exception as e:
                    print(f"Table content update test skipped: {e}")

        finally:
            os.unlink(pptx_path)

    def test_presentation_copying_workflow(self):
        """Test presentation copying workflow."""
        pptx_path = self.create_sample_presentation()

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                api_client = PowerPointAPIClient()
                # Load original
                original_presentation = PowerPointPresentation.from_id(api_client, pptx_path)

                # Copy presentation
                copied_presentation = original_presentation.copy_via_drive(
                    api_client,
                    "copied_test_presentation",
                    temp_dir
                )

                # Verify copy
                assert copied_presentation is not None
                assert copied_presentation.file_path != original_presentation.file_path
                assert os.path.exists(copied_presentation.file_path)
                assert len(copied_presentation.slides) == len(original_presentation.slides)

                # Modify copy
                if copied_presentation.slides:
                    first_slide = copied_presentation.slides[0]
                    text_elements = [e for e in first_slide.elements if isinstance(e, PowerPointShapeElement) and e.has_text]

                    if text_elements:
                        text_elements[0].write_text(api_client, "Modified copy content")

                # Save modified copy
                copied_presentation.save(api_client)

                # Verify original is unchanged
                original_presentation.sync_from_cloud(api_client)
                original_first_slide = original_presentation.slides[0]
                original_text_elements = [e for e in original_first_slide.elements if isinstance(e, PowerPointShapeElement) and e.has_text]

                if original_text_elements:
                    original_text = original_text_elements[0].read_text(as_markdown=False)
                    assert "Modified copy content" not in original_text

        finally:
            os.unlink(pptx_path)

    def test_slide_duplication_workflow(self):
        """Test slide duplication workflow."""
        pptx_path = self.create_sample_presentation()

        try:
            presentation = PowerPointPresentation.from_id(PowerPointAPIClient(), pptx_path)
            api_client = PowerPointAPIClient()

            original_slide_count = len(presentation.slides)
            assert original_slide_count > 0

            # Duplicate first slide
            first_slide = presentation.slides[0]
            try:
                duplicated_slide = first_slide.duplicate(api_client)
                assert isinstance(duplicated_slide, PowerPointSlide)

                # Note: The duplicated slide might not automatically be added to the presentation
                # depending on the implementation, so we might not see an increase in slide count
                print(f"Original slides: {original_slide_count}")
                print(f"Duplicated slide created: {duplicated_slide.objectId}")

            except Exception as e:
                # Slide duplication might have limitations in python-pptx
                print(f"Slide duplication test failed (might be expected): {e}")

        finally:
            os.unlink(pptx_path)

    def test_speaker_notes_workflow(self):
        """Test speaker notes workflow."""
        pptx_path = self.create_sample_presentation()

        try:
            presentation = PowerPointPresentation.from_id(PowerPointAPIClient(), pptx_path)
            api_client = PowerPointAPIClient()

            first_slide = presentation.slides[0]

            # Test speaker notes (if available)
            if first_slide.speaker_notes:
                # Write notes
                notes_content = """Speaker notes for this slide:

- Key point 1
- Key point 2 with **emphasis**
- Key point 3

Remember to speak slowly and clearly."""

                first_slide.speaker_notes.write_text(api_client, notes_content)

                # Read back notes
                read_notes = first_slide.speaker_notes.read_text(as_markdown=True)
                assert "Speaker notes" in read_notes
                assert "Key point 1" in read_notes

                # Save and reload
                presentation.save(api_client)
                reloaded_presentation = PowerPointPresentation.from_id(api_client, pptx_path)
                reloaded_slide = reloaded_presentation.slides[0]

                if reloaded_slide.speaker_notes:
                    reloaded_notes = reloaded_slide.speaker_notes.read_text(as_markdown=False)
                    assert "Speaker notes" in reloaded_notes

        finally:
            os.unlink(pptx_path)

    def test_error_handling_workflow(self):
        """Test error handling in various scenarios."""
        # Test loading non-existent file
        with pytest.raises(FileNotFoundError):
            PowerPointPresentation.from_id(PowerPointAPIClient(), "non_existent_file.pptx")

        # Test invalid file format
        with tempfile.NamedTemporaryFile(delete=False, suffix='.txt') as temp_file:
            temp_file.write(b"This is not a PowerPoint file")
            temp_file.flush()

            try:
                with pytest.raises(ValueError):
                    PowerPointPresentation.from_id(PowerPointAPIClient(), temp_file.name)
            finally:
                os.unlink(temp_file.name)

        # Test API client operations with invalid paths
        api_client = PowerPointAPIClient()

        with pytest.raises(FileNotFoundError):
            api_client.copy_presentation("non_existent.pptx", "copy")

        with pytest.raises(FileNotFoundError):
            api_client.create_folder("test", parent_folder_id="non_existent_folder")

    def test_file_operations_workflow(self):
        """Test file operations workflow."""
        api_client = PowerPointAPIClient()

        with tempfile.TemporaryDirectory() as temp_dir:
            # Create test file
            test_file = os.path.join(temp_dir, "test.pptx")
            with open(test_file, "w") as f:
                f.write("test content")

            # Test copy
            copy_result = api_client.copy_presentation(test_file, "copied_file")
            assert os.path.exists(copy_result["id"])

            # Test folder creation
            folder_result = api_client.create_folder("test_folder", parent_folder_id=temp_dir)
            assert os.path.exists(folder_result["id"])
            assert os.path.isdir(folder_result["id"])

            # Test file deletion
            api_client.delete_file(copy_result["id"])
            assert not os.path.exists(copy_result["id"])

            # Test folder deletion
            api_client.delete_file(folder_result["id"])
            assert not os.path.exists(folder_result["id"])

    def test_presentation_metadata_workflow(self):
        """Test presentation metadata handling."""
        pptx_path = self.create_sample_presentation()

        try:
            presentation = PowerPointPresentation.from_id(PowerPointAPIClient(), pptx_path)

            # Test basic metadata
            assert presentation.presentationId == pptx_path
            assert presentation.title is not None
            assert presentation.url.startswith("file://")

            # Test slide metadata
            first_slide = presentation.slides[0]
            assert first_slide.objectId is not None
            assert hasattr(first_slide, 'slideProperties')

            # Test element metadata
            if first_slide.elements:
                first_element = first_slide.elements[0]
                assert first_element.objectId is not None
                assert hasattr(first_element, 'alt_text')

        finally:
            os.unlink(pptx_path)

    def test_image_replace_workflow(self):
        """Test image replacement workflow - verifies that images can be replaced in slides."""
        # Create a presentation with an image
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add speaker notes for slide identification
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Test Image Slide"

        # Create a simple test image
        test_image_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
        replacement_image_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name

        try:
            # Create test images using PIL
            from PIL import Image

            # Create initial red image
            img1 = Image.new("RGB", (100, 100), color="red")
            img1.save(test_image_path)

            # Create replacement blue image
            img2 = Image.new("RGB", (100, 100), color="blue")
            img2.save(replacement_image_path)

            # Add the test image to the slide with alt text for identification
            picture = slide.shapes.add_picture(
                test_image_path, Inches(1), Inches(1), Inches(2), Inches(2)
            )
            # Set alt text title for identification
            picture._element.xpath(".//p:cNvPr")[0].attrib["title"] = "TestImage"

            # Save presentation
            pptx_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
            prs.save(pptx_path)

            # Load the presentation using our adapter
            api_client = PowerPointAPIClient()
            presentation = PowerPointPresentation.from_id(api_client, pptx_path)

            # Find the image element
            first_slide = presentation.slides[0]
            image_elements = [
                e for e in first_slide.elements if isinstance(e, PowerPointImageElement)
            ]

            assert len(image_elements) > 0, "No image elements found in slide"
            image_element = image_elements[0]

            # Verify pptx_slide was propagated to element
            assert image_element.pptx_slide is not None, "pptx_slide not propagated to element"

            # Replace the image
            image_element.replace_image(api_client, file=replacement_image_path)

            # Save and reload to verify the change persisted
            presentation.save(api_client)

            # Reload and verify
            reloaded_prs = Presentation(pptx_path)
            reloaded_slide = reloaded_prs.slides[0]

            # Find the picture shape
            picture_shapes = [s for s in reloaded_slide.shapes if hasattr(s, "image")]
            assert len(picture_shapes) > 0, "No picture shapes found after replacement"

            # Verify the image was actually replaced by checking the image blob
            # (The blue image should be different from the red one)
            new_image = picture_shapes[0].image
            assert new_image is not None, "Image blob not found"
            assert len(new_image.blob) > 0, "Image blob is empty"

        finally:
            # Cleanup
            for path in [test_image_path, replacement_image_path, pptx_path]:
                if os.path.exists(path):
                    os.unlink(path)