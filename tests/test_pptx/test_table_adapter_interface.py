"""
Tests for table adapter interface conformance.

These tests verify that all table adapters (PowerPoint, HTML, GSlides) correctly
implement the AbstractTableElement interface:
- resize() returns float (font scale factor)
- update_content() accepts font_scale_factor parameter
- get_column_count() returns int
"""

import pytest
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches

from gslides_api.agnostic.element import MarkdownTableElement

from gslides_api.adapters.html_adapter import (
    HTMLAPIClient,
    HTMLTableElement,
)
from gslides_api.adapters.pptx_adapter import (
    PowerPointAPIClient,
    PowerPointTableElement,
)

# ============================================================================
# Fixtures
# ============================================================================

@pytest.fixture
def pptx_api_client():
    """Create a PowerPointAPIClient instance."""
    return PowerPointAPIClient()


@pytest.fixture
def html_api_client():
    """Create an HTMLAPIClient instance."""
    return HTMLAPIClient()


@pytest.fixture
def pptx_table_element():
    """Create a PowerPoint table element for testing."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)

    table_shape = slide.shapes.add_table(rows=3, cols=3, left=left, top=top, width=width, height=height)
    return PowerPointTableElement(pptx_element=table_shape)


@pytest.fixture
def html_table_element():
    """Create an HTML table element for testing."""
    html_content = """
    <table>
        <thead>
            <tr><th>A</th><th>B</th><th>C</th></tr>
        </thead>
        <tbody>
            <tr><td>1</td><td>2</td><td>3</td></tr>
            <tr><td>4</td><td>5</td><td>6</td></tr>
        </tbody>
    </table>
    """
    soup = BeautifulSoup(html_content, 'lxml')
    table_tag = soup.find('table')
    return HTMLTableElement(html_element=table_tag, objectId="test-html-table")


@pytest.fixture
def markdown_table_content():
    """Create a MarkdownTableElement for testing update_content."""
    markdown_input = """| A | B | C |
|---|---|---|
| 1 | 2 | 3 |
| 4 | 5 | 6 |"""
    return MarkdownTableElement(name="Test Table", content=markdown_input)


# ============================================================================
# PowerPointTableElement Tests
# ============================================================================

class TestPowerPointTableElementInterface:
    """Test PowerPointTableElement conforms to AbstractTableElement interface."""

    def test_resize_returns_float(self, pptx_api_client, pptx_table_element):
        """Test that resize() returns a float (font scale factor)."""
        result = pptx_table_element.resize(
            api_client=pptx_api_client,
            rows=4,
            cols=3,
            fix_width=True,
            fix_height=True,
        )

        assert result is not None, "resize() should return a value, not None"
        assert isinstance(result, float), f"resize() should return float, got {type(result)}"
        assert result > 0, "Font scale factor should be positive"

    def test_update_content_accepts_font_scale_factor(self, pptx_api_client, pptx_table_element, markdown_table_content):
        """Test that update_content() accepts font_scale_factor parameter."""
        # This should not raise TypeError about unexpected keyword argument
        pptx_table_element.update_content(
            api_client=pptx_api_client,
            markdown_content=markdown_table_content,
            check_shape=False,
            font_scale_factor=0.8,
        )

    def test_get_column_count_returns_int(self, pptx_table_element):
        """Test that get_column_count() exists and returns int."""
        result = pptx_table_element.get_column_count()

        assert isinstance(result, int), f"get_column_count() should return int, got {type(result)}"
        assert result == 3, f"Expected 3 columns, got {result}"

    def test_get_row_count_returns_int(self, pptx_table_element):
        """Test that get_row_count() exists and returns int (sanity check)."""
        result = pptx_table_element.get_row_count()

        assert isinstance(result, int), f"get_row_count() should return int, got {type(result)}"
        assert result == 3, f"Expected 3 rows, got {result}"


# ============================================================================
# HTMLTableElement Tests
# ============================================================================

class TestHTMLTableElementInterface:
    """Test HTMLTableElement conforms to AbstractTableElement interface."""

    def test_resize_returns_float(self, html_api_client, html_table_element):
        """Test that resize() returns a float (font scale factor)."""
        result = html_table_element.resize(
            api_client=html_api_client,
            rows=4,
            cols=3,
            fix_width=True,
            fix_height=True,
        )

        assert result is not None, "resize() should return a value, not None"
        assert isinstance(result, float), f"resize() should return float, got {type(result)}"
        assert result > 0, "Font scale factor should be positive"

    def test_update_content_accepts_font_scale_factor(self, html_api_client, html_table_element, markdown_table_content):
        """Test that update_content() accepts font_scale_factor parameter."""
        # This should not raise TypeError about unexpected keyword argument
        html_table_element.update_content(
            api_client=html_api_client,
            markdown_content=markdown_table_content,
            check_shape=False,
            font_scale_factor=0.8,
        )

    def test_get_column_count_returns_int(self, html_table_element):
        """Test that get_column_count() exists and returns int."""
        result = html_table_element.get_column_count()

        assert isinstance(result, int), f"get_column_count() should return int, got {type(result)}"
        assert result == 3, f"Expected 3 columns, got {result}"

    def test_get_row_count_returns_int(self, html_table_element):
        """Test that get_row_count() exists and returns int (sanity check)."""
        result = html_table_element.get_row_count()

        assert isinstance(result, int), f"get_row_count() should return int, got {type(result)}"
        assert result == 2, f"Expected 2 rows (tbody only), got {result}"


# ============================================================================
# Edge Cases
# ============================================================================

class TestTableAdapterEdgeCases:
    """Test edge cases for table adapter interface."""

    def test_pptx_resize_with_none_element(self, pptx_api_client):
        """Test resize returns float even with invalid element."""
        elem = PowerPointTableElement()
        elem.pptx_element = None

        result = elem.resize(
            api_client=pptx_api_client,
            rows=2,
            cols=2,
        )

        assert isinstance(result, float), f"resize() should return float even with None element, got {type(result)}"

    def test_html_resize_with_none_element(self, html_api_client):
        """Test resize returns float even with invalid element."""
        elem = HTMLTableElement(objectId="test-empty")
        elem.html_element = None

        result = elem.resize(
            api_client=html_api_client,
            rows=2,
            cols=2,
        )

        assert isinstance(result, float), f"resize() should return float even with None element, got {type(result)}"

    def test_pptx_get_column_count_with_none_element(self):
        """Test get_column_count returns 0 with invalid element."""
        elem = PowerPointTableElement()
        elem.pptx_element = None

        result = elem.get_column_count()

        assert result == 0, f"get_column_count() should return 0 for None element, got {result}"

    def test_html_get_column_count_with_none_element(self):
        """Test get_column_count returns 0 with invalid element."""
        elem = HTMLTableElement(objectId="test-empty")
        elem.html_element = None

        result = elem.get_column_count()

        assert result == 0, f"get_column_count() should return 0 for None element, got {result}"


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
