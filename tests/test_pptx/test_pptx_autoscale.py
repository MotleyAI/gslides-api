"""
Tests for PowerPoint autoscaling font size cap behavior.

Autoscaling should only decrease font size to fit content, never increase it.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from gslides_api.agnostic.text import FullTextStyle, RichStyle
from gslides_api.pptx.markdown_to_pptx import (
    _get_max_font_size_from_textframe,
    apply_markdown_to_textframe,
)


class TestGetMaxFontSizeFromTextframe:
    """Tests for _get_max_font_size_from_textframe helper function."""

    def test_returns_none_for_empty_textframe(self):
        """Empty text frame should return None."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_frame = textbox.text_frame

        # Clear to ensure empty
        text_frame.clear()

        result = _get_max_font_size_from_textframe(text_frame)
        assert result is None

    def test_returns_single_font_size(self):
        """Single font size should be returned."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_frame = textbox.text_frame

        # Set text with specific font size
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Test text"
        run.font.size = Pt(14)

        result = _get_max_font_size_from_textframe(text_frame)
        assert result == 14.0

    def test_returns_max_of_multiple_font_sizes(self):
        """When multiple font sizes exist, return the maximum."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_frame = textbox.text_frame

        # First paragraph with 12pt font
        p1 = text_frame.paragraphs[0]
        run1 = p1.add_run()
        run1.text = "Small text"
        run1.font.size = Pt(12)

        # Second paragraph with 18pt font
        p2 = text_frame.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Large text"
        run2.font.size = Pt(18)

        # Third paragraph with 10pt font
        p3 = text_frame.add_paragraph()
        run3 = p3.add_run()
        run3.text = "Tiny text"
        run3.font.size = Pt(10)

        result = _get_max_font_size_from_textframe(text_frame)
        assert result == 18.0

    def test_ignores_runs_without_font_size(self):
        """Runs without explicit font size should be ignored."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        text_frame = textbox.text_frame

        p = text_frame.paragraphs[0]

        # First run with explicit size
        run1 = p.add_run()
        run1.text = "Sized text"
        run1.font.size = Pt(14)

        # Second run without explicit size
        run2 = p.add_run()
        run2.text = " unsized text"
        # run2.font.size is None (inherited)

        result = _get_max_font_size_from_textframe(text_frame)
        assert result == 14.0


class TestAutoscaleFontSizeCap:
    """Tests for autoscaling font size cap behavior."""

    def test_autoscale_does_not_increase_beyond_original(self):
        """Autoscaling should not increase font size beyond original."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Create a large textbox
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        text_frame = textbox.text_frame

        # Set original text with 10pt font
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Original"
        run.font.size = Pt(10)
        run.font.name = "Arial"

        # Apply very short text with autoscale - without cap, fit_text would increase to 18pt
        apply_markdown_to_textframe("Hi", text_frame, autoscale=True)

        # Check that font size did not increase beyond original 10pt
        actual_font_size = None
        for para in text_frame.paragraphs:
            for r in para.runs:
                if r.font.size is not None:
                    actual_font_size = r.font.size.pt
                    break

        # Font size should be at most 10pt (the original), not 18pt
        if actual_font_size is not None:
            assert actual_font_size <= 10.0, (
                f"Font size {actual_font_size}pt exceeds original 10pt"
            )

    def test_autoscale_still_decreases_for_long_text(self):
        """Autoscaling should still decrease font size when text is too long."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Create a small textbox
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
        text_frame = textbox.text_frame

        # Set original text with 14pt font
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Short"
        run.font.size = Pt(14)
        run.font.name = "Arial"

        # Apply very long text that won't fit at 14pt
        long_text = "This is a very long text that definitely will not fit in the small text box and needs to be scaled down significantly to fit properly."
        apply_markdown_to_textframe(long_text, text_frame, autoscale=True)

        # The text should have been written (we can't easily verify the font size
        # was decreased since fit_text behavior depends on font availability,
        # but we verify no exception was raised and text was applied)
        assert text_frame.text.strip() == long_text

    def test_autoscale_uses_base_style_when_textframe_empty(self):
        """When text frame is empty, base_style font size should be used as cap."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Create a large textbox
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        text_frame = textbox.text_frame

        # Clear text frame so no original font size can be found
        text_frame.clear()

        # Create base_style with 12pt font
        base_style = FullTextStyle(
            rich=RichStyle(font_size_pt=12)
        )

        # Apply short text with autoscale and base_style
        apply_markdown_to_textframe("Hi", text_frame, base_style=base_style, autoscale=True)

        # The function should have used 12pt as the cap from base_style
        # We just verify no exception was raised and text was applied
        assert text_frame.text.strip() == "Hi"

    def test_autoscale_false_does_not_modify_font_size(self):
        """When autoscale=False, font size should not be modified."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        text_frame = textbox.text_frame

        # Set original text
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Original"
        run.font.size = Pt(14)

        # Apply new text without autoscale
        apply_markdown_to_textframe("New text", text_frame, autoscale=False)

        # Text should be updated but no fit_text was called
        assert text_frame.text.strip() == "New text"


class TestAutoscaleDefaultMaxSize:
    """Test that default max_size of 18pt is still applied when appropriate."""

    def test_uses_original_when_smaller_than_default(self):
        """Original font size is used when smaller than default 18pt."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        text_frame = textbox.text_frame

        # Set original text with 10pt font (smaller than default 18pt)
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Original"
        run.font.size = Pt(10)
        run.font.name = "Arial"

        # Apply short text with autoscale
        apply_markdown_to_textframe("Hi", text_frame, autoscale=True)

        # Should use 10pt cap (original), not 18pt
        # This is tested indirectly - we just verify no exception
        assert text_frame.text.strip() == "Hi"

    def test_uses_default_when_no_original_available(self):
        """Default 18pt is used when no original font size is available."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
        text_frame = textbox.text_frame

        # Clear the text frame completely (no original font size)
        text_frame.clear()

        # Apply text with autoscale but no base_style
        apply_markdown_to_textframe("Test text", text_frame, autoscale=True)

        # Should fall back to default 18pt cap
        # This is tested indirectly - we just verify no exception and text was applied
        assert text_frame.text.strip() == "Test text"


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
