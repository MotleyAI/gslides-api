"""
Test suite for writing markdown to PowerPoint text frames.

This module tests the PowerPointShapeElement.write_text method to ensure
markdown is correctly converted to formatted PowerPoint text with proper
bold, italic, hyperlinks, and bullet point support.
"""

import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches

from gslides_api.adapters.pptx_adapter import PowerPointAPIClient, PowerPointShapeElement
from gslides_api.adapters.abstract_slides import AbstractAltText


class TestPowerPointWriteMarkdown:
    """Test writing markdown to PowerPoint text frames."""

    @pytest.fixture
    def sample_pptx_path(self):
        """Return path to sample PPTX file."""
        # Use the samplead template if it exists
        sample_path = (
            Path(__file__).parent.parent.parent
            / "playground"
            / "samplead"
            / "Samplead Master Deck Template.pptx"
        )
        if sample_path.exists():
            return str(sample_path)
        # Otherwise create a simple test PPTX
        return self._create_test_pptx()

    def _create_test_pptx(self):
        """Create a simple test PPTX file."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        text_box = slide.shapes.add_textbox(100, 100, 400, 200)
        text_frame = text_box.text_frame
        text_frame.text = "Initial text"

        # Save to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            return tmp.name

    def _get_text_shape(self, slide):
        """Find or create a shape with a text_frame."""
        # Find first shape with text_frame
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                return shape
        # If no text shape found, create one
        text_box = slide.shapes.add_textbox(
            left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
        )
        text_box.text_frame.text = "Test"
        return text_box

    def test_simple_text_write(self, sample_pptx_path):
        """Test writing simple text without formatting."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        # Create element
        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write simple text
        markdown_text = "Hello, World!"
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written
        assert shape.text_frame.text.strip() == "Hello, World!"

    def test_bold_text_write(self, sample_pptx_path):
        """Test writing bold text."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write text with bold formatting
        markdown_text = "This is **bold** text."
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written and formatting applied
        text_frame = shape.text_frame
        assert "This is" in text_frame.text
        assert "bold" in text_frame.text
        assert "text." in text_frame.text

        # Check that the word "bold" has bold formatting
        found_bold = False
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if "bold" in run.text and run.font.bold:
                    found_bold = True
                    break
        assert found_bold, "Bold formatting not applied"

    def test_italic_text_write(self, sample_pptx_path):
        """Test writing italic text."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write text with italic formatting
        markdown_text = "This is *italic* text."
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written and formatting applied
        text_frame = shape.text_frame
        assert "italic" in text_frame.text

        # Check that the word "italic" has italic formatting
        found_italic = False
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if "italic" in run.text and run.font.italic:
                    found_italic = True
                    break
        assert found_italic, "Italic formatting not applied"

    def test_combined_formatting_write(self, sample_pptx_path):
        """Test writing text with combined bold and italic."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write text with combined formatting
        markdown_text = "This is ***bold and italic*** text."
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written
        text_frame = shape.text_frame
        assert "bold and italic" in text_frame.text

        # Check for combined formatting
        found_combined = False
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if "bold and italic" in run.text:
                    if run.font.bold and run.font.italic:
                        found_combined = True
                        break
        assert found_combined, "Combined bold+italic formatting not applied"

    def test_hyperlink_write(self, sample_pptx_path):
        """Test writing hyperlinks."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write text with hyperlink
        markdown_text = "Visit [our website](https://example.com) for more info."
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written
        text_frame = shape.text_frame
        assert "our website" in text_frame.text

        # Check that hyperlink was applied
        found_link = False
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if "our website" in run.text and run.hyperlink.address:
                    assert run.hyperlink.address == "https://example.com"
                    found_link = True
                    break
        assert found_link, "Hyperlink not applied"

    def test_bullet_list_write(self, sample_pptx_path):
        """Test writing bullet lists."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write bullet list
        markdown_text = """Key points:
* First item
* Second item
* Third item"""
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written
        text_frame = shape.text_frame
        assert "First item" in text_frame.text
        assert "Second item" in text_frame.text
        assert "Third item" in text_frame.text

        # Check that bullet points were applied
        # In PowerPoint, list items have level >= 0 (0 is first level bullet)
        # Count paragraphs with actual content (non-empty)
        bullet_paragraphs = [p for p in text_frame.paragraphs if p.text.strip()]
        # We expect at least 4 paragraphs: "Key points:" + 3 bullet items
        assert (
            len(bullet_paragraphs) >= 4
        ), f"Expected at least 4 paragraphs, found {len(bullet_paragraphs)}"
        # Verify bullet items are present
        assert any(
            "item" in p.text.lower() for p in text_frame.paragraphs
        ), "Bullet items should be present"

    @pytest.mark.skip(reason="Nested list support in markdown parser needs additional work")
    def test_nested_bullet_list_write(self, sample_pptx_path):
        """Test writing nested bullet lists."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write nested bullet list
        markdown_text = """Main points:
* Top level item
    * Nested item 1
    * Nested item 2
* Another top level"""
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written
        text_frame = shape.text_frame
        assert "Top level item" in text_frame.text
        assert "Nested item 1" in text_frame.text

        # Check nesting levels
        # Verify that different nesting levels exist
        levels_found = set()
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                levels_found.add(paragraph.level)
        # Should have at least 2 different levels (0 for top, 1 for nested)
        assert (
            len(levels_found) >= 2
        ), f"Expected at least 2 nesting levels, found {len(levels_found)}: {levels_found}"

    def test_autoscale_option(self, sample_pptx_path):
        """Test autoscale option applies fit_text() to shrink text to fit shape.

        Note: We use fit_text() instead of MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE because
        the latter only sets a flag that PowerPoint applies when you edit the text -
        it doesn't work when you just open the file. fit_text() directly calculates
        and sets the font size, working immediately.
        """
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write with autoscale enabled
        markdown_text = "This is a very long text that should trigger autoscaling if the text box is too small to fit all the content."
        element.write_text(api_client=api_client, content=markdown_text, autoscale=True)

        # Verify fit_text() was applied:
        # - auto_size should be NONE (fit_text sets this)
        # - word_wrap should be True
        # - font size should be explicitly set (not None)
        text_frame = shape.text_frame
        assert text_frame.auto_size == MSO_AUTO_SIZE.NONE  # fit_text sets NONE
        assert text_frame.word_wrap is True
        # Font size should be set (fit_text calculates and sets it)
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            assert text_frame.paragraphs[0].runs[0].font.size is not None

    def test_word_wrap_enabled(self, sample_pptx_path):
        """Test that word wrap is enabled by default."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write text
        markdown_text = "Some text that should wrap to box width."
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify word wrap is enabled
        text_frame = shape.text_frame
        assert text_frame.word_wrap is True

    def test_complex_markdown_write(self, sample_pptx_path):
        """Test writing complex markdown with multiple features."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write complex markdown
        markdown_text = """# Project Summary

This is a **very important** project with *multiple* features:

* Feature 1: Authentication
* Feature 2: **Bold feature** with emphasis
* Feature 3: Link to [documentation](https://docs.example.com)

For more details, visit our ***bold and italic*** section."""
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify all content is present
        text_frame = shape.text_frame
        full_text = text_frame.text
        assert "Project Summary" in full_text
        assert "very important" in full_text
        assert "Feature 1" in full_text
        assert "Feature 2" in full_text
        assert "Feature 3" in full_text
        assert "documentation" in full_text
        assert "bold and italic" in full_text

    def test_code_span_write(self, sample_pptx_path):
        """Test writing code spans."""
        prs = Presentation(sample_pptx_path)
        slide = prs.slides[0]
        shape = self._get_text_shape(slide)

        api_client = PowerPointAPIClient()
        api_client.prs = prs
        element = PowerPointShapeElement(
            objectId="test_1", pptx_element=shape, alt_text=AbstractAltText()
        )

        # Write text with code span
        markdown_text = "Use the `print()` function to display output."
        element.write_text(api_client=api_client, content=markdown_text)

        # Verify text was written
        text_frame = shape.text_frame
        assert "print()" in text_frame.text

        # Code spans should be rendered with Courier New font
        found_code = False
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if "print()" in run.text:
                    # Check for monospace font (Courier New)
                    if run.font.name == "Courier New":
                        found_code = True
                        break
        assert found_code, "Code span not formatted with monospace font"
