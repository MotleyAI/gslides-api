"""
Concrete implementation of abstract slides using python-pptx.
This module provides the actual implementation that maps abstract slide operations to python-pptx calls.
"""

import io
import logging
import os
import shutil
import tempfile
from typing import Annotated, Any, List, Optional, Union

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import NotesSlide, Slide, SlideLayout
from pptx.table import Table
from pptx.text.text import TextFrame
from pptx.util import Emu, Inches, Pt
from pydantic import BaseModel, ConfigDict, Discriminator, Field, Tag, TypeAdapter, model_validator
from typing_extensions import TypedDict

from gslides_api.agnostic.domain import ImageData
from gslides_api.agnostic.element import MarkdownTableElement, TableData
from gslides_api.agnostic.ir import FormattedDocument, FormattedParagraph, FormattedTextRun
from gslides_api.agnostic.ir_to_markdown import ir_to_markdown
from gslides_api.agnostic.text import (
    AbstractColor,
    FullTextStyle,
    MarkdownRenderableStyle,
    RichStyle,
)
from gslides_api.agnostic.units import EMU_PER_INCH, OutputUnit, from_emu, to_emu

from gslides_api.common.download import download_binary_file
from gslides_api.pptx.converters import _paragraph_has_bullet, pptx_table_to_markdown
from gslides_api.pptx.slide_copier import _remove_layout_placeholders
from gslides_api.pptx.chart_renderer import render_slide_to_image

from gslides_api.adapters.abstract_slides import (
    AbstractAltText,
    AbstractCredentials,
    AbstractElement,
    AbstractElementKind,
    AbstractImageElement,
    AbstractPresentation,
    AbstractShapeElement,
    AbstractSize,
    AbstractSlide,
    AbstractSlideProperties,
    AbstractSlidesAPIClient,
    AbstractSpeakerNotes,
    AbstractTableElement,
    AbstractThumbnail,
)

logger = logging.getLogger(__name__)

# Type alias for python-pptx Presentation to avoid name collision with our class
PptxPresentation = Presentation

# Union type for discriminator and model_validator inputs
PptxElementInput = Union[BaseShape, dict, "PowerPointElementParent"]


# TypedDict for style information returned by styles() method
class StyleInfo(TypedDict, total=False):
    """Style information for a text run."""

    text: str
    bold: Optional[bool]
    italic: Optional[bool]
    font_name: Optional[str]
    font_size: Optional[int]
    color_rgb: Optional[tuple]


# Map MSO_THEME_COLOR index to XML element names in theme color scheme
_THEME_COLOR_MAP = {
    1: "dk1",  # Dark 1 (typically text/background dark)
    2: "lt1",  # Light 1 (typically text/background light)
    3: "dk2",  # Dark 2
    4: "lt2",  # Light 2
    5: "accent1",  # Accent 1
    6: "accent2",  # Accent 2
    7: "accent3",  # Accent 3
    8: "accent4",  # Accent 4
    9: "accent5",  # Accent 5
    10: "accent6",  # Accent 6
    11: "hlink",  # Hyperlink
    12: "folHlink",  # Followed Hyperlink
}


def _textframe_to_ir(text_frame: TextFrame) -> FormattedDocument:
    """Convert a PowerPoint TextFrame to platform-agnostic IR.

    This enables using the shared ir_to_markdown function which handles
    run consolidation and proper space placement outside markdown markers.

    Args:
        text_frame: The python-pptx TextFrame to convert

    Returns:
        FormattedDocument with paragraphs and text runs
    """
    from gslides_api.pptx.converters import _paragraph_has_bullet

    if not text_frame:
        return FormattedDocument()

    elements = []

    for paragraph in text_frame.paragraphs:
        runs = []
        has_bullet = _paragraph_has_bullet(paragraph)
        level = getattr(paragraph, "level", 0) or 0

        for run in paragraph.runs:
            # Build the style from run properties
            md_style = MarkdownRenderableStyle(
                bold=run.font.bold or False,
                italic=run.font.italic or False,
                hyperlink=(
                    run.hyperlink.address if run.hyperlink and run.hyperlink.address else None
                ),
            )
            style = FullTextStyle(markdown=md_style)

            runs.append(FormattedTextRun(content=run.text, style=style))

        # Create paragraph
        para = FormattedParagraph(runs=runs)

        # Handle bullet points by prepending indent and bullet marker
        if has_bullet and runs:
            indent = "  " * level
            # Prepend bullet marker to first run's content
            if para.runs:
                para.runs[0].content = f"{indent}- {para.runs[0].content}"

        elements.append(para)

    return FormattedDocument(elements=elements)


def _resolve_font_color_rgb(run, slide: Optional[Slide] = None) -> Optional[tuple]:
    """Resolve font color to RGB tuple, handling both direct RGB and theme colors.

    Args:
        run: A python-pptx text run object
        slide: Optional slide for resolving theme colors. If not provided,
               theme colors cannot be resolved and will return None.

    Returns:
        RGB tuple (r, g, b) with values 0-255, or None if color cannot be resolved.
    """
    try:
        color = run.font.color
        if color.type is None:
            return None

        # Direct RGB color - straightforward case
        if color.type == MSO_COLOR_TYPE.RGB:
            rgb = color.rgb
            if rgb:
                return (rgb[0], rgb[1], rgb[2])

        # Theme/scheme color - need to resolve from theme XML
        elif color.type == MSO_COLOR_TYPE.SCHEME:
            theme_color_idx = color.theme_color
            if theme_color_idx is None or slide is None:
                return None

            try:
                # Navigate: slide -> slide_layout -> slide_master -> theme
                slide_master = slide.slide_layout.slide_master
                theme_part = slide_master.part.part_related_by(RT.THEME)
                theme_xml = parse_xml(theme_part.blob)

                color_name = _THEME_COLOR_MAP.get(theme_color_idx)
                if color_name:
                    # XPath to find the srgbClr value for this theme color
                    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                    # Try srgbClr first (explicit RGB), then sysClr (system color)
                    xpath_srgb = f".//a:clrScheme/a:{color_name}//a:srgbClr/@val"
                    results = theme_xml.xpath(xpath_srgb, namespaces=ns)
                    if results:
                        hex_val = results[0]
                        return (
                            int(hex_val[0:2], 16),
                            int(hex_val[2:4], 16),
                            int(hex_val[4:6], 16),
                        )
                    # Try system color (e.g., windowText, window)
                    xpath_sys = f".//a:clrScheme/a:{color_name}//a:sysClr/@lastClr"
                    results = theme_xml.xpath(xpath_sys, namespaces=ns)
                    if results:
                        hex_val = results[0]
                        return (
                            int(hex_val[0:2], 16),
                            int(hex_val[2:4], 16),
                            int(hex_val[4:6], 16),
                        )
            except Exception:
                pass  # Fall through to None

    except Exception:
        pass

    return None


def _extract_base_style_from_textframe(
    text_frame: TextFrame,
    preserve_bold_italic: bool = False,
    slide: Optional[Slide] = None,
) -> Optional[FullTextStyle]:
    """Extract FullTextStyle from the first non-empty text run.

    Used to preserve font formatting when replacing template variables.
    Extracts font size, color, font family, and underline from the first
    text run that contains non-whitespace text.

    Args:
        text_frame: The PowerPoint text frame to extract style from
        preserve_bold_italic: If True, extract bold/italic from template formatting.
            Set to True when content doesn't have markdown formatting (no * characters),
            so the original template's bold/italic is preserved.
            Set to False (default) when content has markdown formatting, so that
            **bold** and *italic* syntax controls the formatting instead.
        slide: Optional slide reference for resolving theme colors to RGB.
            If provided, theme colors (like accent1, dk1, etc.) will be resolved
            to their actual RGB values from the presentation theme.

    Returns:
        FullTextStyle if a styled run was found, None otherwise
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue

            # Extract RichStyle properties
            font_size_pt = None
            if run.font.size is not None:
                font_size_pt = run.font.size.pt

            font_family = run.font.name

            foreground_color = None
            try:
                # Use helper to resolve both RGB and theme colors
                rgb_tuple = _resolve_font_color_rgb(run, slide=slide)
                if rgb_tuple:
                    foreground_color = AbstractColor.from_rgb_tuple(rgb_tuple)
            except Exception:
                pass

            rich = RichStyle(
                font_family=font_family,
                font_size_pt=font_size_pt,
                foreground_color=foreground_color,
                underline=run.font.underline or False,
            )

            # Extract bold/italic based on preserve_bold_italic flag
            if preserve_bold_italic:
                # Preserve template's bold/italic when content has no markdown formatting
                # Try python-pptx API first, fall back to XML attributes
                bold = run.font.bold
                italic = run.font.italic
                if bold is None:
                    # python-pptx may return None for inherited values, check XML
                    try:
                        rPr = run._r.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                        )
                        if rPr is not None:
                            bold = rPr.get("b") == "1"
                    except Exception:
                        bold = False
                if italic is None:
                    try:
                        rPr = run._r.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                        )
                        if rPr is not None:
                            italic = rPr.get("i") == "1"
                    except Exception:
                        italic = False
                markdown = MarkdownRenderableStyle(
                    bold=bold or False,
                    italic=italic or False,
                )
            else:
                # Let markdown content determine bold/italic
                # This allows mixed formatting like "Text **bold** more text"
                markdown = MarkdownRenderableStyle(
                    bold=False,
                    italic=False,
                )

            return FullTextStyle(rich=rich, markdown=markdown)

    return None


def pptx_element_discriminator(v: PptxElementInput) -> str:
    """Discriminator to determine which PowerPointElement subclass based on type field."""
    # First check if it's a direct pptx shape with shape_type
    if hasattr(v, "shape_type"):
        if hasattr(v, "text_frame"):
            return "shape"
        if v.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX]:
            return "shape"
        elif v.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif v.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        else:
            return "generic"

    # Then check if it's already wrapped with pptx_element
    elif hasattr(v, "pptx_element"):
        pptx_element = v.pptx_element
        if hasattr(pptx_element, "shape_type"):
            shape_type = pptx_element.shape_type
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "shape"
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                return "image"
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                return "table"
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(pptx_element, "text_frame"):
                    return "shape"
                return "generic"
        return "generic"

    # Finally check for type field
    else:
        element_type = getattr(v, "type", None)
        if element_type in [AbstractElementKind.SHAPE, "SHAPE", "shape"]:
            return "shape"
        elif element_type in [AbstractElementKind.IMAGE, "IMAGE", "image"]:
            return "image"
        elif element_type in [AbstractElementKind.TABLE, "TABLE", "table"]:
            return "table"
        else:
            return "generic"


class PowerPointAPIClient(AbstractSlidesAPIClient):
    """PowerPoint API client implementation."""

    def __init__(self):
        # No initialization needed for filesystem-based operations
        self._auto_flush = True

    @property
    def auto_flush(self):
        return self._auto_flush

    @auto_flush.setter
    def auto_flush(self, value: bool):
        # Just store this for consistency with abstract interface
        self._auto_flush = value

    def flush_batch_update(self):
        # No batching needed for filesystem operations
        pass

    def copy_presentation(
        self, presentation_id: str, copy_title: str, folder_id: Optional[str] = None
    ) -> dict:
        """Copy a presentation file to another location."""
        if not os.path.exists(presentation_id):
            raise FileNotFoundError(f"Presentation file not found: {presentation_id}")

        # Determine destination folder
        if folder_id is None:
            # Copy to same folder as source
            dest_folder = os.path.dirname(presentation_id)
        else:
            # Validate folder exists
            if not os.path.exists(folder_id) or not os.path.isdir(folder_id):
                raise FileNotFoundError(f"Destination folder not found: {folder_id}")
            dest_folder = folder_id

        if copy_title.endswith(".pptx"):
            copy_title = copy_title[:-5]

        # Create destination path
        base_name = os.path.splitext(os.path.basename(presentation_id))[0]
        dest_path = os.path.join(dest_folder, f"{copy_title}.pptx")

        # Copy file
        shutil.copy2(presentation_id, dest_path)

        return {
            "id": dest_path,
            "name": copy_title,
            "parents": [dest_folder] if folder_id else [os.path.dirname(presentation_id)],
        }

    def create_folder(
        self, name: str, ignore_existing: bool = True, parent_folder_id: Optional[str] = None
    ) -> dict:
        """Create a folder in the filesystem."""
        if parent_folder_id is None:
            parent_folder_id = os.getcwd()

        if not os.path.exists(parent_folder_id):
            raise FileNotFoundError(f"Parent folder not found: {parent_folder_id}")

        folder_path = os.path.join(parent_folder_id, name)

        try:
            os.makedirs(folder_path, exist_ok=ignore_existing)
        except FileExistsError:
            if not ignore_existing:
                raise

        return {"id": folder_path, "name": name, "parents": [parent_folder_id]}

    def delete_file(self, file_id: str):
        """Delete a file from the filesystem."""
        if os.path.exists(file_id):
            if os.path.isdir(file_id):
                shutil.rmtree(file_id)
            else:
                os.remove(file_id)

    def set_credentials(self, credentials: AbstractCredentials):
        # Do nothing as this is filesystem-based
        pass

    def get_credentials(self) -> Optional[AbstractCredentials]:
        # Return None as no credentials needed for filesystem operations
        return None

    def replace_text(
        self, slide_ids: list[str], match_text: str, replace_text: str, presentation_id: str
    ):
        """Replace text across specified slides in a presentation."""
        if not os.path.exists(presentation_id):
            raise FileNotFoundError(f"Presentation file not found: {presentation_id}")

        prs = Presentation(presentation_id)

        # Convert slide_ids to indices (assuming they are string indices)
        slide_indices = []
        for slide_id in slide_ids:
            try:
                slide_indices.append(int(slide_id))
            except ValueError:
                # If slide_id is not a number, skip it
                continue

        # Replace text in specified slides
        for slide_idx in slide_indices:
            if slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                self._replace_text_in_slide(slide, match_text, replace_text)

        # Save the presentation
        prs.save(presentation_id)

    def _replace_text_in_slide(self, slide: Slide, match_text: str, replace_text: str):
        """Replace text in all shapes of a slide."""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if match_text in run.text:
                            run.text = run.text.replace(match_text, replace_text)

    @classmethod
    def get_default_api_client(cls) -> "PowerPointAPIClient":
        """Get the default API client instance."""
        return cls()

    async def get_presentation_as_pdf(self, presentation_id: str) -> bytes:
        """Get PDF from presentation."""
        raise NotImplementedError("PDF export not implemented for PowerPoint")


class PowerPointSpeakerNotes(AbstractSpeakerNotes):
    """PowerPoint speaker notes (belonging to a particular slide) implementation."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    notes_slide: Any = Field(exclude=True, default=None)

    def __init__(self, notes_slide: NotesSlide, **kwargs):
        super().__init__(**kwargs)
        self.notes_slide = notes_slide

    def read_text(self, as_markdown: bool = True) -> str:
        """Read text from speaker notes."""
        if not self.notes_slide or not self.notes_slide.notes_text_frame:
            return ""

        text_frame = self.notes_slide.notes_text_frame
        if as_markdown:
            return self._convert_to_markdown(text_frame)
        else:
            return text_frame.text

    def write_text(self, api_client: "PowerPointAPIClient", content: str):
        """Write text to speaker notes."""
        if not self.notes_slide:
            return

        # Get or create text frame
        text_frame = self.notes_slide.notes_text_frame
        if not text_frame:
            return

        # Clear existing content
        text_frame.clear()

        # Add content as paragraphs (simple implementation - could be enhanced for markdown)
        lines = content.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                text_frame.text = line
            else:
                p = text_frame.add_paragraph()
                p.text = line

    def _convert_to_markdown(self, text_frame: TextFrame) -> str:
        """Convert text frame content to markdown format."""
        markdown_lines = []

        for paragraph in text_frame.paragraphs:
            line_parts = []
            for run in paragraph.runs:
                text = run.text
                if run.font.bold:
                    text = f"**{text}**"
                if run.font.italic:
                    text = f"*{text}*"
                line_parts.append(text)

            line = "".join(line_parts)
            if line.strip():  # Only add non-empty lines
                markdown_lines.append(line)

        return "\n".join(markdown_lines)


class PowerPointElementParent(AbstractElement):
    """Generic concrete element for PowerPoint slide elements."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    pptx_element: Any = Field(exclude=True, default=None)
    pptx_slide: Any = Field(exclude=True, default=None)

    @model_validator(mode="before")
    @classmethod
    def convert_from_pptx_element(cls, data: PptxElementInput) -> dict:
        """Convert from python-pptx element to our abstract representation."""
        if isinstance(data, dict):
            # Already converted
            return data
        elif isinstance(data, BaseShape):
            pptx_element = data
        elif hasattr(data, "pptx_element"):
            # Already wrapped element
            return data.__dict__
        else:
            raise ValueError(f"Expected BaseShape, got {type(data)}")

        # Extract basic properties
        object_id = str(getattr(pptx_element, "shape_id", ""))
        if not object_id and hasattr(pptx_element, "element"):
            element_attr = getattr(pptx_element, "element", {})
            if hasattr(element_attr, "attrib"):
                object_id = element_attr.attrib.get("id", "")
            elif isinstance(element_attr, dict):
                object_id = element_attr.get("id", "")

        # Get slide ID if available (from parent slide)
        slide_id = ""
        try:
            if hasattr(pptx_element, "_element") and hasattr(pptx_element._element, "getparent"):
                slide_elem = pptx_element._element.getparent()
                while slide_elem is not None and hasattr(slide_elem, "tag"):
                    if "cSld" in slide_elem.tag:
                        slide_id = slide_elem.attrib.get("name", "")
                        break
                    if hasattr(slide_elem, "getparent"):
                        slide_elem = slide_elem.getparent()
                    else:
                        break
        except Exception:
            # Ignore XML traversal errors
            pass

        # Get alt text if available
        # While waiting for https://github.com/scanny/python-pptx/pull/512 to get merged,
        # hack raw XML to get at the alt text, code inspired by
        # https://stackoverflow.com/questions/63802783/check-if-image-is-decorative-in-powerpoint-using-python-pptx
        cnvpr_elements = pptx_element._element.xpath(".//p:cNvPr")
        alt_text_title = cnvpr_elements[0].attrib.get("title", None) if cnvpr_elements else None
        alt_text_descr = cnvpr_elements[0].attrib.get("descr", None) if cnvpr_elements else None

        return {
            "objectId": object_id,
            "presentation_id": "",  # Will be set by parent
            "slide_id": slide_id,
            "alt_text": AbstractAltText(title=alt_text_title, description=alt_text_descr),
            "type": "generic",
            "pptx_element": pptx_element,
        }

    def absolute_size(self, units: OutputUnit = OutputUnit.IN) -> tuple[float, float]:
        """Get the absolute size of the element."""
        if not self.pptx_element:
            return (0.0, 0.0)

        width_emu = float(self.pptx_element.width)
        height_emu = float(self.pptx_element.height)

        return (from_emu(width_emu, units), from_emu(height_emu, units))

    def absolute_position(self, units: OutputUnit = OutputUnit.IN) -> tuple[float, float]:
        """Get the absolute position of the element."""
        if not self.pptx_element:
            return (0.0, 0.0)

        left_emu = float(self.pptx_element.left)
        top_emu = float(self.pptx_element.top)

        return (from_emu(left_emu, units), from_emu(top_emu, units))

    def create_image_element_like(
        self, api_client: PowerPointAPIClient
    ) -> "PowerPointImageElement":
        """Create an image element with the same properties as this element.

        This is used when replacing non-image elements (like charts) with images.
        Creates a placeholder image at the same position/size as this element,
        removes this element from the slide, and returns the new image element.
        """
        if not self.pptx_slide or not self.pptx_element:
            logger.warning("Cannot create image element: missing slide or element reference")
            raise ValueError("Cannot create image element without slide reference")

        # Get position and size from current element
        left = self.pptx_element.left
        top = self.pptx_element.top
        width = self.pptx_element.width
        height = self.pptx_element.height

        # Create a small transparent placeholder image
        # (will be replaced immediately by the actual image via replace_image)
        placeholder_img = Image.new("RGBA", (100, 100), (255, 255, 255, 0))
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        try:
            placeholder_img.save(temp_file.name, "PNG")
            temp_file.close()

            # For chart elements, clean up the chart relationship before removal
            # to prevent orphaned chart files from causing corruption
            self._cleanup_chart_relationship()

            # Remove current element (chart) from the slide
            logger.debug(
                f"create_image_element_like: removing element {self.alt_text.title}, type={type(self.pptx_element)}"
            )
            if hasattr(self.pptx_element, "_element") and hasattr(
                self.pptx_element._element, "getparent"
            ):
                parent = self.pptx_element._element.getparent()
                if parent is not None:
                    parent.remove(self.pptx_element._element)
                    logger.debug(f"create_image_element_like: removed element from parent")
                else:
                    logger.warning(f"create_image_element_like: element has no parent")
            else:
                logger.warning(f"create_image_element_like: element has no _element or getparent")

            # Add new picture at the same position/size
            new_picture = self.pptx_slide.shapes.add_picture(
                image_file=temp_file.name, left=left, top=top, width=width, height=height
            )

            # Create and return a PowerPointImageElement wrapper
            image_element = PowerPointImageElement(
                objectId=str(new_picture.shape_id),
                alt_text=self.alt_text,
                pptx_element=new_picture,
                pptx_slide=self.pptx_slide,
            )
            return image_element

        finally:
            # Clean up temp file
            if os.path.exists(temp_file.name):
                os.unlink(temp_file.name)

    def set_alt_text(
        self,
        api_client: PowerPointAPIClient,
        title: str | None = None,
        description: str | None = None,
    ):
        """Set alt text for the element."""
        if self.pptx_element and (title is not None or description is not None):
            # Set the XML title/descr attributes (the actual alt text used for element matching)
            # This must match how alt text is read in convert_from_pptx_element
            try:
                cnvpr = self.pptx_element._element.xpath(".//p:cNvPr")
                if cnvpr:
                    if title is not None:
                        cnvpr[0].attrib["title"] = title
                    if description is not None:
                        cnvpr[0].attrib["descr"] = description
            except Exception as e:
                logger.debug(f"Could not set alt text in XML: {e}")

            # Also set the name property for compatibility
            if title is not None and hasattr(self.pptx_element, "name"):
                self.pptx_element.name = title

            # Update our alt_text model
            if title is not None:
                self.alt_text.title = title
            if description is not None:
                self.alt_text.description = description

    def _cleanup_chart_relationship(self):
        """Remove orphaned chart relationship when replacing chart with image.

        When a chart graphicFrame is removed from a slide to be replaced with an image,
        the chart relationship and associated files (chartN.xml, styles, colors, Excel)
        would otherwise remain in the package as orphaned references, potentially
        causing file corruption.
        """
        try:
            if not hasattr(self.pptx_element, "has_chart") or not self.pptx_element.has_chart:
                return

            if not hasattr(self.pptx_element, "chart"):
                return

            chart = self.pptx_element.chart
            if not hasattr(chart, "part"):
                return

            chart_part = chart.part
            slide_part = self.pptx_slide.part

            # Find and remove the relationship pointing to this chart
            for rId, rel in list(slide_part.rels.items()):
                if rel._target is chart_part:
                    slide_part.drop_rel(rId)
                    logger.debug(f"Removed orphaned chart relationship {rId}")
                    break

        except Exception as e:
            logger.warning(f"Could not clean up chart relationship: {e}")


class PowerPointShapeElement(AbstractShapeElement, PowerPointElementParent):
    """PowerPoint shape element implementation, in particular a text box."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    pptx_element: Any = Field(exclude=True, default=None)

    @model_validator(mode="before")
    @classmethod
    def convert_from_pptx_element(cls, data: PptxElementInput) -> dict:
        """Convert from python-pptx shape element."""
        base_data = PowerPointElementParent.convert_from_pptx_element(data)
        base_data["type"] = AbstractElementKind.SHAPE
        return base_data

    @property
    def has_text_frame(self) -> bool:
        """Check if the shape can contain text (has a text_frame)."""
        if not self.pptx_element:
            return False
        return hasattr(self.pptx_element, "text_frame") and self.pptx_element.text_frame is not None

    @property
    def has_text(self) -> bool:
        """Check if the shape has text content."""
        return self.has_text_frame and bool(self.pptx_element.text_frame.text)

    def write_text(
        self,
        api_client: PowerPointAPIClient,
        content: str,
        autoscale: bool = False,
    ):
        """Write text to the shape, parsing markdown and applying formatting."""
        if not self.has_text_frame:
            return

        from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

        text_frame = self.pptx_element.text_frame

        # Check if content has markdown formatting indicators (* for bold/italic)
        # If not, preserve the template's bold/italic formatting
        has_markdown_formatting = "*" in content

        # Extract base style from existing text to preserve formatting
        # (font size, color, bold, italic, etc.) when replacing template variables
        base_style = _extract_base_style_from_textframe(
            text_frame, preserve_bold_italic=not has_markdown_formatting
        )

        # Use shared markdown parser to convert markdown to formatted PowerPoint text
        apply_markdown_to_textframe(
            markdown_text=content,
            text_frame=text_frame,
            base_style=base_style,
            autoscale=autoscale,
        )

    def read_text(self, as_markdown: bool = True) -> str:
        """Read text from the shape."""
        if not self.has_text:
            return ""

        text_frame = self.pptx_element.text_frame
        if as_markdown:
            return self._convert_text_frame_to_markdown(text_frame)
        else:
            return text_frame.text

    def _convert_text_frame_to_markdown(self, text_frame: TextFrame) -> str:
        """Convert PowerPoint text frame to markdown format.

        Uses the shared ir_to_markdown function which handles:
        1. Consolidation of adjacent runs with identical formatting
        2. Proper placement of spaces outside markdown markers

        This fixes the bug where adjacent bold runs produce `****` instead of
        a single consolidated bold section.
        """
        if not text_frame:
            return ""

        # Convert to platform-agnostic IR
        ir_doc = _textframe_to_ir(text_frame)

        # Convert IR to markdown using the shared function
        return ir_to_markdown(ir_doc)

    def styles(self, skip_whitespace: bool = True) -> Optional[List[StyleInfo]]:
        """Extract style information from text runs.

        Returns a list of style dictionaries containing formatting information
        (bold, italic, color, font_name, font_size) separate from text content.
        This follows the gslides-api pattern of separating style from content.

        Args:
            skip_whitespace: If True, skip runs that are only whitespace

        Returns:
            List of style dictionaries, or None if element has no text
        """
        if not self.has_text:
            return None

        text_frame = self.pptx_element.text_frame
        styles = []

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Skip whitespace-only runs if requested
                if skip_whitespace and not run.text.strip():
                    continue

                # Extract style information from the run
                style_info = {
                    "text": run.text,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "font_name": run.font.name,
                    "font_size": run.font.size,
                }

                # Extract color if available (handles both direct RGB and theme colors)
                style_info["color_rgb"] = _resolve_font_color_rgb(run, slide=self.pptx_slide)

                styles.append(style_info)

        return styles if styles else None


class PowerPointImageElement(AbstractImageElement, PowerPointElementParent):
    """PowerPoint image element implementation."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    pptx_element: Any = Field(exclude=True, default=None)

    @model_validator(mode="before")
    @classmethod
    def convert_from_pptx_element(cls, data: PptxElementInput) -> dict:
        """Convert from python-pptx image element."""
        base_data = PowerPointElementParent.convert_from_pptx_element(data)
        base_data["type"] = AbstractElementKind.IMAGE
        return base_data

    def replace_image(
        self,
        api_client: PowerPointAPIClient,
        file: str | None = None,
        url: str | None = None,
    ):
        """Replace the image in this element."""
        logger.debug(f"replace_image called: pptx_element={type(self.pptx_element)}, file={file}")
        if not self.pptx_element:
            logger.warning("replace_image: pptx_element is None")
            return
        if not isinstance(self.pptx_element, Picture):
            logger.warning(
                f"replace_image: pptx_element is not Picture, it's {type(self.pptx_element)}"
            )
            return

        if file and os.path.exists(file):
            # Get current position and size
            left = self.pptx_element.left
            top = self.pptx_element.top
            width = self.pptx_element.width
            height = self.pptx_element.height
            logger.debug(f"replace_image: position=({left}, {top}), size=({width}, {height})")

            # Get the slide containing this shape
            slide = self._get_parent_slide()
            if not slide:
                logger.warning(
                    f"Cannot replace image: no parent slide reference for element {self.alt_text.title}"
                )
                return

            logger.debug(f"replace_image: removing old picture from slide {slide}")
            # Remove current picture
            self._remove_from_slide()

            # Add new picture with same position and size
            logger.debug(f"replace_image: adding new picture from {file}")
            new_picture = slide.shapes.add_picture(file, left, top, width, height)
            self.pptx_element = new_picture
            logger.info(f"replace_image: successfully replaced image with {file}")

        elif url:
            # For URL images, download first using utility with retries
            content, _ = download_binary_file(url)
            # Save to temp file and use file path method
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_file:
                temp_file.write(content)
                temp_file_path = temp_file.name

            try:
                self.replace_image(api_client, file=temp_file_path)
            finally:
                os.unlink(temp_file_path)

    def get_image_data(self):
        """Get the image data from the PowerPoint element."""
        if not self.pptx_element or not isinstance(self.pptx_element, Picture):
            return None

        try:
            # Extract the image data from the picture element
            image_part = self.pptx_element.image
            return ImageData(
                content=image_part.blob,
                mime_type=image_part.content_type,
            )
        except AttributeError:
            return None

    def _get_parent_slide(self):
        """Get the slide that contains this shape."""
        return self.pptx_slide

    def _remove_from_slide(self):
        """Remove this shape from its parent slide."""
        if hasattr(self.pptx_element, "_element") and hasattr(
            self.pptx_element._element, "getparent"
        ):
            parent = self.pptx_element._element.getparent()
            if parent is not None:
                parent.remove(self.pptx_element._element)


class PowerPointTableElement(AbstractTableElement, PowerPointElementParent):
    """PowerPoint table element implementation."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    pptx_element: Any = Field(exclude=True, default=None)

    @model_validator(mode="before")
    @classmethod
    def convert_from_pptx_element(cls, data: PptxElementInput) -> dict:
        """Convert from python-pptx table element."""
        base_data = PowerPointElementParent.convert_from_pptx_element(data)
        base_data["type"] = AbstractElementKind.TABLE
        return base_data

    def resize(
        self,
        api_client: PowerPointAPIClient,
        rows: int,
        cols: int,
        fix_width: bool = True,
        fix_height: bool = True,
        target_height_in: float | None = None,
    ) -> float:
        """Resize the table to the specified dimensions.

        Args:
            target_height_in: If provided, constrain total table height to this value in inches.

        Returns:
            Font scale factor (1.0 since PPTX doesn't support font scaling during resize).
        """
        if not self.pptx_element or not isinstance(self.pptx_element, GraphicFrame):
            return 1.0

        try:
            table = self.pptx_element.table
            current_rows = len(table.rows)
            current_cols = len(table.columns)

            # Adjust rows
            if rows > current_rows:
                for _ in range(rows - current_rows):
                    self._add_table_row(table)
            elif rows < current_rows:
                # Remove excess rows (not directly supported, need to work with XML)
                for i in range(current_rows - 1, rows - 1, -1):
                    self._remove_table_row(table, i)

            # Adjust columns
            if cols > current_cols:
                for _ in range(cols - current_cols):
                    self._add_table_column(table)

                # Copy header row (row 0) styling from rightmost existing cell to new columns
                if len(table.rows) > 0 and current_cols > 0:
                    # Get the rightmost existing header cell (before columns were added)
                    rightmost_header_cell = table.cell(0, current_cols - 1)

                    # Copy fill/background to new header cells
                    for new_col_idx in range(current_cols, cols):
                        new_header_cell = table.cell(0, new_col_idx)

                        # Copy cell fill (background color)
                        try:
                            src_fill = rightmost_header_cell.fill
                            dst_fill = new_header_cell.fill
                            if src_fill.type is not None:
                                # Copy solid fill
                                if (
                                    hasattr(src_fill, "fore_color")
                                    and src_fill.fore_color
                                    and src_fill.fore_color.type is not None
                                ):
                                    dst_fill.solid()
                                    if src_fill.fore_color.type == MSO_COLOR_TYPE.RGB:
                                        dst_fill.fore_color.rgb = src_fill.fore_color.rgb
                                    elif src_fill.fore_color.type == MSO_COLOR_TYPE.SCHEME:
                                        dst_fill.fore_color.theme_color = (
                                            src_fill.fore_color.theme_color
                                        )
                        except Exception as e:
                            logger.debug(f"Could not copy header cell fill: {e}")

            elif cols < current_cols:
                # Remove excess columns (not directly supported, need to work with XML)
                for i in range(current_cols - 1, cols - 1, -1):
                    self._remove_table_column(table, i)

            # Optionally adjust dimensions
            if fix_width:
                # Distribute width evenly among columns
                total_width = self.pptx_element.width
                col_width = total_width // cols
                for col in table.columns:
                    col.width = col_width

            if fix_height or target_height_in is not None:
                # Distribute height evenly among rows
                # Use target_height_in if provided, otherwise use current element height
                if target_height_in is not None:
                    # Convert inches to EMU (914400 EMU = 1 inch)
                    total_height = Inches(target_height_in)
                else:
                    total_height = self.pptx_element.height
                row_height = total_height // rows
                for row in table.rows:
                    row.height = row_height

        except Exception as e:
            logger.warning(f"Error resizing table: {e}")

        return 1.0

    def update_content(
        self,
        api_client: PowerPointAPIClient,
        markdown_content: MarkdownTableElement,
        check_shape: bool = True,
        font_scale_factor: float = 1.0,
    ):
        """Update the table content with markdown data.

        Args:
            font_scale_factor: Font scale factor (currently unused for PPTX, but kept for interface conformance).
        """
        if not self.pptx_element or not isinstance(self.pptx_element, GraphicFrame):
            return

        try:
            table = self.pptx_element.table

            # Get table data from markdown content (MarkdownTableElement has TableData in content field)
            if hasattr(markdown_content, "content") and hasattr(
                markdown_content.content, "headers"
            ):
                headers = markdown_content.content.headers
                data_rows = markdown_content.content.rows
                # Combine headers and data rows for the full table
                all_rows_data = [headers] + data_rows
            else:
                # Fallback for old interface
                all_rows_data = markdown_content.rows if hasattr(markdown_content, "rows") else []

            if not all_rows_data:
                return

            # Get slide reference for theme color resolution
            slide = None
            try:
                slide = self.pptx_element.part.slide
            except (AttributeError, TypeError):
                pass  # Fall back to None - theme colors won't be resolved

            # Extract canonical styles BEFORE resizing to ensure consistency
            # Header style from row 0, body style from row 1 (or row 0 if only 1 row)
            header_style = None
            body_style = None

            if len(table.rows) > 0 and len(table.columns) > 0:
                # Extract header style from first cell of first row
                header_cell = table.cell(0, 0)
                header_style = _extract_base_style_from_textframe(
                    header_cell.text_frame,
                    preserve_bold_italic=True,
                    slide=slide,
                )

                # Extract body style from first cell of second row (if exists)
                if len(table.rows) > 1:
                    body_cell = table.cell(1, 0)
                    body_style = _extract_base_style_from_textframe(
                        body_cell.text_frame,
                        preserve_bold_italic=True,
                        slide=slide,
                    )
                else:
                    # Only one row, use header style for body too
                    body_style = header_style

            # Ensure table has enough rows and columns
            required_rows = len(all_rows_data)
            required_cols = len(all_rows_data[0]) if all_rows_data else 0

            if check_shape:
                current_rows = len(table.rows)
                current_cols = len(table.columns)

                if required_rows != current_rows or required_cols != current_cols:
                    self.resize(api_client, required_rows, required_cols)

            # Import markdown conversion utility (same as used for text boxes)
            from gslides_api.pptx.markdown_to_pptx import apply_markdown_to_textframe

            # Fill table with data using canonical styles for consistency
            for row_idx, row_data in enumerate(all_rows_data):
                if row_idx < len(table.rows):
                    # Use header style for row 0, body style for all other rows
                    canonical_style = header_style if row_idx == 0 else body_style

                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < len(table.columns):
                            cell = table.cell(row_idx, col_idx)
                            cell_text = str(cell_data) if cell_data is not None else ""

                            # Check if content has markdown formatting indicators
                            has_markdown_formatting = "*" in cell_text

                            # Use canonical style, but allow markdown to override bold/italic
                            if canonical_style and has_markdown_formatting:
                                # Create a copy without bold/italic so markdown can control them
                                from copy import deepcopy

                                base_style = deepcopy(canonical_style)
                                base_style.bold = False
                                base_style.italic = False
                            else:
                                base_style = canonical_style

                            # Apply markdown text with formatting to cell
                            apply_markdown_to_textframe(
                                markdown_text=cell_text,
                                text_frame=cell.text_frame,
                                base_style=base_style,
                            )

        except Exception as e:
            logger.warning(f"Error updating table content: {e}")

    def _add_table_row(self, table: Table):
        """Add a row to the table by cloning the last row (working with XML)."""
        try:
            from copy import deepcopy

            num_rows = len(table.rows)
            if num_rows == 0:
                logger.warning("Cannot add row to empty table")
                return

            # Clone the last row's XML element (use explicit index, not -1)
            last_row = table.rows[num_rows - 1]._tr
            new_row = deepcopy(last_row)

            # Clear text content from all cells in the new row
            for tc in new_row.iter(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}tc"
            ):
                for txBody in tc.iter(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}txBody"
                ):
                    for p in txBody.iter(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
                    ):
                        for t in list(
                            p.iter(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
                            )
                        ):
                            t.text = ""

            # Append the new row to the table
            last_row.getparent().append(new_row)
        except Exception as e:
            logger.warning(f"Error adding table row: {e}")

    def _add_table_column(self, table: Table):
        """Add a column to the table (working with XML)."""
        try:
            from copy import deepcopy

            num_cols = len(table.columns)
            if num_cols == 0:
                logger.warning("Cannot add column to empty table")
                return

            # Add a new column definition to the grid
            # Clone the last column's gridCol (use explicit index, not -1)
            last_col = table.columns[num_cols - 1]._gridCol
            new_col = deepcopy(last_col)
            last_col.getparent().append(new_col)

            # Add a new cell to each row by cloning the last cell
            for row in table.rows:
                num_cells = len(row.cells)
                if num_cells > 0:
                    last_cell = row.cells[num_cells - 1]._tc
                    new_cell = deepcopy(last_cell)

                    # Clear text content from the new cell
                    for txBody in new_cell.iter(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}txBody"
                    ):
                        for p in txBody.iter(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
                        ):
                            for t in list(
                                p.iter(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
                                )
                            ):
                                t.text = ""

                    # Append the new cell to the row
                    last_cell.getparent().append(new_cell)
        except Exception as e:
            logger.warning(f"Error adding table column: {e}")

    def _remove_table_row(self, table: Table, row_index: int):
        """Remove a row from the table (working with XML)."""
        try:
            if row_index < len(table.rows):
                # python-pptx _Row uses ._tr for the underlying CT_TableRow XML element
                row_element = table.rows[row_index]._tr
                row_element.getparent().remove(row_element)
        except Exception as e:
            logger.warning(f"Error removing table row {row_index}: {e}")

    def _remove_table_column(self, table: Table, col_index: int):
        """Remove a column from the table (working with XML)."""
        try:
            if col_index < len(table.columns):
                # Remove cells from all rows for this column
                # python-pptx _Cell uses ._tc for the underlying CT_TableCell XML element
                for row in table.rows:
                    if col_index < len(row.cells):
                        cell_element = row.cells[col_index]._tc
                        cell_element.getparent().remove(cell_element)

                # Remove column definition
                # python-pptx _Column uses ._gridCol for the underlying CT_TableCol XML element
                col_element = table.columns[col_index]._gridCol
                col_element.getparent().remove(col_element)
        except Exception as e:
            logger.warning(f"Error removing table column {col_index}: {e}")

    def get_horizontal_border_weight(self, units: OutputUnit = OutputUnit.IN) -> float:
        """Get weight of horizontal borders in specified units.

        For PowerPoint tables, borders are handled via table styles and don't have
        a fixed weight that contributes to layout height the same way as Google Slides,
        so we return 0.
        """
        return 0.0

    def get_row_count(self) -> int:
        """Get current number of rows."""
        if not self.pptx_element or not isinstance(self.pptx_element, GraphicFrame):
            return 0
        try:
            return len(self.pptx_element.table.rows)
        except (AttributeError, TypeError):
            return 0

    def get_column_count(self) -> int:
        """Get current number of columns."""
        if not self.pptx_element or not isinstance(self.pptx_element, GraphicFrame):
            return 0
        try:
            return len(self.pptx_element.table.columns)
        except (AttributeError, TypeError):
            return 0

    def to_markdown_element(self, name: str | None = None) -> MarkdownTableElement:
        """Convert PowerPoint table to markdown table element."""

        if not self.pptx_element or not isinstance(self.pptx_element, GraphicFrame):
            raise ValueError("PowerPointTableElement has no valid GraphicFrame element")

        if not hasattr(self.pptx_element, "table"):
            raise ValueError("GraphicFrame does not contain a table")

        # Convert table to markdown using the converters module
        markdown_table_str = pptx_table_to_markdown(self.pptx_element.table)

        # Create MarkdownTableElement with the markdown string
        markdown_elem = MarkdownTableElement(
            name=name or self.alt_text.title or "Table",
            content=markdown_table_str,
        )

        # Add metadata
        metadata = {
            "objectId": self.objectId,
            "rows": len(self.pptx_element.table.rows),
            "columns": len(self.pptx_element.table.columns) if self.pptx_element.table.rows else 0,
        }

        # Add size metadata if available
        if self.pptx_element.width and self.pptx_element.height:
            metadata["size"] = {
                "width": self.pptx_element.width,
                "height": self.pptx_element.height,
            }

        # Add position metadata if available
        if hasattr(self.pptx_element, "left") and hasattr(self.pptx_element, "top"):
            metadata["position"] = {"left": self.pptx_element.left, "top": self.pptx_element.top}

        # Add alt text title if available
        if self.alt_text and self.alt_text.title:
            metadata["title"] = self.alt_text.title

        # Update metadata
        markdown_elem.metadata.update(metadata)

        return markdown_elem


# Discriminated union type for concrete elements
PowerPointElement = Annotated[
    Union[
        Annotated[PowerPointShapeElement, Tag("shape")],
        Annotated[PowerPointImageElement, Tag("image")],
        Annotated[PowerPointTableElement, Tag("table")],
        Annotated[PowerPointElementParent, Tag("generic")],
    ],
    Discriminator(pptx_element_discriminator),
]

# TypeAdapter for validating the discriminated union
_pptx_element_adapter = TypeAdapter(PowerPointElement)


def validate_pptx_element(pptx_element: BaseShape) -> PowerPointElement:
    """Create the appropriate concrete element from a python-pptx element."""
    element_type = pptx_element_discriminator(pptx_element)

    if element_type == "shape":
        data = PowerPointShapeElement.convert_from_pptx_element(pptx_element)
        return PowerPointShapeElement(**data)
    elif element_type == "image":
        data = PowerPointImageElement.convert_from_pptx_element(pptx_element)
        return PowerPointImageElement(**data)
    elif element_type == "table":
        data = PowerPointTableElement.convert_from_pptx_element(pptx_element)
        return PowerPointTableElement(**data)
    else:
        data = PowerPointElementParent.convert_from_pptx_element(pptx_element)
        return PowerPointElementParent(**data)


class PowerPointSlide(AbstractSlide):
    """PowerPoint slide implementation."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    pptx_slide: Any = Field(exclude=True, default=None)
    pptx_presentation: Any = Field(exclude=True, default=None)
    # Store the presentation ID for reference, will be propagated to elements
    presentation_id: Optional[str] = Field(default=None, exclude=True)

    def _propagate_presentation_id(self, presentation_id: Optional[str] = None) -> None:
        """Set presentation_id on all elements."""
        target_id = presentation_id if presentation_id is not None else self.presentation_id
        if target_id is not None:
            for element in self.elements:
                element.presentation_id = target_id
            # Also propagate via page_elements_flat in case of nested elements
            for element in self.page_elements_flat:
                element.presentation_id = target_id

    def _propagate_pptx_slide(self) -> None:
        """Set pptx_slide on all elements so they can perform slide operations."""
        if self.pptx_slide is not None:
            for element in self.elements:
                element.pptx_slide = self.pptx_slide

    def __setattr__(self, name: str, value) -> None:
        """Override setattr to propagate presentation_id when it's set directly."""
        super().__setattr__(name, value)
        # If presentation_id was just set, propagate it to elements
        if name == "presentation_id" and hasattr(self, "elements"):
            self._propagate_presentation_id(value)

    def __init__(self, pptx_slide: Slide, **kwargs):
        # Convert python-pptx elements to abstract elements
        slide_name = pptx_slide.notes_slide.notes_text_frame.text
        logger.info(f"Processing slide {pptx_slide.slide_id} with title {slide_name}")

        elements = []

        for shape in pptx_slide.shapes:
            try:
                pptx_element = validate_pptx_element(shape)
                elements.append(pptx_element)
            except Exception as e:
                print(f"Warning: Could not convert shape {shape}: {e}")
                continue

        # Get speaker notes
        speaker_notes = None
        if pptx_slide.has_notes_slide:
            speaker_notes = PowerPointSpeakerNotes(pptx_slide.notes_slide)

        # Get slide properties
        slide_properties = AbstractSlideProperties(
            isSkipped=False
        )  # PowerPoint doesn't have skip property

        super().__init__(
            elements=elements,
            objectId=str(pptx_slide.slide_id),
            slideProperties=slide_properties,
            speaker_notes=speaker_notes,
        )

        self.pptx_slide = pptx_slide
        self._propagate_pptx_slide()

    @property
    def page_elements_flat(self) -> list[PowerPointElementParent]:
        """Flatten the elements tree into a list."""
        return self.elements

    def _get_slide_index(self) -> int:
        """Get zero-based index of this slide within the presentation.

        Returns:
            Zero-based slide index

        Raises:
            ValueError: If slide not found in presentation
        """
        if not self.pptx_presentation:
            raise ValueError("No presentation reference available")

        slide_id = self.pptx_slide.slide_id
        for i, slide in enumerate(self.pptx_presentation.slides):
            if slide.slide_id == slide_id:
                return i

        raise ValueError(f"Slide {slide_id} not found in presentation")

    def thumbnail(
        self, api_client: PowerPointAPIClient, size: str, include_data: bool = False
    ) -> AbstractThumbnail:
        """Generate a thumbnail of the slide using LibreOffice rendering."""

        # Get slide index within presentation
        try:
            slide_index = self._get_slide_index()
        except ValueError as e:
            logger.warning(f"Cannot determine slide index: {e}")
            # Return placeholder on failure
            return AbstractThumbnail(
                contentUrl="placeholder_thumbnail.png", width=320, height=240, mime_type="image/png"
            )

        # Check if presentation file exists
        if not self.presentation_id or not os.path.exists(self.presentation_id):
            logger.warning(f"Presentation file not found: {self.presentation_id}")
            return AbstractThumbnail(
                contentUrl="placeholder_thumbnail.png", width=320, height=240, mime_type="image/png"
            )

        # Render slide to PNG bytes
        logger.debug(f"Rendering thumbnail for slide {slide_index} from {self.presentation_id}")
        png_bytes = render_slide_to_image(
            presentation_path=self.presentation_id,
            slide_index=slide_index,
            crop_bounds=None,  # Full slide, not cropped
        )

        if png_bytes is None:
            logger.warning(
                f"Failed to render thumbnail for slide {self.objectId} "
                f"(LibreOffice/pdftoppm not available or rendering failed)"
            )
            # Return placeholder on rendering failure
            return AbstractThumbnail(
                contentUrl="placeholder_thumbnail.png", width=320, height=240, mime_type="image/png"
            )

        # Get image dimensions
        try:
            img = Image.open(io.BytesIO(png_bytes))
            width, height = img.size
        except Exception as e:
            logger.warning(f"Failed to read image dimensions: {e}")
            width, height = 320, 240

        # Save to temp file and return file:// URL
        # (Consistent with PowerPointPresentation.url pattern at line 1480)
        temp_path = None
        try:
            with tempfile.NamedTemporaryFile(
                delete=False, suffix=".png", prefix="slide_thumb_"
            ) as tmp:
                tmp.write(png_bytes)
                temp_path = tmp.name

            return AbstractThumbnail(
                contentUrl=f"file://{os.path.abspath(temp_path)}",
                width=width,
                height=height,
                mime_type="image/png",
                content=png_bytes if include_data else None,
            )
        except Exception as e:
            # Clean up temp file if it was created
            if temp_path:
                try:
                    os.unlink(temp_path)
                    logger.debug(f"Cleaned up temp file after error: {temp_path}")
                except Exception:
                    pass
            logger.error(f"Failed to save thumbnail to temp file: {e}")
            return AbstractThumbnail(
                contentUrl="placeholder_thumbnail.png",
                width=width,
                height=height,
                mime_type="image/png",
            )

    def _copy_shape_to_slide(self, source_shape: BaseShape, target_slide: Slide):
        """Copy a shape from one slide to another."""
        if source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Copy text box or auto shape
            if hasattr(source_shape, "text_frame") and source_shape.text_frame:
                # Add text box
                textbox = target_slide.shapes.add_textbox(
                    source_shape.left, source_shape.top, source_shape.width, source_shape.height
                )
                textbox.text_frame.text = source_shape.text_frame.text

        elif source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Copy image - this is complex as we need the image data
            # For now, skip image copying
            print("Image copying not fully implemented")

        elif source_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Copy table
            if hasattr(source_shape, "table"):
                table = source_shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                new_table_shape = target_slide.shapes.add_table(
                    rows,
                    cols,
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height,
                )

                new_table = new_table_shape.table

                # Copy cell contents
                for row_idx in range(rows):
                    for col_idx in range(cols):
                        if row_idx < len(table.rows) and col_idx < len(table.columns):
                            source_cell = table.cell(row_idx, col_idx)
                            target_cell = new_table.cell(row_idx, col_idx)
                            target_cell.text = source_cell.text


class PowerPointPresentation(AbstractPresentation):
    """PowerPoint presentation implementation."""

    model_config = ConfigDict(arbitrary_types_allowed=True)

    pptx_presentation: Any = Field(exclude=True, default=None)
    file_path: Optional[str] = None
    uploaded_url: Optional[str] = None

    def __init__(
        self,
        pptx_presentation: Presentation,
        file_path: Optional[str] = None,
        uploaded_url: Optional[str] = None,
    ):
        # Convert python-pptx slides to abstract slides
        slides = []
        for pptx_slide in pptx_presentation.slides:
            try:
                slide = PowerPointSlide(pptx_slide)
                slides.append(slide)
            except Exception as e:
                print(f"Warning: Could not convert slide: {e}")
                continue

        # Extract presentation metadata
        presentation_id = file_path or "untitled.pptx"
        title = getattr(pptx_presentation.core_properties, "title", None) or os.path.basename(
            presentation_id
        ).replace(".pptx", "")

        super().__init__(
            slides=slides,
            presentationId=presentation_id,
            revisionId=None,  # PowerPoint doesn't have revision IDs like Google Slides
            title=title,
        )

        self.pptx_presentation = pptx_presentation
        self.file_path = file_path
        self.uploaded_url = uploaded_url

        # Propagate presentation_id and pptx_presentation to all slides and their elements
        for slide in self.slides:
            slide.presentation_id = presentation_id
            slide.pptx_presentation = pptx_presentation

    @property
    def url(self) -> str:
        """Return the file path as URL (file system based)."""
        if self.file_path:
            return f"file://{os.path.abspath(self.file_path)}"
        else:
            raise ValueError("No file path specified for presentation")

    def slide_height(self, units: OutputUnit = OutputUnit.IN) -> float:
        """Return slide height in specified units."""
        # python-pptx stores slide dimensions in EMU
        height_emu = float(self.pptx_presentation.slide_height)
        return from_emu(height_emu, units)

    @classmethod
    def from_id(
        cls,
        api_client: PowerPointAPIClient,
        presentation_id: str,
        uploaded_url: Optional[str] = None,
    ) -> "PowerPointPresentation":
        """Load presentation from file path/ID."""
        # In the context of pptx, ID is simply a filename
        if not os.path.exists(presentation_id):
            raise FileNotFoundError(f"Presentation file not found: {presentation_id}")

        try:
            pptx_presentation = Presentation(presentation_id)
            return cls(pptx_presentation, presentation_id, uploaded_url=uploaded_url)
        except Exception as e:
            raise ValueError(f"Could not load presentation from {presentation_id}: {e}")

    def copy_via_drive(
        self,
        api_client: PowerPointAPIClient,
        copy_title: str,
        folder_id: Optional[str] = None,
    ) -> "PowerPointPresentation":
        """Copy presentation to another location."""

        if not self.file_path:
            raise ValueError("Cannot copy presentation without a file path")

        # Use the API client to copy the file
        copy_result = api_client.copy_presentation(self.file_path, copy_title, folder_id)

        # Load the copied presentation
        copied_presentation = PowerPointPresentation.from_id(api_client, copy_result["id"])

        return copied_presentation

    def sync_from_cloud(self, api_client: PowerPointAPIClient):
        """Re-read presentation from filesystem."""
        if not self.file_path or not os.path.exists(self.file_path):
            return

        # Reload the presentation from file
        pptx_presentation = Presentation(self.file_path)

        # Update our internal representation
        self.pptx_presentation = pptx_presentation

        # Rebuild slides
        slides = []
        for pptx_slide in pptx_presentation.slides:
            try:
                slide = PowerPointSlide(pptx_slide)
                slides.append(slide)
            except Exception as e:
                print(f"Warning: Could not convert slide during sync: {e}")
                continue

        self.slides = slides

        # Propagate presentation_id to all slides and their elements
        for slide in self.slides:
            slide.presentation_id = self.presentationId

        # Update metadata
        self.title = getattr(pptx_presentation.core_properties, "title", None) or os.path.basename(
            self.file_path
        ).replace(".pptx", "")

    def _update_app_properties_metadata(self) -> None:
        """
        Update docProps/app.xml to reflect current slide count.

        Python-pptx doesn't update this metadata after slide changes,
        which causes PowerPoint Online to reject the file as corrupted.
        This is called automatically before save().
        """
        from lxml import etree

        pkg = self.pptx_presentation.part.package

        # Find the extended-properties (app.xml) part
        app_part = None
        for rel in pkg._rels.values():
            if "extended-properties" in rel.reltype:
                app_part = rel._target
                break

        if app_part is None:
            return

        # Parse the XML
        root = etree.fromstring(app_part.blob)

        ns = {
            "ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties",
            "vt": "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes",
        }

        actual_slide_count = len(self.pptx_presentation.slides)

        # Fix <Slides> count
        slides_elem = root.find(".//ep:Slides", ns)
        if slides_elem is not None:
            slides_elem.text = str(actual_slide_count)

        # Fix HeadingPairs - find the slide titles count
        heading_pairs = root.find(".//ep:HeadingPairs/vt:vector", ns)
        if heading_pairs is not None:
            variants = heading_pairs.findall("vt:variant", ns)
            for i, var in enumerate(variants):
                lpstr = var.find("vt:lpstr", ns)
                if lpstr is not None and lpstr.text == "Slide Titles":
                    if i + 1 < len(variants):
                        count_var = variants[i + 1]
                        i4 = count_var.find("vt:i4", ns)
                        if i4 is not None:
                            i4.text = str(actual_slide_count)

        # Fix TitlesOfParts - update vector size
        titles_vector = root.find(".//ep:TitlesOfParts/vt:vector", ns)
        if titles_vector is not None:
            new_size = 1 + actual_slide_count  # 1 theme + N slides
            titles_vector.set("size", str(new_size))

            # Remove excess lpstr elements (keep first one which is theme)
            lpstrs = titles_vector.findall("vt:lpstr", ns)
            for lpstr in lpstrs[new_size:]:
                titles_vector.remove(lpstr)

        # Set the modified XML back
        new_blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone="yes")
        app_part._blob = new_blob

    def _renumber_slides_in_file(self, file_path: str) -> None:
        """
        Post-process a saved PPTX file to renumber slides sequentially.

        PowerPoint Online requires slide files to be numbered sequentially starting from 1.
        Python-pptx may create non-sequential slide numbers (e.g., slide16.xml through slide24.xml)
        when slides are copied and originals deleted. This method fixes that.
        """
        import re
        import zipfile

        # Read the existing file
        with zipfile.ZipFile(file_path, "r") as zf:
            file_contents = {}
            for name in zf.namelist():
                file_contents[name] = zf.read(name)

        # Find all slide files and create renumbering map
        slide_pattern = re.compile(r"^ppt/slides/slide(\d+)\.xml$")
        slide_files = []
        for name in file_contents.keys():
            match = slide_pattern.match(name)
            if match:
                slide_files.append((int(match.group(1)), name))

        # Sort by current number and create mapping to sequential numbers
        slide_files.sort(key=lambda x: x[0])
        renumber_map = {}  # old_num -> new_num
        for new_num, (old_num, _) in enumerate(slide_files, start=1):
            if old_num != new_num:
                renumber_map[old_num] = new_num

        if not renumber_map:
            # Slides are already sequential, nothing to do
            return

        logger.info(f"Renumbering slides: {renumber_map}")

        # Create new file contents with renumbered files
        new_contents = {}
        for name, content in file_contents.items():
            new_name = name
            new_content = content

            # Rename slide files
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                match = slide_pattern.match(name)
                if match:
                    old_num = int(match.group(1))
                    if old_num in renumber_map:
                        new_name = f"ppt/slides/slide{renumber_map[old_num]}.xml"

            # Rename slide relationship files
            elif name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels"):
                match = re.match(r"^ppt/slides/_rels/slide(\d+)\.xml\.rels$", name)
                if match:
                    old_num = int(match.group(1))
                    if old_num in renumber_map:
                        new_name = f"ppt/slides/_rels/slide{renumber_map[old_num]}.xml.rels"

            # Rename notesSlide files
            elif name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml"):
                match = re.match(r"^ppt/notesSlides/notesSlide(\d+)\.xml$", name)
                if match:
                    old_num = int(match.group(1))
                    if old_num in renumber_map:
                        new_name = f"ppt/notesSlides/notesSlide{renumber_map[old_num]}.xml"

            # Rename notesSlide relationship files
            elif name.startswith("ppt/notesSlides/_rels/notesSlide") and name.endswith(".xml.rels"):
                match = re.match(r"^ppt/notesSlides/_rels/notesSlide(\d+)\.xml\.rels$", name)
                if match:
                    old_num = int(match.group(1))
                    if old_num in renumber_map:
                        new_name = (
                            f"ppt/notesSlides/_rels/notesSlide{renumber_map[old_num]}.xml.rels"
                        )

            # Update references in XML content
            if isinstance(content, bytes):
                try:
                    text_content = content.decode("utf-8")
                    modified = False

                    # Update slide references in relationships and content types
                    for old_num, new_num in renumber_map.items():
                        # Update slide file references
                        old_ref = f"slide{old_num}.xml"
                        new_ref = f"slide{new_num}.xml"
                        if old_ref in text_content:
                            text_content = text_content.replace(old_ref, new_ref)
                            modified = True

                        # Update notesSlide references
                        old_notes_ref = f"notesSlide{old_num}.xml"
                        new_notes_ref = f"notesSlide{new_num}.xml"
                        if old_notes_ref in text_content:
                            text_content = text_content.replace(old_notes_ref, new_notes_ref)
                            modified = True

                    if modified:
                        new_content = text_content.encode("utf-8")
                    else:
                        new_content = content
                except UnicodeDecodeError:
                    # Binary file, don't modify
                    new_content = content

            new_contents[new_name] = new_content

        # Write the new file
        with zipfile.ZipFile(file_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for name, content in new_contents.items():
                zf.writestr(name, content)

        logger.info(f"Slide renumbering complete. Renumbered {len(renumber_map)} slides.")

    def save(self, api_client: PowerPointAPIClient) -> None:
        """Save/persist all changes made to this presentation."""
        if not self.file_path:
            raise ValueError("No file path specified for saving")

        # Update app.xml metadata before saving (fixes slide count after add/delete)
        self._update_app_properties_metadata()

        # Ensure directory exists
        dir_path = os.path.dirname(self.file_path)
        if dir_path:
            os.makedirs(dir_path, exist_ok=True)

        # Save the presentation
        self.pptx_presentation.save(self.file_path)

        # Renumber slides to be sequential (fixes PowerPoint Online compatibility)
        self._renumber_slides_in_file(self.file_path)

    def save_as(self, api_client: PowerPointAPIClient, file_path: str):
        """Save the presentation to a new file path."""
        # Update app.xml metadata before saving (fixes slide count after add/delete)
        self._update_app_properties_metadata()

        # Ensure directory exists
        dir_path = os.path.dirname(file_path)
        if dir_path:
            os.makedirs(dir_path, exist_ok=True)

        # Save the presentation
        self.pptx_presentation.save(file_path)

        # Renumber slides to be sequential (fixes PowerPoint Online compatibility)
        self._renumber_slides_in_file(file_path)

        # Update our file path
        self.file_path = file_path
        self.presentationId = file_path

    def insert_copy(
        self,
        source_slide: AbstractSlide,
        api_client: PowerPointAPIClient,
        insertion_index: int | None = None,
    ) -> AbstractSlide:
        """Insert a copy of a slide from another presentation into this presentation."""
        if not isinstance(source_slide, PowerPointSlide):
            raise ValueError("Can only copy PowerPointSlide instances")

        # Import the robust slide copier
        from gslides_api.pptx.slide_copier import SlideCopierManager

        # Use the robust slide copier implementation
        with SlideCopierManager(self.pptx_presentation) as copier:
            try:
                # Attempt robust copying with full error handling
                new_pptx_slide = copier.copy_slide_safe(
                    source_slide.pptx_slide,
                    insertion_index=insertion_index,
                    fallback_to_layout_only=True,
                )

                if new_pptx_slide is None:
                    raise RuntimeError("All slide copying strategies failed")

                # Get copy statistics for logging
                stats = copier.get_copy_statistics()
                logger.info(f"Slide copy completed. Stats: {stats}")

                # Create and return a PowerPointSlide wrapper
                new_slide = PowerPointSlide(new_pptx_slide)

                # DEBUG: Log alt_text values from copied slide
                for elem in new_slide.page_elements_flat:
                    logger.debug(
                        f"Copied slide element: alt_text.title='{elem.alt_text.title}' "
                        f"type={type(elem).__name__}"
                    )

                # Propagate presentation_id and pptx_presentation to the new slide and its elements
                new_slide.presentation_id = self.presentationId
                new_slide.pptx_presentation = self.pptx_presentation

                # Update our slides list
                self.slides.append(new_slide)

                # Set parent references for the new slide and its elements
                new_slide._parent_presentation = self
                for element in new_slide.elements:
                    element._parent_presentation = self

                return new_slide

            except Exception as e:
                logger.error(f"Robust slide copying failed: {e}")
                # Fall back to the original simple implementation as last resort
                return self._insert_copy_fallback(source_slide, api_client, insertion_index)

    def _insert_copy_fallback(
        self,
        source_slide: PowerPointSlide,
        api_client: PowerPointAPIClient,
        insertion_index: int | None = None,
    ) -> AbstractSlide:
        """Fallback implementation using the original simple copying method."""
        logger.warning("Using fallback slide copying method")

        # Get the source slide's pptx slide
        source_pptx_slide = source_slide.pptx_slide

        # Get the layout from the source slide
        try:
            layout = source_pptx_slide.slide_layout
        except Exception:
            # Use blank layout if we can't get the source layout
            layout_items_count = [
                len(layout.placeholders) for layout in self.pptx_presentation.slide_layouts
            ]
            min_items = min(layout_items_count)
            blank_layout_id = layout_items_count.index(min_items)
            layout = self.pptx_presentation.slide_layouts[blank_layout_id]

        # Create a new slide with the layout
        new_slide = self.pptx_presentation.slides.add_slide(layout)
        _remove_layout_placeholders(new_slide)

        # Copy speaker notes if they exist
        if source_pptx_slide.has_notes_slide and source_pptx_slide.notes_slide.notes_text_frame:
            try:
                notes_text = source_pptx_slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    new_slide.notes_slide.notes_text_frame.text = notes_text
            except Exception as e:
                logger.warning(f"Could not copy notes in fallback method: {e}")

        # Create and return a PowerPointSlide wrapper
        new_slide_obj = PowerPointSlide(new_slide)
        # Propagate presentation_id and pptx_presentation to the new slide and its elements
        new_slide_obj.presentation_id = self.presentationId
        new_slide_obj.pptx_presentation = self.pptx_presentation

        # Set parent references for the new slide and its elements
        new_slide_obj._parent_presentation = self
        for element in new_slide_obj.elements:
            element._parent_presentation = self

        return new_slide_obj

    def delete_slide(self, slide: Union["PowerPointSlide", int], api_client: PowerPointAPIClient):
        """Delete a slide from the presentation using robust SlideDeleter."""
        from gslides_api.pptx.slide_deleter import SlideDeleter

        if isinstance(slide, int):
            slide_idx = slide
        else:
            slide_idx = self.slides.index(slide)

        # Use the robust SlideDeleter implementation
        deleter = SlideDeleter(self.pptx_presentation)
        result = deleter.delete_slide(slide_idx)

        if result.success:
            # Update our slides list to reflect the deletion
            self.slides.pop(slide_idx)
            logger.info(f"Successfully deleted slide at index {slide_idx}")

            # Log any warnings
            for warning in result.warnings:
                logger.warning(f"Slide deletion warning: {warning}")
        else:
            # Deletion failed, raise an error
            error_msg = f"Failed to delete slide at index {slide_idx}: {result.error_message}"
            logger.error(error_msg)
            raise ValueError(error_msg)

    def delete_slides(
        self, slides: List[Union["PowerPointSlide", int]], api_client: PowerPointAPIClient
    ):
        """Delete multiple slides from the presentation using robust SlideDeleter."""
        from gslides_api.pptx.slide_deleter import SlideDeleter

        # Convert all slides to indices
        slide_indices = []
        for slide in slides:
            if isinstance(slide, int):
                slide_indices.append(slide)
            else:
                # Convert PowerPointSlide to index
                try:
                    slide_idx = self.slides.index(slide)
                    slide_indices.append(slide_idx)
                except ValueError:
                    logger.warning(f"Slide {slide} not found in presentation")
                    continue

        if not slide_indices:
            logger.info("No valid slides to delete")
            return

        # Use SlideDeleter which handles reverse order deletion internally
        deleter = SlideDeleter(self.pptx_presentation)
        results = deleter.delete_slides(slide_indices)

        # Process results and update local slides list
        successfully_deleted = []
        for result in results:
            if result.success:
                successfully_deleted.append(result.slide_index)
                # Log any warnings
                for warning in result.warnings:
                    logger.warning(f"Slide deletion warning: {warning}")
            else:
                logger.error(
                    f"Failed to delete slide at index {result.slide_index}: {result.error_message}"
                )

        # Remove successfully deleted slides from our list (in reverse order)
        for idx in sorted(successfully_deleted, reverse=True):
            self.slides.pop(idx)

        logger.info(f"Deleted {len(successfully_deleted)}/{len(slide_indices)} slides successfully")

    def move_slide(
        self,
        slide: Union["PowerPointSlide", int],
        insertion_index: int,
        api_client: PowerPointAPIClient,
    ):
        """Move a slide to a new position within the presentation."""
        if isinstance(slide, int):
            slide = self.slides[slide]

        # Note: python-pptx doesn't support direct slide reordering
        # Log the operation for now
        current_index = self.slides.index(slide)
        logger.debug(f"Moving slide from index {current_index} to {insertion_index}")
        # TODO: Implement actual XML-based slide reordering

        # For now, just update our local slides list
        self.slides.remove(slide)
        self.slides.insert(insertion_index, slide)

    def duplicate_slide(
        self, slide: Union["PowerPointSlide", int], api_client: PowerPointAPIClient
    ) -> "PowerPointSlide":
        """Duplicate a slide within the presentation."""
        if isinstance(slide, int):
            slide = self.slides[slide]

        # Use the robust SlideCopierManager
        from gslides_api.pptx.slide_copier import SlideCopierManager

        with SlideCopierManager(self.pptx_presentation) as copier:
            new_pptx_slide = copier.copy_slide_safe(slide.pptx_slide, fallback_to_layout_only=True)

            if new_pptx_slide:
                new_slide = PowerPointSlide(new_pptx_slide)
                # Propagate presentation_id to the new slide and its elements
                new_slide.presentation_id = self.presentationId
                new_slide.pptx_presentation = self.pptx_presentation
                self.slides.append(new_slide)
                return new_slide
            else:
                raise RuntimeError("Failed to duplicate slide")

    async def get_slide_thumbnails(
        self,
        api_client: "PowerPointAPIClient",
        slides: Optional[List["AbstractSlide"]] = None,
    ) -> List[AbstractThumbnail]:
        """Get thumbnails for slides using efficient batch PDF conversion.

        Converts PPTX to PDF once and extracts all slide images in a single operation,
        which is much more efficient than converting per-slide.

        Args:
            api_client: The PowerPoint API client
            slides: Optional list of slides to get thumbnails for. If None, uses all slides.

        Returns:
            List of AbstractThumbnail objects with image data
        """
        from io import BytesIO

        from PIL import Image

        from gslides_api.pptx.chart_renderer import render_all_slides_to_images

        if not self.file_path:
            logger.warning("Cannot generate thumbnails: no file path")
            return []

        # Determine which slides to get thumbnails for
        target_slides = slides if slides is not None else self.slides

        # If we need specific slides, we need to map their indices
        if slides is not None:
            slide_indices = []
            for slide in slides:
                try:
                    idx = self.slides.index(slide)
                    slide_indices.append(idx)
                except ValueError:
                    logger.warning(f"Slide not found in presentation: {slide}")
                    slide_indices.append(-1)  # Mark as not found
        else:
            slide_indices = list(range(len(self.slides)))

        # Save presentation to ensure latest changes are on disk
        self.save(api_client)

        # Render all slides at once (efficient: single PPTX->PDF conversion)
        all_png_bytes = await render_all_slides_to_images(
            presentation_path=self.file_path,
            dpi=150,
        )

        if not all_png_bytes:
            logger.warning("Batch rendering produced no images")
            return []

        # Build result list, extracting only the requested slides
        thumbnails = []
        for idx in slide_indices:
            if idx < 0 or idx >= len(all_png_bytes):
                # Slide not found or out of range, use placeholder
                thumbnails.append(AbstractThumbnail(
                    contentUrl="",
                    width=320,
                    height=240,
                    mime_type="image/png",
                    content=b"",
                    file_size=0,
                ))
                continue

            png_bytes = all_png_bytes[idx]

            # Get dimensions from the image
            try:
                img = Image.open(BytesIO(png_bytes))
                width, height = img.size
                img.close()
            except Exception as e:
                logger.warning(f"Failed to read image dimensions: {e}")
                width, height = 320, 240

            thumbnails.append(AbstractThumbnail(
                contentUrl="",
                width=width,
                height=height,
                mime_type="image/png",
                content=png_bytes,
                file_size=len(png_bytes),
            ))

        return thumbnails


#
# class PowerPointLayoutMatcher:
#     """PowerPoint implementation of AbstractLayoutMatcher using python-pptx."""
#
#     # DO NOT implement this - I'll move the matcher to abstract_slides.py
#
#     def __init__(self, presentation, matching_rule: Optional[str] = None):
#         # Layout matching will be implemented in the abstract base class
#         pass
#
#     def match(self, layout, matching_rule: Optional[str] = None):
#         # Layout matching will be implemented in the abstract base class
#         pass
